# -*- coding: utf-8 -*-
from pathlib import Path

from default_template import (
    DEFAULT_LAYOUT_NAME,
    DEFAULT_SHAPE_NAMES,
    META_TABLE_NAMES,
    PLACEHOLDER_IDX,
    build_default_template,
    default_template_mapping,
)
from pptx import Presentation
from slide_layout import P_NS

# 실제 화면설계서 '내용설명연결' 레이아웃에서 잰 값. 소스의 상수를 그대로
# 끌어다 쓰면 값이 바뀌어도 통과하는 검사가 되므로 여기에 따로 적는다.
MEASURED_PLACEHOLDERS = {
    1: (7496632, 188657, 1635346, 144000),    # 화면ID
    11: (0, 6750725, 2444995, 75085),          # 문서제목
    12: (4370029, 6716266, 403943, 144000),    # 쪽번호
}
MEASURED_SHELL = {
    "본문박스": (-1, 404664, 9131979, 6264697),
    "화면ID배경": (7496632, 188657, 1635346, 144000),
    "하단바": (0, 6716266, 9144000, 144000),
}
MEASURED_META_TABLES = {
    "메타표1": {
        "box": (1, 0, 9134931, 116632),
        "cols": 18,
        "labels": ["프로젝트명", "산출물명", "화면명", "버전", "작성자",
                   "검토자", "작성일", "수정일", "ID"],
    },
    "메타표2": {
        "box": (1, 195617, 7430165, 116632),
        "cols": 6,
        "labels": ["네비게이션", "화면유형", "알림여부"],
    },
}


def _layout(path: Path):
    prs = Presentation(str(path))
    return next(lay for master in prs.slide_masters
                for lay in master.slide_layouts
                if lay.name == DEFAULT_LAYOUT_NAME)


def test_default_template_uses_measured_size(tmp_path: Path):
    prs = Presentation(str(build_default_template(tmp_path / "d.pptx")))
    assert prs.slide_width == 9144000
    assert prs.slide_height == 6858000


def test_default_template_size_is_configurable(tmp_path: Path):
    prs = Presentation(str(build_default_template(
        tmp_path / "d.pptx", slide_width_emu=12192000, slide_height_emu=6858000)))
    assert prs.slide_width == 12192000


def test_default_template_has_no_slides(tmp_path: Path):
    """슬라이드는 빌드가 레이아웃으로 만든다. 템플릿에 예시가 필요 없다."""
    prs = Presentation(str(build_default_template(tmp_path / "d.pptx")))
    assert len(list(prs.slides)) == 0


def test_default_template_layout_exists(tmp_path: Path):
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    assert lay.name == DEFAULT_LAYOUT_NAME


def test_default_template_placeholders_match_measurements(tmp_path: Path):
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    got = {ph.placeholder_format.idx: (ph.left, ph.top, ph.width, ph.height)
           for ph in lay.placeholders}
    for idx, geom in MEASURED_PLACEHOLDERS.items():
        assert got[idx] == geom, "idx=%d" % idx


def test_default_template_shell_is_on_the_layout(tmp_path: Path):
    """껍데기는 레이아웃이 담당한다. 슬라이드마다 다시 그리지 않는다."""
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    by_name = {s.name: s for s in lay.shapes if not s.is_placeholder}
    for name, geom in MEASURED_SHELL.items():
        assert name in by_name, name
        s = by_name[name]
        assert (s.left, s.top, s.width, s.height) == geom, name


def test_default_template_meta_tables_match_measurements(tmp_path: Path):
    """상단 메타 표는 라벨과 값이 번갈아 놓인 1행짜리 표다."""
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    by_name = {s.name: s for s in lay.shapes if s.has_table}
    for name, spec in MEASURED_META_TABLES.items():
        assert name in by_name, name
        shp = by_name[name]
        assert (shp.left, shp.top, shp.width) == spec["box"][:3], name
        table = shp.table
        assert len(table.columns) == spec["cols"], name
        assert len(table.rows) == 1, name
        assert [table.cell(0, i).text
                for i in range(0, spec["cols"], 2)] == spec["labels"], name
        # 값 칸은 비어 있되 런은 있어야 한다 — 런이 없으면 빌드가 값을 채울 때
        # 글자 크기가 기본값(18pt)으로 잡혀 칸을 넘는다.
        for i in range(1, spec["cols"], 2):
            para = table.cell(0, i).text_frame.paragraphs[0]
            assert table.cell(0, i).text == ""
            assert para.runs and para.runs[0].font.size is not None


def test_default_template_placeholders_start_empty(tmp_path: Path):
    """표지가 없는 Excel에서 자리표시 문구가 산출물에 찍히면 안 된다."""
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    for ph in lay.placeholders:
        if ph.placeholder_format.idx == 12:
            continue  # 쪽번호는 자동 필드다
        assert ph.text_frame.text == "", ph.placeholder_format.idx


def test_default_template_shell_sits_behind_every_placeholder(tmp_path: Path):
    """껍데기는 전부 placeholder보다 뒤에 깔려야 한다 — 화면ID배경 위에 화면ID
    글자가 보이는 것이 대표적인 경우다.

    예전 테스트는 '화면ID배경'과 idx=1 placeholder 한 쌍의 순서만 봤다. 그런데
    build_default_template은 껍데기를 전부 붙인 뒤 *모든* placeholder를 spTree
    맨 뒤로 다시 붙이므로, 그 한 쌍의 비교는 SHELL_ORDER를 어떻게 바꾸든 항상
    참이라 깨질 수가 없었다. 실제 불변식은 '껍데기 전부가 placeholder 전부보다
    앞'이고, 껍데기를 placeholder 재배치 뒤에 붙이는 순간 이 검사는 실패한다.
    """
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    shells = [i for i, s in enumerate(lay.shapes) if not s.is_placeholder]
    phs = [i for i, s in enumerate(lay.shapes) if s.is_placeholder]

    # 껍데기와 메타 표가 전부 레이아웃에 있어야 아래 비교가 의미를 갖는다
    assert ({s.name for s in lay.shapes if not s.is_placeholder}
            == set(MEASURED_SHELL) | set(MEASURED_META_TABLES))
    assert len(phs) == len(MEASURED_PLACEHOLDERS)
    assert max(shells) < min(phs), [s.name for s in lay.shapes]


def test_default_template_layout_shape_ids_are_unique(tmp_path: Path):
    """cNvPr/@id는 한 파트 안에서 유일해야 한다 — 겹치면 PowerPoint가 파일을
    열 때 복구를 요구한다. 껍데기 도형을 임시 슬라이드에서 deepcopy로 이식할 때
    슬라이드에서 매겨진 id가 따라와 레이아웃 placeholder의 id와 겹쳤다."""
    lay = _layout(build_default_template(tmp_path / "d.pptx"))
    ids = [el.get("id") for el
           in lay.shapes._spTree.iter("{%s}cNvPr" % P_NS)]
    dupes = sorted({i for i in ids if ids.count(i) > 1})
    assert dupes == [], "중복된 cNvPr/@id: %s" % dupes


def test_default_template_carries_no_third_party_copyright(tmp_path: Path):
    """기본 템플릿은 코드로 배포되므로 남의 저작권 표기가 들어가면 안 된다."""
    prs = Presentation(str(build_default_template(tmp_path / "d.pptx")))
    texts = []
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            for s in lay.shapes:
                if s.has_text_frame:
                    texts.append(s.text_frame.text)
    blob = "\n".join(texts)
    for banned in ("Copyright", "ⓒ", "Nurimedia", "All rights reserved"):
        assert banned not in blob, banned


def test_default_template_mapping_is_layout_mode(tmp_path: Path):
    path = build_default_template(tmp_path / "d.pptx")
    m = default_template_mapping(path)
    assert m["mode"] == "layout"
    assert m["layout"] == DEFAULT_LAYOUT_NAME
    assert "source_slide" not in m
    assert m["placeholders"] == PLACEHOLDER_IDX
    assert m["detail_tables"] == {"count": 5, "rows": 4}
    assert m["shapes"]["screen_id"] == DEFAULT_SHAPE_NAMES["screen_id"]
    # 화면명·작성일은 placeholder가 아니라 메타 표가 담당한다
    assert m["meta_table"]["tables"] == META_TABLE_NAMES
    assert m["meta_table"]["labels"]["화면명"] == "title"
    assert m["meta_table"]["labels"]["ID"] == "screen_id"
    assert "title" not in m["shapes"]
    assert m["shapes"]["detail_tables"] == [
        "상세표1", "상세표2", "상세표3", "상세표4", "상세표5"]
    assert m["table_columns"] == {"no": 0, "text": 1}


def test_default_template_mapping_is_buildable(tmp_path: Path):
    """제안 매핑을 그대로 build에 넘겨 슬라이드가 나와야 한다."""
    from build import build
    from common import Warnings

    path = build_default_template(tmp_path / "d.pptx")
    mapping = {
        "version": 1,
        "excel": {"layout": "sheet-per-screen"},
        "template": default_template_mapping(path),
        "options": {"detail_text_source": "desc", "clear_unused_slots": True},
    }
    screens = {
        "meta": {"문서제목": "화면설계서"},
        "screens": [{"id": "SCR001", "name": "목록", "images": [], "fields": {},
                     "details": [{"no": "1", "desc": "설명"}]}],
    }
    out = tmp_path / "out.pptx"
    report = build(screens, mapping, tmp_path, out, Warnings())
    assert report["slides"] == 1
    slide = Presentation(str(out)).slides[0]
    by_name = {s.name: s for s in slide.shapes}
    assert by_name["문서제목"].text_frame.text == "화면설계서"
    # 화면명은 상단 메타 표의 '화면명' 칸 자리에 얹힌 글자로 들어간다
    from slide_layout import meta_slot_name
    assert by_name[meta_slot_name("화면명")].text_frame.text == "목록"
    assert by_name[meta_slot_name("ID")].text_frame.text == "SCR001"
