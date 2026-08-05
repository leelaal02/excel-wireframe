from pathlib import Path

from build import build
from common import Warnings
from fixtures import make_png, make_template_pptx
from pptx import Presentation
from verify import verify_output


def _mapping(template: Path) -> dict:
    return {
        "template": {
            "file": str(template),
            "mode": "clone",
            "source_slide": 0,
            "shapes": {
                "title": "제목 13",
                "screen_id": "텍스트 개체 틀 14",
                "image": "그림 18",
                "detail_tables": ["표 7", "표 8", "표 9", "표 10", "표 11"],
            },
            "table_columns": {"no": 0, "text": 1},
        },
        "options": {"detail_text_source": "desc", "clear_unused_slots": True},
    }


def _data(image: str | None) -> dict:
    return {
        "meta": {"title": "화면설계서"},
        "screens": [
            {"id": "SCR001", "name": "이용기관 목록",
             "images": [image] if image else [], "fields": {},
             "details": [{"no": "1", "desc": "등록한다"}]}
        ],
    }


def test_verify_counts_slot_filled_when_only_number_is_present(tmp_path: Path):
    """회귀 1: 상세 항목 수 검사가 설명 칸만 봐서, desc가 빈 문자열인 상세
    행(다른 칸 — no, 요소명 등 — 은 차 있는)이 있는 정상 생성물을 실패로
    보고했다. _read_details는 매핑된 칸 중 하나라도 값이 있으면 그 행을
    유지하므로 이런 상세는 실제로 흔하다. 번호 칸은 Excel에서 왔으니
    설명이 비어도 항상 있다 — 번호 또는 설명 중 하나만 있어도 슬롯이
    채워진 것으로 세어야 한다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")
    data["screens"][0]["details"] = [
        {"no": "1", "desc": "등록한다"},
        {"no": "2", "desc": ""},
        {"no": "3", "desc": "삭제한다"},
    ]
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())

    result = verify_output(out, data, _mapping(tpl), report["slides"])
    assert result["ok"] is True
    assert all(c["ok"] for c in result["checks"])


def test_verify_does_not_double_count_when_screen_name_is_a_substring(tmp_path: Path):
    """회귀 2: '목록'과 '이용기관 목록'처럼 한 화면명이 다른 화면명의 부분
    문자열이면, 부분일치 매칭은 '이용기관 목록' 슬라이드를 '목록' 화면에도
    매치시켜 상세 항목 수 검사가 두 화면 몫을 합산해 오탐했다. 화면명 반영/
    이미지 배치 검사는 매치 유무만 보므로 이 과다매치가 예전엔 무해했다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "S1.png")
    make_png(tmp_path / "images" / "S2.png")
    data = {
        "meta": {"title": "화면설계서"},
        "screens": [
            {"id": "S1", "name": "목록", "images": ["images/S1.png"], "fields": {},
             "details": [{"no": "1", "desc": "a"}]},
            {"id": "S2", "name": "이용기관 목록", "images": ["images/S2.png"], "fields": {},
             "details": [{"no": "1", "desc": "b"}, {"no": "2", "desc": "c"}]},
        ],
    }
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())

    result = verify_output(out, data, _mapping(tpl), report["slides"])
    assert result["ok"] is True
    assert all(c["ok"] for c in result["checks"])


def test_verify_skips_detail_count_when_clear_unused_slots_is_false(tmp_path: Path):
    """회귀 3: clear_unused_slots가 꺼지면 안 쓰는 슬롯은 템플릿의 예시 텍스트
    ('예시 설명 N')를 그대로 유지한다. 상세 항목 수 검사가 옵션을 모르면 그
    예시 텍스트까지 '채워진 슬롯'으로 세어, 상세가 표 슬롯 수보다 적은
    화면을 전부 영구적으로 실패시켰다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")  # 상세 1건, 표는 20슬롯
    mapping = _mapping(tpl)
    mapping["options"]["clear_unused_slots"] = False
    out = tmp_path / "out.pptx"
    report = build(data, mapping, tmp_path, out, Warnings())

    result = verify_output(out, data, mapping, report["slides"])
    assert result["ok"] is True
    checks = {c["name"]: c for c in result["checks"]}
    assert checks["상세 항목 수"]["ok"] is True
    assert "clear_unused_slots" in checks["상세 항목 수"]["detail"]


def test_verify_passes_for_workdir_relative_template_path(tmp_path: Path):
    """블로킹 발견 1: build()는 template.file이 상대경로면 work_dir 기준으로
    다시 찾지만, verify_output은 예전엔 그렇게 하지 않아 정상적으로 만들어진
    결과물에서도 '템플릿 파일을 찾을 수 없어 비교 불가'로 실패 보고를 했다.
    mapping-schema.md 예시 그대로(work 디렉토리 상대경로)를 재현한다."""
    make_template_pptx(tmp_path / "t.pptx")  # 절대경로로 실제 생성
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")
    mapping = _mapping(tmp_path / "t.pptx")
    mapping["template"]["file"] = "t.pptx"  # work_dir 기준 상대경로로 바꿔 둔다
    out = tmp_path / "out.pptx"
    report = build(data, mapping, tmp_path, out, Warnings())

    result = verify_output(out, data, mapping, report["slides"], tmp_path)
    assert result["ok"] is True
    assert all(c["ok"] for c in result["checks"])


def test_verify_fails_when_screens_are_empty(tmp_path: Path):
    """블로킹 발견 2: sheet_include 오타 등으로 화면이 통째로 0개면
    expected_slides도 같은 빈 입력에서 파생되어 '기대 0장, 실제 0장'으로
    공허하게 통과한다. 화면 데이터 자체가 비어 있다는 사실을 별도로 잡아야
    빈 매핑 사고가 '검증: 통과'로 보고되지 않는다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = {"meta": {"title": "화면설계서"}, "screens": []}
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())
    assert report["slides"] == 0

    result = verify_output(out, data, _mapping(tpl), report["slides"])
    assert result["ok"] is False
    failed = {c["name"]: c for c in result["checks"] if not c["ok"]}
    assert "화면 데이터" in failed
    assert "sheet_include" in failed["화면 데이터"]["detail"]


def test_verify_detects_wrong_detail_table_names(tmp_path: Path):
    """블로킹 발견 5: detail_tables 이름이 실제 템플릿과 다르면 collect_tables가
    빈 목록을 반환해 fill_slots가 전혀 돌지 않는다. 표는 템플릿의 예시 텍스트
    ('예시 설명 1'...)를 그대로 유지한 채 저장되지만, 앞의 네 검사(슬라이드 수/
    화면 데이터/화면명 반영/이미지 배치)는 이 상황에서도 전부 통과한다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")
    bad_mapping = _mapping(tpl)
    bad_mapping["template"]["shapes"]["detail_tables"] = [
        "없는표1", "없는표2", "없는표3", "없는표4", "없는표5"
    ]
    out = tmp_path / "out.pptx"
    warns = Warnings()
    report = build(data, bad_mapping, tmp_path, out, warns)

    result = verify_output(out, data, bad_mapping, report["slides"])
    assert result["ok"] is False
    failed = {c["name"]: c for c in result["checks"] if not c["ok"]}
    assert "상세 항목 수" in failed
    assert "SCR001" in failed["상세 항목 수"]["detail"]
    # 나머지 네 검사는 이 상황에서도 여전히 통과해야 한다 — 그래서 이 검사가
    # 없으면 사고가 통째로 묻힌다는 리뷰 지적이 성립한다.
    assert failed.keys() == {"상세 항목 수"}


def test_verify_passes_for_good_output(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())

    result = verify_output(out, data, _mapping(tpl), report["slides"])
    assert result["ok"] is True
    assert all(c["ok"] for c in result["checks"])


def test_verify_passes_for_layout_mode_output(tmp_path: Path):
    """최종 리뷰 지적 2: 이 파일도 test_sample_e2e.py도 전부 clone 모드였다.
    설계는 "verify.py는 손댈 필요가 없다"를 명시적 약속으로 두는데, 그 약속은
    layout 모드에만 있는 장치에 기대고 있다 — name_placeholders가 placeholder
    이름을 shapes의 이름으로 바꿔 주기 때문에 _title_texts가 제목 도형만 보는
    엄격한 경로를 탄다. 그 장치가 무너지면 검증은 '아무 도형의 텍스트나 본다'는
    느슨한 경로로 조용히 내려앉고, 테스트는 그대로 초록이었다."""
    from default_template import build_default_template, default_template_mapping

    template = build_default_template(tmp_path / "d.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    mapping = {
        "version": 1,
        "excel": {"layout": "sheet-per-screen"},
        "template": default_template_mapping(template),
        "options": {"detail_text_source": "desc", "clear_unused_slots": True},
    }
    data = {
        "meta": {"문서제목": "화면설계서"},
        "screens": [
            {"id": "SCR001", "name": "이용기관 목록",
             "images": ["images/SCR001.png"], "fields": {},
             "details": [{"no": "1", "desc": "등록한다"},
                         {"no": "2", "desc": "삭제한다"}]},
        ],
    }
    out = tmp_path / "out.pptx"
    report = build(data, mapping, tmp_path, out, Warnings())

    result = verify_output(out, data, mapping, report["slides"], tmp_path)

    assert {c["name"]: c["ok"] for c in result["checks"]} == {
        "슬라이드 수": True,
        "화면 데이터": True,
        "화면명 반영": True,
        "이미지 배치": True,
        "상세 항목 수": True,
        "슬라이드 크기": True,
    }
    assert result["ok"] is True
    # 엄격 경로를 실제로 탔는지 확인한다: 제목 도형이 매핑이 정한 이름을 달고
    # 있어야 _title_texts가 그 도형만 본다. 이름이 'Title 1'로 남아 있으면 위
    # 검사들은 느슨한 경로로도 전부 통과하므로 이 확인이 따로 필요하다.
    slide = Presentation(str(out)).slides[0]
    title_name = mapping["template"]["shapes"]["title"]
    assert [s.name for s in slide.shapes if s.name == title_name] == [title_name]
    assert next(s for s in slide.shapes if s.name == title_name)\
        .text_frame.text == "이용기관 목록"


def test_verify_detects_slide_count_mismatch(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _data(None)
    out = tmp_path / "out.pptx"
    build(data, _mapping(tpl), tmp_path, out, Warnings())

    result = verify_output(out, data, _mapping(tpl), 5)
    assert result["ok"] is False
    failed = [c["name"] for c in result["checks"] if not c["ok"]]
    assert "슬라이드 수" in failed


def test_verify_fails_size_check_when_template_missing(tmp_path: Path):
    """비교 대상 템플릿이 없으면 크기 검사를 건너뛰지 말고 실패로 보고해야 한다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())

    bad_mapping = _mapping(tpl)
    bad_mapping["template"]["file"] = str(tmp_path / "does-not-exist.pptx")

    result = verify_output(out, data, bad_mapping, report["slides"])
    assert result["ok"] is False
    failed = {c["name"]: c for c in result["checks"] if not c["ok"]}
    assert "슬라이드 크기" in failed
    assert "does-not-exist.pptx" in failed["슬라이드 크기"]["detail"]
    # 다른 검사는 여전히 정상이어야 한다 — 크기 검사만 실패로 좁혀졌는지 확인
    assert failed.keys() == {"슬라이드 크기"}


def test_verify_detects_missing_title_for_screen(tmp_path: Path):
    """화면명 반영 검사의 성공 경로만 지금까지 실행됐다 — 실패 경로(제목이 실제로
    슬라이드 어디에도 없는 경우)는 아무 테스트도 거치지 않았다. screens.json에
    빌드가 전혀 모르는 화면을 하나 더 끼워 넣어 재현한다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png")
    data = _data("images/SCR001.png")
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())

    # 검증 시점에만 화면을 하나 추가한다 — build()가 만든 실제 출력물에는
    # 이 화면에 해당하는 슬라이드가 존재하지 않는다.
    data_with_ghost = _data("images/SCR001.png")
    data_with_ghost["screens"].append(
        {"id": "GHOST", "name": "존재하지 않는 화면", "images": [], "fields": {},
         "details": []}
    )

    result = verify_output(out, data_with_ghost, _mapping(tpl), report["slides"])
    assert result["ok"] is False
    failed = {c["name"]: c for c in result["checks"] if not c["ok"]}
    assert "화면명 반영" in failed
    assert "존재하지 않는 화면" in failed["화면명 반영"]["detail"]


def test_verify_detects_missing_image_for_screen(tmp_path: Path):
    """상세에 이미지가 있다고 되어 있지만 실제 파일이 없어 배치가 실패하는 상황.

    build()는 화면 단위로 예외를 격리하므로 슬라이드는 만들어지지만(제목은
    '[생성 실패] ...'로 바뀐다) 그림 도형은 존재하지 않는다. 이미지 배치 검사가
    이 화면을 정확히 짚어내야 한다 — 전체 그림 개수만 세면 놓치는 사례다.
    """
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _data("images/missing.png")  # 파일을 만들지 않아 place_image가 실패한다
    out = tmp_path / "out.pptx"
    report = build(data, _mapping(tpl), tmp_path, out, Warnings())

    result = verify_output(out, data, _mapping(tpl), report["slides"])
    assert result["ok"] is False
    failed = {c["name"]: c for c in result["checks"] if not c["ok"]}
    assert "이미지 배치" in failed
    assert "SCR001" in failed["이미지 배치"]["detail"]
