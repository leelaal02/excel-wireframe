# -*- coding: utf-8 -*-
from pathlib import Path

from common import Warnings
from fixtures import make_png, make_template_pptx
from pptx import Presentation
from slide_fill import collect_tables, count_slots, fill_slots, find_shape, place_image

COLS = {"no": 0, "text": 1}


def _details(n: int) -> list[dict]:
    return [{"no": str(i + 1), "desc": "설명 %d" % (i + 1)} for i in range(n)]


def test_collect_tables_by_name(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], ["표 8", "표 7"])
    assert [t.name for t in tables] == ["표 8", "표 7"]


def test_collect_tables_left_to_right_when_no_names(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    assert [t.name for t in tables] == ["표 7", "표 8", "표 9", "표 10", "표 11"]
    assert count_slots(tables) == 20


def test_fill_slots_fills_in_order(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    filled = fill_slots(tables, _details(6), COLS, "desc", True, Warnings(), "S1")

    assert filled == 6
    assert tables[0].table.cell(0, 0).text == "1"
    assert tables[0].table.cell(0, 1).text == "설명 1"
    assert tables[1].table.cell(1, 1).text == "설명 6"


def test_fill_slots_clears_unused(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    fill_slots(tables, _details(6), COLS, "desc", True, Warnings(), "S1")
    assert tables[4].table.cell(3, 0).text == ""
    assert tables[4].table.cell(3, 1).text == ""


def test_fill_slots_keeps_unused_when_flag_off(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    fill_slots(tables, _details(2), COLS, "desc", False, Warnings(), "S1")
    assert tables[4].table.cell(3, 1).text == "예시 설명 20"


def test_fill_slots_warns_on_shortage(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    warns = Warnings()
    filled = fill_slots(tables, _details(25), COLS, "desc", True, warns, "S1")
    assert filled == 20
    assert [w["code"] for w in warns.to_list()] == ["slot-shortage"]


def test_fill_slots_uses_original_no(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    details = [{"no": "17", "desc": "열일곱"}, {"no": "18", "desc": "열여덟"}]
    fill_slots(tables, details, COLS, "desc", True, Warnings(), "S1")
    assert tables[0].table.cell(0, 0).text == "17"
    assert tables[0].table.cell(1, 0).text == "18"


def test_place_image_fits_and_centers(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    slide = prs.slides[0]
    anchor = find_shape(slide, "그림 18")
    left, top, w, h = anchor.left, anchor.top, anchor.width, anchor.height

    png = make_png(tmp_path / "shot.png", size=(1000, 200))
    pic = place_image(slide, anchor, png)

    assert find_shape(slide, "그림 18") is None
    assert pic.width <= w and pic.height <= h
    assert abs((pic.width / pic.height) - 5.0) < 0.05
    assert pic.left >= left and pic.top >= top
    assert abs((pic.left - left) - (left + w - (pic.left + pic.width))) <= 2


def test_collect_tables_warns_on_missing_name(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    warns = Warnings()
    tables = collect_tables(prs.slides[0], ["표 8", "표 999", "표 7"], warns, "S1")
    assert [t.name for t in tables] == ["표 8", "표 7"]
    assert len(warns.to_list()) == 1
    assert warns.to_list()[0]["code"] == "shape-not-found"
    assert "표 999" in warns.to_list()[0]["message"]


def test_fill_slots_no_overflow_warning_when_text_column_missing(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    tables = collect_tables(prs.slides[0], None)
    warns = Warnings()
    # Use a cols dict where text column index is beyond table width
    fill_slots(tables, _details(1), {"no": 0, "text": 5}, "desc", True, warns, "S1")
    # Should not produce text-overflow warning since text column doesn't exist
    overflow_warns = [w for w in warns.to_list() if w["code"] == "text-overflow"]
    assert len(overflow_warns) == 0
