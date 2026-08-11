from pathlib import Path

from analyze import build_report, main, resolve_template
from common import read_json
from fixtures import make_sheet_per_screen_xlsx, make_template_pptx

SCREENS = [
    {
        "id": "SCR001",
        "name": "이용기관 목록",
        "image": True,
        "details": [{"no": "1", "type": "버튼", "element": "[등록]", "desc": "등록한다", "pos": "상단"}],
    }
]


def test_build_report_joins_both_sides(tmp_path: Path):
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    pptx = make_template_pptx(tmp_path / "t.pptx")
    report = build_report(xlsx, pptx)
    assert [s["name"] for s in report["excel"]["sheets"]][1] == "설계_SCR001"
    assert report["template"]["slide_width"] == 9906000
    assert report["suggestion"]["mode"] == "clone"


def test_main_writes_report_file(tmp_path: Path):
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    pptx = make_template_pptx(tmp_path / "t.pptx")
    out = tmp_path / "output" / "structure-report.json"
    code = main(["--excel", str(xlsx), "--template", str(pptx), "--out", str(out)])
    assert code == 0
    data = read_json(out)
    assert data["suggestion"]["source_slide"] == 0
    assert data["template_generated"] is False


def test_main_puts_report_in_work_dir(tmp_path: Path):
    """--output만 주면 리포트 경로는 코드가 정한다."""
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    pptx = make_template_pptx(tmp_path / "t.pptx")
    out_dir = tmp_path / "output"

    code = main(["--excel", str(xlsx), "--template", str(pptx),
                 "--output", str(out_dir)])

    assert code == 0
    assert read_json(out_dir / ".work" / "structure-report.json")["template_generated"] is False
    assert not (out_dir / "structure-report.json").exists()


def test_main_puts_default_template_in_work_dir(tmp_path: Path):
    """만들어 쓰는 기본 템플릿도 작업 파일이다. 결과물 폴더에 pptx를 늘리지 않는다."""
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    out_dir = tmp_path / "output"

    assert main(["--excel", str(xlsx), "--output", str(out_dir)]) == 0

    assert (out_dir / ".work" / "default-template.pptx").exists()
    assert list(out_dir.glob("*.pptx")) == []


def test_main_requires_output_or_out(tmp_path: Path):
    """어디에 쓸지 모르는 채로 스캔을 시작하지 않는다."""
    import pytest

    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    with pytest.raises(SystemExit):
        main(["--excel", str(xlsx)])


def test_main_out_wins_over_output(tmp_path: Path):
    """명시 지정이 유도된 경로를 이긴다 — build.py의 --out-file과 같은 규칙이다."""
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    out_dir = tmp_path / "output"
    target = tmp_path / "리포트.json"

    assert main(["--excel", str(xlsx), "--output", str(out_dir),
                 "--out", str(target)]) == 0

    assert target.exists()
    assert not (out_dir / ".work" / "structure-report.json").exists()


def test_main_migrates_legacy_files(tmp_path: Path):
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    out_dir = tmp_path / "output"
    out_dir.mkdir()
    (out_dir / "mapping.json").write_text("{}", encoding="utf-8")

    assert main(["--excel", str(xlsx), "--output", str(out_dir)]) == 0

    assert (out_dir / ".work" / "mapping.json").exists()
    assert not (out_dir / "mapping.json").exists()


def test_resolve_template_uses_given_path(tmp_path: Path):
    pptx = make_template_pptx(tmp_path / "t.pptx")
    path, generated, cfg = resolve_template(str(pptx),
                                            tmp_path / "output" / "r.json")
    assert path == pptx
    assert generated is False
    assert cfg is None


def test_resolve_template_generates_when_missing(tmp_path: Path):
    out = tmp_path / "output" / "r.json"
    path, generated, cfg = resolve_template(None, out)
    assert generated is True
    assert cfg is None
    assert path == tmp_path / "output" / "default-template.pptx"
    assert path.exists()


def test_resolve_template_prefers_user_default_over_generating(tmp_path: Path):
    """사용자가 조직 템플릿을 등록해 두면 --template 없이도 그것을 쓴다."""
    from common import write_json
    from default_template import build_default_template

    base = tmp_path / "skill"
    (base / "assets").mkdir(parents=True)
    tpl = build_default_template(base / "assets" / "t.pptx")
    write_json(base / "user-default.json",
               {"template": "assets/t.pptx", "layout": "화면"})

    path, generated, cfg = resolve_template(
        None, tmp_path / "output" / "r.json", skill_base=base)

    assert path == tpl
    assert generated is False
    assert cfg["layout"] == "화면"
    assert not (tmp_path / "output" / "default-template.pptx").exists()


def test_resolve_template_given_path_wins_over_user_default(tmp_path: Path):
    """명시 지정이 사용자 기본 설정보다 우선한다."""
    from common import write_json
    from default_template import build_default_template

    base = tmp_path / "skill"
    (base / "assets").mkdir(parents=True)
    build_default_template(base / "assets" / "t.pptx")
    write_json(base / "user-default.json",
               {"template": "assets/t.pptx", "layout": "화면"})
    given = make_template_pptx(tmp_path / "given.pptx")

    path, generated, cfg = resolve_template(
        str(given), tmp_path / "output" / "r.json", skill_base=base)

    assert path == given
    assert cfg is None


def test_main_without_template_generates_one(tmp_path: Path):
    """generated=True일 때 실제로 쓰이는 제안은 `suggested_template_mapping`
    이다(build_default_template의 도형 이름을 이미 알고 있으므로 항상
    layout 모드로 확정해 채운다) — `suggestion`(범용 휴리스틱)이 아니다.

    기본 템플릿은 슬라이드가 0장이고 디자인을 레이아웃이 담당하므로
    `suggest_mode`도 자연히 layout으로 판정한다. 둘의 결론이 같아졌지만
    파이프라인이 실제로 쓰는 값은 여전히 `suggested_template_mapping`이다."""
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    out = tmp_path / "output" / "structure-report.json"
    code = main(["--excel", str(xlsx), "--out", str(out)])
    assert code == 0

    data = read_json(out)
    assert data["template_generated"] is True
    assert (tmp_path / "output" / "default-template.pptx").exists()
    assert data["suggestion"]["mode"] == "layout"

    suggested = data["suggested_template_mapping"]
    assert suggested["shapes"]["title"] == "제목"
    assert suggested["shapes"]["detail_tables"][0] == "상세표1"
    assert data["template"]["slide_width"] == 9906000


def test_main_with_given_template_has_no_suggestion(tmp_path: Path):
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    pptx = make_template_pptx(tmp_path / "t.pptx")
    out = tmp_path / "output" / "structure-report.json"
    main(["--excel", str(xlsx), "--template", str(pptx), "--out", str(out)])
    assert "suggested_template_mapping" not in read_json(out)


def test_suggested_mapping_carries_meta_shapes(tmp_path: Path):
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    out = tmp_path / "output" / "structure-report.json"
    main(["--excel", str(xlsx), "--out", str(out)])
    shapes = read_json(out)["suggested_template_mapping"]["shapes"]
    assert shapes["문서제목"] == "문서제목"
    assert shapes["작성일"] == "작성일"


def test_analyze_suggests_layout_mode_mapping(tmp_path: Path):
    from analyze import main
    from common import read_json
    from fixtures import make_sheet_per_screen_xlsx

    xlsx = make_sheet_per_screen_xlsx(
        tmp_path / "in.xlsx",
        [{"id": "SCR001", "name": "목록", "image": False,
          "details": [{"no": "1", "type": "버튼", "element": "[등록]",
                       "desc": "등록한다", "pos": "우상단"}]}],
    )
    out = tmp_path / "report.json"
    assert main(["--excel", str(xlsx), "--out", str(out)]) == 0

    report = read_json(out)
    assert report["template_generated"] is True
    tpl = report["suggested_template_mapping"]
    assert tpl["mode"] == "layout"

    # 최종 리뷰 지적 3: 예전엔 '키가 있고 길이가 4'만 봤다 — [0, 0, W, 0]도
    # 통과하는 검사였다. SKILL.md는 이 값을 그대로 복사하라고 안내하므로
    # 실제 실행에 쓰이는 content_area는 실측 폴백(DEFAULT_CONTENT_AREA)이
    # 아니라 여기 값이다. 껍데기를 실제로 비켜 갔는지 값으로 확인한다.
    from default_template import MEASURED_SHELL
    divider = MEASURED_SHELL["구분선"]      # 상단 껍데기 중 가장 아래
    footer = MEASURED_SHELL["하단바"]
    left, top, width, height = tpl["content_area"]
    assert left == 0
    assert top == divider[1] + divider[3] == 620688      # 구분선 바로 아래
    assert top + height == footer[1] == 6716266          # 하단바 바로 위
    assert height > 0
    assert width == report["template"]["slide_width"] == 9906000


def test_suggested_mapping_builds_with_its_own_content_area(tmp_path: Path):
    """최종 리뷰 지적 3: 실제로 쓰이는 content_area는 analyze.py가 계산한
    값인데(SKILL.md가 suggested_template_mapping을 그대로 복사하라고 안내한다),
    빌드 테스트들은 전부 실측 폴백을 쓰고 있어 analyze→build 이음매를 아무도
    지나가지 않았다. 제안 매핑을 손대지 않고 그대로 build에 넘긴다."""
    from build import build
    from common import Warnings
    from pptx import Presentation

    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SCREENS)
    out = tmp_path / "output" / "structure-report.json"
    assert main(["--excel", str(xlsx), "--out", str(out)]) == 0
    tpl = read_json(out)["suggested_template_mapping"]

    mapping = {
        "version": 1,
        "excel": {"layout": "sheet-per-screen"},
        "template": tpl,  # 그대로 복사 — content_area도 제안값 그대로다
        "options": {"detail_text_source": "desc", "clear_unused_slots": True},
    }
    screens = {
        "meta": {"문서제목": "화면설계서"},
        "screens": [{"id": "SCR001", "name": "이용기관 목록", "images": [],
                     "fields": {}, "details": [{"no": "1", "desc": "등록한다"}]}],
    }
    ppt = tmp_path / "out.pptx"
    warns = Warnings()
    report = build(screens, mapping, tmp_path, ppt, warns)
    assert report["slides"] == 1

    slide = Presentation(str(ppt)).slides[0]
    c_left, c_top, c_width, c_height = tpl["content_area"]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 5
    # 표는 본문 영역 아래쪽에 붙고(하단바 위), 이미지 자리는 그 위 나머지다
    assert tables[0].top + tables[0].height == c_top + c_height
    assert sum(t.width for t in tables) == c_width
    anchor = next(s for s in slide.shapes if s.name == "화면이미지")
    assert (anchor.left, anchor.top) == (c_left, c_top)
    assert anchor.top + anchor.height == tables[0].top
    assert [w["code"] for w in warns.to_list()] == ["no-image"]
