from pathlib import Path

from build import build, chunk_details, page_title
from build import main as build_cli
from common import Warnings, write_json
from fixtures import make_png, make_template_pptx
from pptx import Presentation


def _mapping(template: Path) -> dict:
    return {
        "version": 1,
        "excel": {"layout": "sheet-per-screen"},
        "template": {
            "file": str(template),
            "mode": "clone",
            "source_slide": 0,
            "shapes": {
                "title": "제목 13",
                "screen_id": "텍스트 개체 틀 14",
                "image": "그림 18",
                "detail_tables": ["표 7", "표 8", "표 9", "표 10", "표 11"],
            },
            "table_columns": {"no": 0, "text": 1},
        },
        "options": {
            "detail_text_source": "desc",
            "overflow": "split",
            "clear_unused_slots": True,
        },
    }


def _screens(n_details: int, image: str | None = None) -> dict:
    return {
        "meta": {"title": "화면설계서", "source": "s.xlsx", "template": "t.pptx"},
        "screens": [
            {
                "id": "SCR001",
                "name": "이용기관 목록",
                "images": [image] if image else [],
                "fields": {},
                "details": [
                    {"no": str(i + 1), "desc": "설명 %d" % (i + 1)}
                    for i in range(n_details)
                ],
            }
        ],
    }


def test_chunk_details_splits_by_slot_count():
    d = [{"no": str(i)} for i in range(25)]
    chunks = chunk_details(d, 20)
    assert [len(c) for c in chunks] == [20, 5]
    assert chunk_details([], 20) == [[]]
    assert len(chunk_details(d[:20], 20)) == 1


def test_page_title_marks_split():
    assert page_title("목록", 0, 1) == "목록"
    assert page_title("목록", 0, 2) == "목록 (1/2)"
    assert page_title("목록", 1, 2) == "목록 (2/2)"


def test_build_creates_one_slide_per_screen(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    report = build(_screens(6), _mapping(tpl), tmp_path, out, Warnings())

    assert report["slides"] == 1
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    title = next(s for s in slide.shapes if s.name == "제목 13")
    assert title.text_frame.text == "이용기관 목록"
    sid = next(s for s in slide.shapes if s.name == "텍스트 개체 틀 14")
    assert sid.text_frame.text == "SCR001"


def test_build_keeps_slide_size(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    build(_screens(3), _mapping(tpl), tmp_path, out, Warnings())
    assert Presentation(str(out)).slide_width == 9906000


def test_build_splits_on_overflow(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    warns = Warnings()
    report = build(_screens(25), _mapping(tpl), tmp_path, out, warns)

    assert report["slides"] == 2
    assert report["split"] == ["SCR001"]
    prs = Presentation(str(out))
    titles = [
        next(s for s in sl.shapes if s.name == "제목 13").text_frame.text
        for sl in prs.slides
    ]
    assert titles == ["이용기관 목록 (1/2)", "이용기관 목록 (2/2)"]
    assert "slide-split" in [w["code"] for w in warns.to_list()]


def test_build_places_image(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    make_png(tmp_path / "images" / "SCR001.png", size=(800, 400))
    out = tmp_path / "out.pptx"
    build(_screens(4, "images/SCR001.png"), _mapping(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    pics = [s for s in slide.shapes if s.shape_type == 13]
    assert len(pics) == 1
    assert not any(s.name == "그림 18" for s in slide.shapes)


def test_build_isolates_screen_failure(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _screens(2)
    data["screens"].append(
        {"id": "BROKEN", "name": "깨진 화면", "images": ["images/없는파일.png"],
         "fields": {}, "details": [{"no": "1", "desc": "x"}]}
    )
    out = tmp_path / "out.pptx"
    warns = Warnings()
    report = build(data, _mapping(tpl), tmp_path, out, warns)

    assert report["screens"] == 2
    assert report["failed"] == ["BROKEN"]
    assert out.exists()
    assert "screen-failed" in [w["code"] for w in warns.to_list()]

    prs = Presentation(str(out))
    assert len(prs.slides) == report["slides"]
    titles = [
        next(s for s in sl.shapes if s.name == "제목 13").text_frame.text
        for sl in prs.slides
    ]
    assert titles == ["이용기관 목록", "[생성 실패] 깨진 화면"]


def test_build_isolates_screen_missing_id(tmp_path: Path):
    """블로킹 발견 3: failed_ids.append(scr["id"])가 except 블록 안에 있으면,
    id가 아예 빠진 screen dict 하나가 scr["id"] 참조에서 또 KeyError를 내며
    빌드 전체를 무너뜨린다 — 결과물이 통째로 안 생긴다. 중간 산출물인
    screens.json에서 id가 빠지는 것(추출이 어긋났거나 사람이 손댄 경우)은
    화면 단위 예외 격리가 정확히 보호해야 하는 시나리오다. 다른 화면들은
    끝까지 완성되고 파일이 반드시 저장돼야 한다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _screens(2)
    data["screens"].append(
        {"name": "id 없는 화면", "images": [], "fields": {},
         "details": [{"no": "1", "desc": "x"}]}
    )
    out = tmp_path / "out.pptx"
    warns = Warnings()

    report = build(data, _mapping(tpl), tmp_path, out, warns)

    assert out.exists()
    assert report["screens"] == 2
    codes = [w["code"] for w in warns.to_list()]
    assert "screen-failed" in codes

    prs = Presentation(str(out))
    titles = [
        next(s for s in sl.shapes if s.name == "제목 13").text_frame.text
        for sl in prs.slides
    ]
    assert "이용기관 목록" in titles


def test_build_mode_layout_gives_readable_error_not_typeerror(tmp_path: Path):
    """블로킹 발견 4: source_slide: null인 채로 clone 모드로 빌드를 시도하면
    int(None)이 raw TypeError를 던졌다. 원인과 대안을 알려주는 ValueError로
    바뀌어야 한다. (mode: layout + source_slide: null 조합은 이제 실제로
    지원되는 정상 경로다 — test_build_layout_mode_* 테스트들이 그 경로를
    검증한다. 이 테스트는 clone 모드에서 source_slide 자체가 없는 경우만
    다룬다.)"""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    mapping = _mapping(tpl)
    mapping["template"]["source_slide"] = None
    out = tmp_path / "out.pptx"

    try:
        build(_screens(2), mapping, tmp_path, out, Warnings())
        assert False, "ValueError가 나야 한다"
    except TypeError:
        raise AssertionError("원시 TypeError가 그대로 새고 있습니다")
    except ValueError as exc:
        msg = str(exc)
        assert "source_slide" in msg
        assert "--template" in msg or "기본 템플릿" in msg


def test_build_does_not_import_openpyxl():
    import build as build_mod
    import inspect

    src = inspect.getsource(build_mod)
    assert "openpyxl" not in src


def _mapping_with_meta_shape(template: Path) -> dict:
    m = _mapping(template)
    m["template"]["shapes"]["프로젝트명"] = "텍스트 개체 틀 14"
    return m


def test_build_fills_meta_into_named_shape(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _screens(2)
    data["meta"]["프로젝트명"] = "통합관리시스템"
    out = tmp_path / "out.pptx"
    build(data, _mapping_with_meta_shape(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    shp = next(s for s in slide.shapes if s.name == "텍스트 개체 틀 14")
    assert shp.text_frame.text == "통합관리시스템"


def test_build_screen_fields_beat_document_meta(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _screens(2)
    data["meta"]["프로젝트명"] = "문서 전체 값"
    data["screens"][0]["fields"]["프로젝트명"] = "화면별 값"
    out = tmp_path / "out.pptx"
    build(data, _mapping_with_meta_shape(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    shp = next(s for s in slide.shapes if s.name == "텍스트 개체 틀 14")
    assert shp.text_frame.text == "화면별 값"


def test_build_ignores_meta_without_matching_shape(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _screens(2)
    data["meta"]["아무도_안_쓰는_키"] = "값"
    out = tmp_path / "out.pptx"
    warns = Warnings()
    build(data, _mapping(tpl), tmp_path, out, warns)
    assert "shape-not-found" not in [w["code"] for w in warns.to_list()]


def test_build_meta_does_not_override_title_shape(tmp_path: Path):
    tpl = make_template_pptx(tmp_path / "t.pptx")
    data = _screens(2)
    data["meta"]["title"] = "문서 제목이 화면명을 덮으면 안 된다"
    out = tmp_path / "out.pptx"
    build(data, _mapping(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    shp = next(s for s in slide.shapes if s.name == "제목 13")
    assert shp.text_frame.text == "이용기관 목록"


from pptx.util import Emu


def _layout_template(path: Path) -> Path:
    """placeholder만 있는 레이아웃형 템플릿 (예시 슬라이드 없음)."""
    prs = Presentation()
    prs.slide_width = Emu(9906000)
    prs.slide_height = Emu(6858000)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _layout_mapping(template: Path) -> dict:
    return {
        "version": 1,
        "excel": {"layout": "sheet-per-screen"},
        "template": {
            "file": str(template),
            "mode": "layout",
            "layout": "Title and Content",
            "placeholders": {"title": 0, "screen_id": 1, "작성일": 10,
                             "문서제목": 11},
            "shapes": {
                "title": "제목",
                "screen_id": "화면ID",
                "image": "화면이미지",
                "작성일": "작성일",
                "문서제목": "문서제목",
                "detail_tables": ["상세표1", "상세표2", "상세표3",
                                  "상세표4", "상세표5"],
            },
            "content_area": [-12319, 337940, 9957099, 6331421],
            "detail_tables": {"count": 5, "rows": 4},
            "table_columns": {"no": 0, "text": 1},
        },
        "options": {
            "detail_text_source": "desc",
            "overflow": "split",
            "clear_unused_slots": True,
        },
    }


def test_build_layout_mode_creates_slides(tmp_path: Path):
    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    report = build(_screens(6), _layout_mapping(tpl), tmp_path, out, Warnings())

    assert report["slides"] == 1
    slide = Presentation(str(out)).slides[0]
    by_name = {s.name: s for s in slide.shapes}
    assert by_name["제목"].text_frame.text == "이용기관 목록"
    assert by_name["화면ID"].text_frame.text == "SCR001"


def test_build_layout_mode_fills_detail_tables(tmp_path: Path):
    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    build(_screens(6), _layout_mapping(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 5
    first = next(t for t in tables if t.name == "상세표1").table
    assert first.cell(0, 0).text == "1"
    assert first.cell(0, 1).text == "설명 1"


def test_build_layout_mode_splits_when_details_exceed_slots(tmp_path: Path):
    """슬롯은 5표 × 4행 = 20개다."""
    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    warns = Warnings()
    report = build(_screens(25), _layout_mapping(tpl), tmp_path, out, warns)

    assert report["slides"] == 2
    assert report["split"] == ["SCR001"]
    assert any(w["code"] == "slide-split" for w in warns.to_list())


def test_build_layout_mode_places_image(tmp_path: Path):
    img = make_png(tmp_path / "images" / "SCR001.png")
    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    build(_screens(3, image="images/SCR001.png"), _layout_mapping(tpl),
          tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    assert any("PICTURE" in str(s.shape_type) for s in slide.shapes)
    assert not any(s.name == "화면이미지" for s in slide.shapes)


def test_build_layout_mode_without_image_shape_leaves_no_orphan_anchor(tmp_path: Path):
    """최종 리뷰 지적 1: 이미지 자리를 그리는 쪽(_new_layout_slide)은 기본 이름
    '화면이미지'를 쓰고, 그것을 그림으로 바꾸는 쪽(_fill_page)은 기본값 없이
    shapes.image를 읽었다. 그래서 매핑에서 shapes.image를 빼면 본문 영역만 한
    회색 사각형이 슬라이드마다 그대로 남았다 — 아무 경고도 없이. 레이아웃에는
    이미지 placeholder가 없으니 이름을 지을 게 없다고 판단하는 매핑 작성자가
    실제로 밟을 수 있는 길이다."""
    tpl = _layout_template(tmp_path / "t.pptx")
    mapping = _layout_mapping(tpl)
    del mapping["template"]["shapes"]["image"]
    out = tmp_path / "out.pptx"
    warns = Warnings()

    report = build(_screens(3, image="images/SCR001.png"), mapping, tmp_path,
                   out, warns)

    assert report["slides"] == 1
    slide = Presentation(str(out)).slides[0]
    names = [s.name for s in slide.shapes]
    assert "화면이미지" not in names
    # 이미지 자리로 쓰이던 큰 사각형이 다른 이름으로 남지도 않아야 한다.
    # (표 5개와 placeholder만 남는다)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    autoshapes = [s for s in slide.shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert autoshapes == [], [s.name for s in autoshapes]
    assert len([s for s in slide.shapes if s.has_table]) == 5
    assert warns.to_list() == []


def test_build_layout_mode_drops_empty_placeholders(tmp_path: Path):
    """값이 없는 문서제목 placeholder가 산출물에 남으면 안 된다."""
    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    screens = _screens(3)
    screens["meta"] = {"source": "s.xlsx", "template": "t.pptx"}  # 문서제목 없음
    build(screens, _layout_mapping(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    assert "문서제목" not in [s.name for s in slide.shapes]


def test_build_layout_mode_fills_doc_title_from_meta(tmp_path: Path):
    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    screens = _screens(3)
    screens["meta"]["문서제목"] = "발행기관관리"
    build(screens, _layout_mapping(tpl), tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    by_name = {s.name: s for s in slide.shapes}
    assert by_name["문서제목"].text_frame.text == "발행기관관리"


def test_build_layout_mode_raises_for_unknown_layout(tmp_path: Path):
    import pytest
    tpl = _layout_template(tmp_path / "t.pptx")
    mapping = _layout_mapping(tpl)
    mapping["template"]["layout"] = "없는레이아웃"
    with pytest.raises(ValueError) as exc:
        build(_screens(3), mapping, tmp_path, tmp_path / "out.pptx", Warnings())
    assert "없는레이아웃" in str(exc.value)


def test_build_layout_mode_raises_for_global_row_misconfig_not_per_screen(
    tmp_path: Path,
):
    """rows_per_table이 과도해 이미지 자리가 안 나오는 것은 화면 하나의 사고가
    아니라 mapping 전체에 걸린 설정 오류다. 화면 단위 격리(screen-failed)로
    흩어지면 build()가 조용히 끝나면서 모든 슬라이드가 '[생성 실패]'가 되어
    버린다 — 여기서는 그 대신 ValueError가 루프 밖으로 곧장 터져야 한다."""
    import pytest
    tpl = _layout_template(tmp_path / "t.pptx")
    mapping = _layout_mapping(tpl)
    mapping["template"]["detail_tables"]["rows"] = 40
    warns = Warnings()
    with pytest.raises(ValueError):
        build(_screens(6), mapping, tmp_path, tmp_path / "out.pptx", warns)
    assert "screen-failed" not in [w["code"] for w in warns.to_list()]


def test_build_clone_mode_still_works(tmp_path: Path):
    """layout 모드를 더해도 clone 경로는 그대로여야 한다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    report = build(_screens(6), _mapping(tpl), tmp_path, out, Warnings())
    assert report["slides"] == 1


def _screens_with_cover_date(date_text: str) -> dict:
    d = _screens(3)
    d["meta"]["작성일"] = date_text
    return d


def test_build_stamps_its_own_run_date(tmp_path: Path, monkeypatch):
    """화면설계서의 작성일은 그 PPT를 만든 날이다.

    Excel 표지의 작성일은 Excel을 쓴 날이라 다르다 — 실제 샘플에서 두 달
    차이가 났다. 표지 값이 있어도 생성일이 이겨야 한다.
    """
    import build as build_mod
    monkeypatch.setattr(build_mod, "_today", lambda: "2026-08-05")

    tpl = _layout_template(tmp_path / "t.pptx")
    out = tmp_path / "out.pptx"
    build(_screens_with_cover_date("2026-06-25"), _layout_mapping(tpl),
          tmp_path, out, Warnings())

    slide = Presentation(str(out)).slides[0]
    by_name = {s.name: s for s in slide.shapes}
    assert by_name["작성일"].text_frame.text == "2026-08-05"


def test_build_keeps_cover_date_when_date_field_disabled(tmp_path: Path,
                                                         monkeypatch):
    """options.date_field를 끄면 표지에서 읽은 값을 그대로 쓴다."""
    import build as build_mod
    monkeypatch.setattr(build_mod, "_today", lambda: "2026-08-05")

    tpl = _layout_template(tmp_path / "t.pptx")
    mapping = _layout_mapping(tpl)
    mapping["options"]["date_field"] = None
    out = tmp_path / "out.pptx"
    build(_screens_with_cover_date("2026-06-25"), mapping, tmp_path, out,
          Warnings())

    slide = Presentation(str(out)).slides[0]
    by_name = {s.name: s for s in slide.shapes}
    assert by_name["작성일"].text_frame.text == "2026-06-25"


def test_build_stamps_run_date_in_clone_mode_too(tmp_path: Path, monkeypatch):
    """생성일 도장은 모드와 무관하다."""
    import build as build_mod
    monkeypatch.setattr(build_mod, "_today", lambda: "2026-08-05")

    tpl = make_template_pptx(tmp_path / "t.pptx")
    mapping = _mapping(tpl)
    mapping["template"]["shapes"]["작성일"] = "작성일"
    out = tmp_path / "out.pptx"
    build(_screens_with_cover_date("2026-06-25"), mapping, tmp_path, out,
          Warnings())

    slide = Presentation(str(out)).slides[0]
    by_name = {s.name: s for s in slide.shapes}
    assert by_name["작성일"].text_frame.text == "2026-08-05"


def test_build_does_not_mutate_the_caller_screens_data(tmp_path: Path,
                                                       monkeypatch):
    """생성일을 얹느라 넘겨받은 screens_data를 건드리면 안 된다."""
    import build as build_mod
    monkeypatch.setattr(build_mod, "_today", lambda: "2026-08-05")

    tpl = _layout_template(tmp_path / "t.pptx")
    screens = _screens_with_cover_date("2026-06-25")
    build(screens, _layout_mapping(tpl), tmp_path, tmp_path / "out.pptx",
          Warnings())

    assert screens["meta"]["작성일"] == "2026-06-25"


def test_today_returns_iso_date():
    """_today는 표지 작성일과 같은 표기(YYYY-MM-DD)를 쓴다."""
    from datetime import date

    from build import _today

    assert _today() == date.today().isoformat()


def _details(n: int) -> list[dict]:
    return [{"no": str(i + 1), "desc": "설명 %d" % (i + 1)} for i in range(n)]


def _tall_png(path: Path, badges_at: list[int], w=400, h=1200) -> Path:
    """세로로 긴 가짜 스크린샷. 지정한 y에 노란 뱃지를 찍는다."""
    from PIL import Image, ImageDraw
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(im)
    for y in range(0, h, 300):          # 여백이 아닌 띠를 만들어 자를 곳을 제한한다
        d.rectangle([0, y, w, y + 250], fill=(230, 235, 245))
    for y in badges_at:
        d.ellipse([100, y - 12, 124, y + 12], fill=(255, 210, 40))
    im.save(path)
    return path


def test_plan_pages_splits_details_by_badge_count(tmp_path: Path):
    """조각별 뱃지 수만큼 상세가 배분된다."""
    from build import plan_pages
    img = _tall_png(tmp_path / "images" / "S.png", [100, 200, 700, 800, 900])
    scr = {"id": "S", "name": "화면", "images": ["images/S.png"],
           "fields": {}, "details": _details(5)}

    pages, note = plan_pages(scr, 20, tmp_path, (0, 0, 400, 200), True)

    assert note is None
    assert len(pages) >= 2
    assert sum(len(p) for p, _ in pages) == 5
    assert [d["no"] for d, _ in [(p[0], i) for p, i in pages if p]][0] == "1"
    assert all(i is not None for _, i in pages)


def test_plan_pages_repeats_piece_when_details_exceed_slots(tmp_path: Path):
    """한 조각의 상세가 슬롯을 넘으면 같은 조각을 여러 장에 싣는다."""
    from build import plan_pages
    # 뱃지 사이 틈이 여백 띠로 잡히지 않을 만큼 촘촘히 찍어, 이 구간이 통째로
    # 한 조각에 들어가게 만든다. 사이에서 잘리면 슬롯을 넘지 않아 반복이 안 난다.
    ys = list(range(60, 270, 28))
    img = _tall_png(tmp_path / "images" / "S.png", ys)
    scr = {"id": "S", "name": "화면", "images": ["images/S.png"],
           "fields": {}, "details": _details(len(ys))}

    pages, note = plan_pages(scr, 5, tmp_path, (0, 0, 400, 200), True)

    assert note is None
    imgs = [i for _, i in pages]
    assert len(imgs) > len(set(imgs)), "같은 조각이 반복돼야 한다"
    assert sum(len(p) for p, _ in pages) == len(ys)


def test_plan_pages_falls_back_when_badges_mismatch(tmp_path: Path):
    """뱃지 수가 상세 건수와 다르면 순차 배분으로 물러선다."""
    from build import plan_pages
    _tall_png(tmp_path / "images" / "S.png", [100, 700])
    scr = {"id": "S", "name": "화면", "images": ["images/S.png"],
           "fields": {}, "details": _details(9)}

    pages, note = plan_pages(scr, 4, tmp_path, (0, 0, 400, 200), True)

    assert note is not None and "맞지 않아" in note
    assert sum(len(p) for p, _ in pages) == 9


def test_plan_pages_skips_split_when_disabled(tmp_path: Path):
    from build import plan_pages
    _tall_png(tmp_path / "images" / "S.png", [100, 700])
    scr = {"id": "S", "name": "화면", "images": ["images/S.png"],
           "fields": {}, "details": _details(4)}

    pages, note = plan_pages(scr, 20, tmp_path, (0, 0, 400, 200), False)

    assert note is None
    assert len(pages) == 1
    assert pages[0][1] is None


def test_plan_pages_leaves_short_image_alone(tmp_path: Path):
    """자리보다 짧은 이미지는 나누지 않는다."""
    from build import plan_pages
    _tall_png(tmp_path / "images" / "S.png", [50], w=400, h=100)
    scr = {"id": "S", "name": "화면", "images": ["images/S.png"],
           "fields": {}, "details": _details(3)}

    pages, note = plan_pages(scr, 20, tmp_path, (0, 0, 400, 200), True)

    assert len(pages) == 1
    assert pages[0][1] is None


def test_build_clone_mode_splits_tall_image(tmp_path: Path):
    """clone 모드도 긴 스크린샷을 나눈다 — 이미지 자리는 앵커 도형이 알려 준다.

    상세는 슬롯(20)에 다 들어가므로 장이 늘어난 이유는 이미지 분할뿐이다.
    """
    tpl = make_template_pptx(tmp_path / "t.pptx")
    _tall_png(tmp_path / "images" / "SCR001.png", [100, 400, 700, 900, 1100])
    out = tmp_path / "out.pptx"

    report = build(_screens(5, "images/SCR001.png"), _mapping(tpl), tmp_path,
                   out, Warnings())

    assert report["slides"] > 1, "긴 이미지가 조각으로 나뉘어야 한다"
    prs = Presentation(str(out))
    assert all(len([s for s in sl.shapes if s.shape_type == 13]) == 1
               for sl in prs.slides), "장마다 조각이 한 장씩 놓여야 한다"


def test_build_clone_mode_without_image_anchor_skips_split(tmp_path: Path):
    """앵커 도형을 못 찾으면 분할 없이 한 장으로 끝낸다.

    앵커에서 자리를 읽는 코드가 None을 방어하지 않으면 여기서 AttributeError가
    나고, 그 계산은 화면 루프 밖이라 화면 단위 격리에도 걸리지 않는다.
    """
    tpl = make_template_pptx(tmp_path / "t.pptx")
    _tall_png(tmp_path / "images" / "SCR001.png", [100, 400, 700, 900, 1100])
    mapping = _mapping(tpl)
    mapping["template"]["shapes"]["image"] = "없는 그림"
    out = tmp_path / "out.pptx"

    report = build(_screens(5, "images/SCR001.png"), mapping, tmp_path, out,
                   Warnings())

    assert report["slides"] == 1


# --- 상세표 높이 맞추기 -------------------------------------------------

LONG_DESC = "가" * 300  # 7pt에서 15줄. 실측 행 높이(최대 4줄)로는 어림없다


def _screens_with_desc(descs: list[str]) -> dict:
    return {
        "meta": {"title": "화면설계서", "source": "s.xlsx", "template": "t.pptx"},
        "screens": [
            {
                "id": "SCR001",
                "name": "이용기관 목록",
                "images": [],
                "fields": {},
                "details": [{"no": str(i + 1), "desc": d}
                            for i, d in enumerate(descs)],
            }
        ],
    }


def _tables_of(slide):
    return [s for s in slide.shapes if s.has_table]


def _build_layout(tmp_path: Path, screens: dict, mapping: dict | None = None):
    tpl = _layout_template(tmp_path / "t.pptx")
    mapping = mapping or _layout_mapping(tpl)
    out = tmp_path / "out.pptx"
    build(screens, mapping, tmp_path, out, Warnings())
    return Presentation(str(out))


def test_build_layout_mode_grows_tables_upward_for_long_details(tmp_path: Path):
    """긴 상세가 들어가면 표가 위로 자란다. 아래 끝은 본문 영역에 붙어 있어야 한다."""
    area_top, area_h = 337940, 6331421
    short = _build_layout(tmp_path / "a", _screens_with_desc(["짧게"] * 4))
    tall = _build_layout(tmp_path / "b", _screens_with_desc([LONG_DESC] * 4))

    s_tbl = _tables_of(short.slides[0])[0]
    t_tbl = _tables_of(tall.slides[0])[0]

    assert t_tbl.height > s_tbl.height
    assert t_tbl.top < s_tbl.top
    # 두 경우 모두 아래 끝은 본문 영역 하단이다
    assert s_tbl.top + s_tbl.height == area_top + area_h
    assert t_tbl.top + t_tbl.height == area_top + area_h


def test_build_layout_mode_keeps_tables_inside_the_slide(tmp_path: Path):
    """어떤 길이가 와도 표 하단이 슬라이드를 넘지 않는다 — 이 작업의 목적이다."""
    prs = _build_layout(tmp_path, _screens_with_desc([LONG_DESC] * 20))
    height = int(prs.slide_height)
    for slide in prs.slides:
        for tbl in _tables_of(slide):
            assert tbl.top + tbl.height <= height


def test_build_layout_mode_short_details_keep_measured_heights(tmp_path: Path):
    """짧은 상세만 있으면 실측 고정 높이 그대로다 — 기존 산출물이 바뀌면 안 된다."""
    prs = _build_layout(tmp_path, _screens_with_desc(["짧게"] * 4))
    tbl = _tables_of(prs.slides[0])[0]
    assert [r.height for r in tbl.table.rows] == [382457, 268746, 496168, 268746]
    assert tbl.top == 5253244


def test_build_layout_mode_gives_each_page_its_own_height(tmp_path: Path):
    """장마다 자기 내용에 맞는 높이를 쓴다.

    전 장을 최악값으로 통일하면 긴 상세 하나가 모든 장의 표를 밀어 올려
    스크린샷을 눌러 버린다.
    """
    descs = ["짧게"] * 40
    descs[25] = LONG_DESC          # 2장에만 긴 항목이 있다
    prs = _build_layout(tmp_path, _screens_with_desc(descs))

    assert len(prs.slides) == 2
    first = _tables_of(prs.slides[0])[0]
    second = _tables_of(prs.slides[1])[0]
    # 1장은 짧은 상세뿐이라 실측 하한 그대로다
    assert [r.height for r in first.table.rows] == [382457, 268746, 496168, 268746]
    # 2장만 긴 항목 때문에 자란다
    assert second.height > first.height
    assert second.top < first.top


def test_build_layout_mode_bottom_aligns_every_page(tmp_path: Path):
    """장별 높이가 달라도 아래 끝은 모든 장에서 본문 영역 하단이다."""
    area_bottom = 337940 + 6331421
    descs = ["짧게"] * 40
    descs[25] = LONG_DESC
    prs = _build_layout(tmp_path, _screens_with_desc(descs))

    for slide in prs.slides:
        for tbl in _tables_of(slide):
            assert tbl.top + tbl.height == area_bottom


def test_build_layout_mode_uses_one_font_size_across_pages(tmp_path: Path):
    """높이는 장별이지만 글자 크기는 화면 안에서 통일한다.

    한 장만 작으면 장을 넘길 때 글자가 커졌다 작아졌다 한다.
    """
    prs = _build_layout(tmp_path, _screens_with_desc([LONG_DESC] * 40))

    assert len(prs.slides) == 2
    sizes = set()
    for slide in prs.slides:
        cell = _tables_of(slide)[0].table.cell(0, 1)
        sizes.add(cell.text_frame.paragraphs[0].runs[0].font.size)
    assert len(sizes) == 1


def test_build_layout_mode_aligns_row_heights_across_tables(tmp_path: Path):
    """표 다섯 개가 나란히 놓이므로 행 높이가 어긋나면 안 된다."""
    descs = ["짧게"] * 20
    descs[13] = LONG_DESC          # 슬롯 13 -> 표 3, 행 1
    prs = _build_layout(tmp_path, _screens_with_desc(descs))

    tables = _tables_of(prs.slides[0])
    assert len(tables) == 5
    base = [r.height for r in tables[0].table.rows]
    for t in tables[1:]:
        assert [r.height for r in t.table.rows] == base
    assert base[1] > 268746        # 행 1이 자랐다


def test_build_layout_mode_shrinks_font_when_image_slot_would_vanish(tmp_path: Path):
    """표가 이미지 자리를 다 먹으면 설명 글자를 낮춘다. 경고 코드는 늘리지 않는다."""
    from pptx.util import Pt

    prs = _build_layout(tmp_path, _screens_with_desc([LONG_DESC * 3] * 20))
    tbl = _tables_of(prs.slides[0])[0]
    run = tbl.table.cell(0, 1).text_frame.paragraphs[0].runs[0]

    assert run.font.size < Pt(7)
    assert run.font.size >= Pt(6)
    # 번호 칸은 그대로다
    num = tbl.table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert num.font.size == Pt(6.5)


def test_build_layout_mode_never_lets_tables_eat_the_image_slot(tmp_path: Path):
    """폰트를 낮춰서라도 이미지 자리를 1인치는 남긴다."""
    prs = _build_layout(tmp_path, _screens_with_desc([LONG_DESC * 3] * 20))
    area_top = 337940
    tbl = _tables_of(prs.slides[0])[0]
    assert tbl.top - area_top >= 914400


def _cli_args(tmp_path: Path) -> list[str]:
    """--out-file 없이 build.py를 부를 CLI 인자. screens.json의 meta.source는 s.xlsx다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    out_dir = tmp_path / "output"
    write_json(out_dir / "screens.json", _screens(3))
    write_json(out_dir / "mapping.json", _mapping(tpl))
    return ["--screens", str(out_dir / "screens.json"),
            "--mapping", str(out_dir / "mapping.json"),
            "--output", str(out_dir)]


def _work_cli_args(tmp_path: Path) -> list[str]:
    """중간 산출물이 .work/에 있는 새 구조. build.py는 --output만 받는다."""
    tpl = make_template_pptx(tmp_path / "t.pptx")
    out_dir = tmp_path / "output"
    write_json(out_dir / ".work" / "screens.json", _screens(3))
    write_json(out_dir / ".work" / "mapping.json", _mapping(tpl))
    return ["--output", str(out_dir)]


def test_cli_reads_work_dir_without_path_arguments(tmp_path: Path):
    assert build_cli(_work_cli_args(tmp_path)) == 0
    assert (tmp_path / "output" / "s.pptx").exists()


def test_cli_leaves_only_pptx_in_output_dir(tmp_path: Path):
    """결과물 폴더에 보이는 건 pptx뿐이다."""
    assert build_cli(_work_cli_args(tmp_path)) == 0

    out_dir = tmp_path / "output"
    assert sorted(p.name for p in out_dir.iterdir()) == [".work", "s.pptx"]
    assert not (out_dir / ".work" / "s.pptx").exists()


def test_cli_migrates_legacy_files_then_builds(tmp_path: Path):
    """구버전 폴더를 그대로 줘도 .work/로 올린 뒤 이어서 만든다."""
    args = _cli_args(tmp_path)[-2:]        # --output <디렉토리>만 남긴다
    assert args[0] == "--output"

    assert build_cli(args) == 0

    out_dir = tmp_path / "output"
    assert (out_dir / ".work" / "mapping.json").exists()
    assert (out_dir / ".work" / "screens.json").exists()
    assert (out_dir / "s.pptx").exists()


def test_cli_names_output_from_source_excel(tmp_path: Path):
    args = _cli_args(tmp_path)
    assert build_cli(args) == 0
    assert (tmp_path / "output" / "s.pptx").exists()


def test_cli_numbers_instead_of_overwriting(tmp_path: Path):
    """재실행이 직전 산출물을 덮어쓰면 결과를 비교할 수단이 없어진다."""
    args = _cli_args(tmp_path)
    assert build_cli(args) == 0
    assert build_cli(args) == 0
    assert build_cli(args) == 0

    out_dir = tmp_path / "output"
    assert sorted(p.name for p in out_dir.glob("s*.pptx")) == [
        "s.pptx", "s2.pptx", "s3.pptx",
    ]


def test_cli_out_file_is_taken_as_is(tmp_path: Path):
    """경로를 직접 지정했으면 번호를 붙이지 않는다 — 덮어쓰기가 의도한 동작이다."""
    target = tmp_path / "납품용.pptx"
    args = _cli_args(tmp_path) + ["--out-file", str(target)]
    assert build_cli(args) == 0
    assert build_cli(args) == 0

    assert target.exists()
    assert not (tmp_path / "납품용2.pptx").exists()
    assert not (tmp_path / "output" / "s.pptx").exists()
