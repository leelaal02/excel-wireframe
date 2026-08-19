from pathlib import Path

from fixtures import make_template_pptx
from pptx import Presentation
from pptx.util import Pt
from slide_fill import estimate_overflow, find_shape, set_cell_text, set_text


def test_find_shape_by_name(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    slide = prs.slides[0]
    assert find_shape(slide, "제목 13") is not None
    assert find_shape(slide, "없는도형") is None


def test_set_text_keeps_font_size(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    slide = prs.slides[0]
    shape = find_shape(slide, "제목 13")
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    set_text(shape, "이용기관 목록")

    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.text == "이용기관 목록"
    assert run.font.size == Pt(18)
    assert len(shape.text_frame.paragraphs) == 1


def test_set_cell_text_splits_newlines(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    cell = table.cell(0, 1)
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    set_cell_text(cell, "- 구분값 : 전체, Y, N\n- 디폴트 : 전체")

    paras = cell.text_frame.paragraphs
    assert len(paras) == 2
    assert paras[0].runs[0].text == "- 구분값 : 전체, Y, N"
    assert paras[1].runs[0].text == "- 디폴트 : 전체"
    assert paras[0].runs[0].font.size == Pt(9)
    assert paras[1].runs[0].font.size == Pt(9)


def test_set_cell_text_clears_previous_content(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    cell = table.cell(0, 1)
    set_cell_text(cell, "첫 줄\n둘째 줄")
    set_cell_text(cell, "짧게")
    assert cell.text == "짧게"


def test_set_cell_text_accepts_empty(tmp_path: Path):
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    cell = table.cell(0, 1)
    set_cell_text(cell, "")
    assert cell.text == ""


# 실측 설명 칸: 폭 1810509, 두 번째 행 높이 268746(7pt 두 줄)
DESC_W, ROW_H = 1810509, 268746


def test_estimate_overflow():
    assert estimate_overflow("짧은 글", DESC_W, ROW_H, 7.0) is False
    assert estimate_overflow("가" * 400, DESC_W, ROW_H, 7.0) is True


def test_estimate_overflow_uses_the_row_height():
    """행에 들어가는 줄 수를 넘어야 넘침이다. 40자는 두 줄이라 딱 맞는다."""
    assert estimate_overflow("가" * 40, DESC_W, ROW_H, 7.0) is False
    assert estimate_overflow("가" * 41, DESC_W, ROW_H, 7.0) is True


def test_estimate_overflow_relaxes_in_a_taller_row():
    """같은 글이라도 행이 높으면 넘치지 않는다. 세 번째 행은 네 줄을 담는다."""
    assert estimate_overflow("가" * 41, DESC_W, 496168, 7.0) is False


def test_estimate_overflow_relaxes_with_smaller_font():
    """글자를 낮추면 같은 행에 더 들어간다."""
    assert estimate_overflow("가" * 41, DESC_W, ROW_H, 7.0) is True
    assert estimate_overflow("가" * 41, DESC_W, ROW_H, 6.0) is False


def test_estimate_overflow_ignores_empty_text():
    assert estimate_overflow("", DESC_W, ROW_H, 7.0) is False
    assert estimate_overflow(None, DESC_W, ROW_H, 7.0) is False


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _add_date_field(shape, shown: str) -> None:
    """자동 날짜 필드를 심는다. 사용자 템플릿의 DATE placeholder가 이 꼴이다."""
    from lxml import etree

    p = shape.text_frame.paragraphs[0]._p
    for r in list(p.findall("{%s}r" % A_NS)):
        p.remove(r)
    fld = etree.SubElement(p, "{%s}fld" % A_NS)
    fld.set("id", "{7DE4667A-1AB8-4DCD-A2F5-D16009686B1A}")
    fld.set("type", "datetime1")
    t = etree.SubElement(fld, "{%s}t" % A_NS)
    t.text = shown


def test_set_text_replaces_auto_field(tmp_path: Path):
    """자동 날짜 필드가 있는 자리에 값을 쓰면 필드 값이 남아 이어 붙으면 안 된다.

    a:fld는 런이 아니라서 paragraph.runs에 잡히지 않는다. 그 사실을 모르고
    런만 갈아끼우면 '2026-05-292026-06-25'처럼 두 값이 붙어 나온다.
    """
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    shape = find_shape(prs.slides[0], "제목 13")
    _add_date_field(shape, "2026-05-29")
    assert shape.text_frame.text == "2026-05-29"

    set_text(shape, "2026-06-25")

    assert shape.text_frame.text == "2026-06-25"
    assert shape._element.find(".//{%s}fld" % A_NS) is None


def test_set_cell_text_keeps_paragraph_spacing_on_extra_lines(tmp_path: Path):
    """줄바꿈으로 늘어난 문단도 첫 문단의 줄간격을 물려받아야 한다.

    안 물려받으면 둘째 줄부터 폰트 기본 줄높이가 쓰여, 계산한 행 높이보다
    실제 렌더링이 커진다 — 표가 슬라이드를 넘는 원인이다.
    """
    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    cell = table.cell(0, 1)
    cell.text_frame.paragraphs[0].line_spacing = 0.95

    set_cell_text(cell, "첫 줄\n둘째 줄\n셋째 줄")

    paras = cell.text_frame.paragraphs
    assert len(paras) == 3
    assert all(p.line_spacing == 0.95 for p in paras)


# --- 미입력 표시 ---

def test_set_text_or_required_writes_the_value_when_there_is_one(tmp_path: Path):
    from slide_fill import set_text_or_required

    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    shape = find_shape(prs.slides[0], "제목 13")
    assert set_text_or_required(shape, "이용기관 목록") is True
    assert shape.text_frame.text == "이용기관 목록"


def test_set_text_or_required_marks_missing_values_in_bold_red(tmp_path: Path):
    """비워 두면 채워야 할 자리가 있다는 사실 자체가 산출물에서 사라진다."""
    from pptx.dml.color import RGBColor
    from slide_fill import INPUT_REQUIRED, set_text_or_required

    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    shape = find_shape(prs.slides[0], "제목 13")
    for empty in (None, "", "   "):
        assert set_text_or_required(shape, empty) is False
        assert shape.text_frame.text == INPUT_REQUIRED
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.bold is True
        assert run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)


def test_set_text_keeps_size_while_marking_required(tmp_path: Path):
    """강조는 굵기와 색만 바꾼다 — 크기·글꼴은 템플릿 값을 지킨다."""
    from slide_fill import set_text_or_required

    prs = Presentation(str(make_template_pptx(tmp_path / "t.pptx")))
    shape = find_shape(prs.slides[0], "제목 13")
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    set_text_or_required(shape, None)
    assert shape.text_frame.paragraphs[0].runs[0].font.size == Pt(11)


# --- 메타 표 위에 얹는 글자 자리 ---

def _meta_slide(tmp_path: Path):
    from default_template import DEFAULT_LAYOUT_NAME, build_default_template
    from slide_layout import add_meta_text_slots, find_layout

    prs = Presentation(str(build_default_template(tmp_path / "d.pptx")))
    layout = find_layout(prs, DEFAULT_LAYOUT_NAME)
    slide = prs.slides.add_slide(layout)
    made = add_meta_text_slots(slide, layout, ["메타표1", "메타표2"],
                               ["프로젝트명", "화면명", "ID", "알림여부"])
    return slide, made


def test_meta_slots_are_created_for_each_label(tmp_path: Path):
    from slide_layout import meta_slot_name

    slide, made = _meta_slide(tmp_path)
    assert sorted(made) == sorted(["프로젝트명", "화면명", "ID", "알림여부"])
    names = {s.name for s in slide.shapes}
    for label in made:
        assert meta_slot_name(label) in names
    # 원본 표를 복제하지 않았다 — 슬라이드에 표가 없어야 한다
    assert [s for s in slide.shapes if s.has_table] == []


def test_meta_slot_sits_exactly_on_the_cell_right_of_its_label(tmp_path: Path):
    from default_template import DEFAULT_LAYOUT_NAME, build_default_template
    from slide_layout import add_meta_text_slots, find_layout, meta_slot_name

    prs = Presentation(str(build_default_template(tmp_path / "d.pptx")))
    layout = find_layout(prs, DEFAULT_LAYOUT_NAME)
    slide = prs.slides.add_slide(layout)
    add_meta_text_slots(slide, layout, ["메타표1", "메타표2"], ["화면명"])

    frame = next(s for s in layout.shapes
                 if s.has_table and len(s.table.columns) == 18)
    table = frame.table
    col = [table.cell(0, i).text for i in range(len(table.columns))].index("화면명") + 1
    left = int(frame.left) + sum(int(table.columns[i].width) for i in range(col))

    shp = next(s for s in slide.shapes if s.name == meta_slot_name("화면명"))
    assert shp.left == left
    assert shp.top == frame.top
    assert shp.width == table.columns[col].width
    assert shp.height == table.rows[0].height


def test_meta_slot_carries_the_cells_own_text_format(tmp_path: Path):
    """칸이 품고 있던 서식을 그대로 써야 값을 넣어도 표가 원래 모습으로 남는다."""
    from pptx.util import Pt as _Pt
    from slide_fill import set_text
    from slide_layout import meta_slot_name

    slide, _ = _meta_slide(tmp_path)
    shp = next(s for s in slide.shapes if s.name == meta_slot_name("화면명"))
    set_text(shp, "이용기관 목록")

    run = shp.text_frame.paragraphs[0].runs[0]
    assert run.font.size == _Pt(6.5)
    assert run.font.bold is True


def test_meta_slot_is_transparent(tmp_path: Path):
    """원본 표가 비쳐 보여야 하므로 채움과 테두리가 없어야 한다."""
    from pptx.enum.dml import MSO_FILL
    from slide_layout import meta_slot_name

    slide, _ = _meta_slide(tmp_path)
    shp = next(s for s in slide.shapes if s.name == meta_slot_name("ID"))
    assert shp.fill.type == MSO_FILL.BACKGROUND
    assert shp.line.fill.type == MSO_FILL.BACKGROUND


def test_meta_slots_skip_labels_that_are_not_in_the_table(tmp_path: Path):
    from default_template import DEFAULT_LAYOUT_NAME, build_default_template
    from slide_layout import add_meta_text_slots, find_layout

    prs = Presentation(str(build_default_template(tmp_path / "d.pptx")))
    layout = find_layout(prs, DEFAULT_LAYOUT_NAME)
    slide = prs.slides.add_slide(layout)
    assert add_meta_text_slots(slide, layout, ["메타표1"], ["없는라벨"]) == []
