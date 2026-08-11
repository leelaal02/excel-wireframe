import subprocess
import sys
from pathlib import Path

import pytest
from common import read_json, write_json
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
SCRIPTS = ROOT / "skills" / "excel-wireframe" / "scripts"
SAMPLE_XLSX = ROOT / "짧은 버전.xlsx"
SAMPLE_PPTX = ROOT / "화면설계서_저작권_발행기관.정산처_발행기관관리_v1.0_20260427.pptx"

pytestmark = pytest.mark.skipif(
    not (SAMPLE_XLSX.exists() and SAMPLE_PPTX.exists()),
    reason="실제 샘플 파일이 없습니다",
)

MAPPING = {
    "version": 1,
    "excel": {
        "layout": "sheet-per-screen",
        "sheet_include": "^설계_",
        "screen_meta": {
            "cell": "A1",
            "pattern": r"화면설계서\s*-\s*(?P<id>\S+)\s*\((?P<name>.+)\)",
        },
        "detail": {
            "header_scan_column": "A",
            "header_marker": "No.",
            "columns": {"no": "A", "type": "B", "element": "C", "desc": "D", "pos": "E"},
        },
    },
    "template": {
        "file": str(SAMPLE_PPTX),
        "mode": "clone",
        "source_slide": 1,
        "shapes": {
            "title": "제목 13",
            "screen_id": "텍스트 개체 틀 14",
            "image": "그림 18",
            "detail_tables": ["표 7", "표 10", "표 11", "표 12", "표 15"],
        },
        "table_columns": {"no": 0, "text": 1},
    },
    "options": {
        "detail_text_source": "desc",
        "overflow": "split",
        "clear_unused_slots": True,
    },
}


def _run(script: str, *args: str):
    cmd = [sys.executable, str(SCRIPTS / script), *args]
    return subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", cwd=ROOT)


def test_sample_end_to_end(tmp_path: Path):
    out_dir = tmp_path / "output"
    write_json(out_dir / ".work" / "mapping.json", MAPPING)

    # 스킬이 실제로 돌리는 명령 그대로다 — 경로는 스크립트가 유도한다.
    r = _run("extract.py", "--excel", str(SAMPLE_XLSX), "--output", str(out_dir))
    assert r.returncode == 0, r.stderr

    screens = read_json(out_dir / ".work" / "screens.json")
    assert len(screens["screens"]) == 1
    scr = screens["screens"][0]
    assert scr["id"] == "B2BISMT1001"
    assert scr["name"] == "이용기관 목록"
    assert len(scr["details"]) == 16
    assert scr["images"] == ["images/B2BISMT1001.png"]
    # 중요 발견 5: extract가 만든 meta가 표지의 실제 작성일과 일치하는지
    # 사람이 수동으로만 확인해 오던 것을 코드화한다.
    assert screens["meta"]["작성일"] == "2026-06-11"

    # --out-file 없이 돌린다. 이름은 build.py가 원본 Excel 파일명에서 정한다.
    out = out_dir / (SAMPLE_XLSX.stem + ".pptx")
    r = _run("build.py", "--output", str(out_dir))
    assert r.returncode == 0, r.stderr
    assert out.exists()

    # 같은 명령을 다시 돌리면 덮어쓰지 않고 번호를 붙인다.
    r = _run("build.py", "--output", str(out_dir))
    assert r.returncode == 0, r.stderr
    assert (out_dir / (SAMPLE_XLSX.stem + "2.pptx")).exists()
    assert out.exists()

    # 결과물 폴더에서 눈에 보이는 건 pptx뿐이다.
    assert sorted(p.name for p in out_dir.iterdir()) == [
        ".work", SAMPLE_XLSX.stem + ".pptx", SAMPLE_XLSX.stem + "2.pptx",
    ]

    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert prs.slide_width == 9906000

    slide = prs.slides[0]
    title = next(s for s in slide.shapes if s.name == "제목 13")
    assert title.text_frame.text == "이용기관 목록"

    tables = sorted((s for s in slide.shapes if s.has_table), key=lambda s: s.left)
    assert len(tables) == 5
    assert tables[0].table.cell(0, 0).text == "1"
    # 상세 16건 → 표1~4(16슬롯)까지 꽉 차고 마지막 표(17~20번 슬롯)는 통째로 비어야 한다.
    # 마지막 행 한 칸만 확인하면 "마지막 행만 지우는" 버그도 통과해버리므로 4행 x 2열을 모두 본다.
    last_table = tables[4].table
    for row in range(4):
        assert last_table.cell(row, 0).text == ""
        assert last_table.cell(row, 1).text == ""

    pics = [s for s in slide.shapes if s.shape_type == 13]
    assert len(pics) == 1


def test_analyze_on_sample(tmp_path: Path):
    out = tmp_path / "structure-report.json"
    r = _run("analyze.py", "--excel", str(SAMPLE_XLSX),
             "--template", str(SAMPLE_PPTX), "--out", str(out))
    assert r.returncode == 0, r.stderr
    report = read_json(out)
    assert [s["name"] for s in report["excel"]["sheets"]] == [
        "표지", "설계_B2BISMT1001", "테스트_B2BISMT1001", "비교결과요약",
    ]
    assert report["suggestion"]["mode"] == "clone"
