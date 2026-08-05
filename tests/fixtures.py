# -*- coding: utf-8 -*-
"""테스트용 xlsx / pptx / png 생성기."""
from __future__ import annotations

from pathlib import Path

from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Emu

DETAIL_HEADER = ["No.", "요소타입", "요소명", "상세 설명", "위치"]


def make_png(path: Path, size=(400, 300), color=(200, 220, 255)) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    PILImage.new("RGB", size, color).save(path)
    return path


def make_sheet_per_screen_xlsx(path: Path, screens: list[dict]) -> Path:
    """1시트 = 1화면 양식. 화면이 아닌 시트도 섞어 시트 필터를 검증할 수 있게 한다."""
    wb = Workbook()
    cover = wb.active
    cover.title = "표지"
    cover["A1"] = "화면설계서"

    for scr in screens:
        ws = wb.create_sheet("설계_%s" % scr["id"])
        ws["A1"] = "화면설계서 - %s (%s)" % (scr["id"], scr["name"])
        ws["A3"] = "[웹 스크린샷 (SoM 뱃지)]"
        if scr.get("image"):
            img_path = path.parent / ("_fx_%s.png" % scr["id"])
            # 화면마다 색을 달리한다. 내용이 같으면 저장 과정에서 하나로 합쳐져
            # 이미지 귀속 테스트가 무의미해질 수 있다.
            tint = 40 * (ord(scr["id"][-1]) % 5)
            make_png(img_path, color=(200, 220 - tint, 255 - tint))
            ws.add_image(XLImage(str(img_path)), "A4")
        header_row = 28
        for i, name in enumerate(DETAIL_HEADER):
            ws.cell(row=header_row, column=i + 1, value=name)
        for j, d in enumerate(scr["details"]):
            r = header_row + 1 + j
            ws.cell(row=r, column=1, value=d["no"])
            ws.cell(row=r, column=2, value=d["type"])
            ws.cell(row=r, column=3, value=d["element"])
            ws.cell(row=r, column=4, value=d["desc"])
            ws.cell(row=r, column=5, value=d["pos"])

    test_ws = wb.create_sheet("테스트_무시대상")
    test_ws["A1"] = "이 시트는 무시되어야 한다"

    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(path)
    return path


def make_table_xlsx(path: Path) -> Path:
    """1행 = 1화면 양식. 화면ID가 빈 후속 행은 직전 화면의 상세 행이다."""
    wb = Workbook()
    ws = wb.active
    ws.title = "화면목록"
    ws["A1"] = "화면 정의서"
    headers = ["화면ID", "화면명", "화면설명", "No.", "요소명", "상세 설명"]
    for i, h in enumerate(headers):
        ws.cell(row=3, column=i + 1, value=h)

    rows = [
        ["SCR001", "이용기관 목록", "기관을 조회한다", "1", "[등록] 버튼", "등록 팝업을 연다"],
        ["", "", "", "2", "[삭제] 버튼", "선택 항목을 삭제한다"],
        ["SCR002", "이용기관 상세", "상세를 본다", "1", "[저장] 버튼", "변경 내용을 저장한다"],
    ]
    for j, row in enumerate(rows):
        for i, v in enumerate(row):
            ws.cell(row=4 + j, column=i + 1, value=v)

    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(path)
    return path


def make_template_pptx(path: Path) -> Path:
    """실제 샘플과 같은 도형 이름·슬라이드 크기를 가진 *예시 슬라이드형* 템플릿.

    기본 템플릿은 슬라이드가 0장이므로, 여기서 layout 모드의 슬라이드 생성
    경로로 한 장을 만들어 예시가 채워진 실제 템플릿을 흉내 낸다 —
    clone 판정(suggest_mode)과 텍스트 교체 테스트들이 그런 파일을 가정한다.

    표 개수·행 수는 실측 조합(5표 × 4행 = 20슬롯)으로 고정한다. 예전엔 인자로
    받았지만 그 인자를 쓰던 테스트가 사라져 아무도 넘기지 않는다 — 쓰이지 않는
    설정 손잡이는 남겨 두면 지켜지는지 아무도 확인하지 않는 약속이 된다.
    """
    table_count, rows_per_table = 5, 4
    from default_template import (
        DEFAULT_LAYOUT_NAME,
        DEFAULT_SHAPE_NAMES,
        build_default_template,
        default_template_mapping,
    )
    from slide_layout import (
        add_detail_table,
        add_image_anchor,
        find_layout,
        inherit_placeholders,
        name_placeholders,
        split_content_area,
        DEFAULT_CONTENT_AREA,
    )
    from common import Warnings

    build_default_template(path)
    prs = Presentation(str(path))
    layout = find_layout(prs, DEFAULT_LAYOUT_NAME)
    tpl = default_template_mapping(path, table_count, rows_per_table)

    slide = prs.slides.add_slide(layout)
    inherit_placeholders(slide, layout)
    name_placeholders(slide, tpl["placeholders"], tpl["shapes"],
                      Warnings(), "fixture")
    image_box, table_boxes = split_content_area(
        DEFAULT_CONTENT_AREA, table_count, rows_per_table)
    add_image_anchor(slide, image_box, DEFAULT_SHAPE_NAMES["image"])
    for i, box in enumerate(table_boxes):
        add_detail_table(slide, box, rows_per_table, "상세표%d" % (i + 1))

    # 실제 샘플의 도형 이름과 예시 텍스트를 흉내 낸다
    rename = {"제목": "제목 13", "화면ID": "텍스트 개체 틀 14", "화면이미지": "그림 18"}
    example_text = {"제목": "화면명", "화면ID": "SCR000", "작성일": "2024-01-01"}
    for shp in slide.shapes:
        if shp.name in example_text and shp.has_text_frame:
            shp.text_frame.text = example_text[shp.name]
        if shp.name in rename:
            shp.name = rename[shp.name]
        elif shp.name.startswith("상세표"):
            shp.name = "표 %d" % (6 + int(shp.name[len("상세표"):]))

    # 상세 표에 번호와 예시 문구를 넣는다 — clone 모드는 이것을 덮어써야 한다.
    # 슬롯마다 *다른* 값을 써야 한다: 20칸에 같은 문자열을 넣으면 "안 쓴 슬롯을
    # 비웠는가"를 보는 테스트가 아무 일도 안 일어난 경우와 구별하지 못한다.
    # 번호 열까지 채우는 이유도 같다.
    # add_detail_table이 만들어 둔 빈 런에 글자만 넣는다 — text_frame.text에
    # 대입하면 런이 새로 생기면서 실측 서식(6.5/7pt, 맑은 고딕)이 날아간다.
    tables = sorted((s for s in slide.shapes if s.has_table), key=lambda s: s.left)
    for i, shp in enumerate(tables):
        for r in range(len(shp.table.rows)):
            n = i * rows_per_table + r + 1
            shp.table.cell(r, 0).text_frame.paragraphs[0].runs[0].text = str(n)
            shp.table.cell(r, 1).text_frame.paragraphs[0].runs[0].text = (
                "예시 설명 %d" % n)

    prs.save(str(path))
    return path


def make_empty_layout_pptx(path: Path) -> Path:
    """빈 레이아웃형 템플릿 — clone 모드로 오판되면 안 된다."""
    prs = Presentation()
    prs.slide_width = Emu(9906000)
    prs.slide_height = Emu(6858000)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
