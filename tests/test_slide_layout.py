# -*- coding: utf-8 -*-
import pytest
from common import Warnings
from text_metrics import CELL_MARGIN
from pptx import Presentation
from pptx.util import Pt
from slide_layout import (
    DEFAULT_CONTENT_AREA,
    P_NS,
    add_detail_table,
    add_image_anchor,
    drop_empty_placeholders,
    find_layout,
    inherit_placeholders,
    name_placeholders,
    split_content_area,
)


def _prs():
    """Title and Content 레이아웃(TITLE/OBJECT/DATE/FOOTER/SLIDE_NUMBER)을 가진 프레젠테이션."""
    return Presentation()


def test_find_layout_by_name():
    prs = _prs()
    lay = find_layout(prs, "Title and Content")
    assert lay.name == "Title and Content"


def test_find_layout_by_index():
    prs = _prs()
    assert find_layout(prs, 1).name == prs.slide_layouts[1].name


def test_find_layout_raises_for_unknown_name():
    prs = _prs()
    with pytest.raises(ValueError) as exc:
        find_layout(prs, "없는레이아웃")
    assert "없는레이아웃" in str(exc.value)


def test_inherit_placeholders_adds_date_and_footer():
    """python-pptx는 date/footer/slidenumber를 복제하지 않는다 — 우리가 채운다."""
    prs = _prs()
    lay = find_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(lay)
    before = {ph.placeholder_format.idx for ph in slide.placeholders}
    assert 10 not in before

    added = inherit_placeholders(slide, lay)

    after = {ph.placeholder_format.idx for ph in slide.placeholders}
    assert {10, 11, 12} <= after
    assert set(added) == after - before


def test_inherit_placeholders_assigns_unique_shape_ids():
    """레이아웃 id를 그대로 쓰면 사용자 템플릿에서 중복 id가 생겨 PPT가 복구를 요구한다."""
    prs = _prs()
    lay = find_layout(prs, "Title and Content")

    # 레이아웃의 날짜 placeholder(idx=10)에, 슬라이드가 실제로 쓸 id(제목 placeholder의
    # id)를 일부러 심어 충돌을 재현한다. 이렇게 강제하지 않으면 python-pptx 기본
    # 템플릿에서는 번호가 우연히 겹치지 않아 회귀를 못 잡는다.
    date_ph = next(ph for ph in lay.placeholders if ph.placeholder_format.idx == 10)
    date_ph._element.find(".//{%s}cNvPr" % P_NS).set("id", "2")

    slide = prs.slides.add_slide(lay)
    assert slide.placeholders[0].shape_id == 2  # 제목이 이미 id=2를 쓰고 있음을 확인

    inherit_placeholders(slide, lay)

    ids = [shape.shape_id for shape in slide.shapes]
    assert len(ids) == len(set(ids)), ids


def test_inherit_placeholders_is_idempotent():
    prs = _prs()
    lay = find_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(lay)
    inherit_placeholders(slide, lay)
    n = len(list(slide.placeholders))
    assert inherit_placeholders(slide, lay) == []
    assert len(list(slide.placeholders)) == n


def test_name_placeholders_renames_by_idx():
    prs = _prs()
    lay = find_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(lay)
    inherit_placeholders(slide, lay)

    warns = Warnings()
    name_placeholders(
        slide,
        {"title": 0, "screen_id": 1, "작성일": 10},
        {"title": "제목", "screen_id": "화면ID"},
        warns,
        "SCR001",
    )

    names = [s.name for s in slide.shapes]
    assert "제목" in names
    assert "화면ID" in names
    assert "작성일" in names  # shapes에 없으면 키를 그대로 이름으로 쓴다
    assert len(warns) == 0


def test_name_placeholders_warns_for_missing_idx():
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    warns = Warnings()
    name_placeholders(slide, {"title": 0, "없는것": 99}, {}, warns, "SCR001")
    items = warns.to_list()
    assert len(items) == 1
    assert items[0]["code"] == "shape-not-found"
    assert "99" in items[0]["message"]


def test_drop_empty_placeholders_removes_blank_ones():
    prs = _prs()
    lay = find_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(lay)
    slide.placeholders[0].text_frame.text = "제목 있음"

    removed = drop_empty_placeholders(slide)

    assert removed >= 1
    remaining = [ph.placeholder_format.idx for ph in slide.placeholders]
    assert 0 in remaining
    assert 1 not in remaining


def test_drop_empty_placeholders_keeps_field_placeholders():
    """쪽번호는 자동 번호 필드라 텍스트가 비어 보여도 지우면 안 된다."""
    prs = _prs()
    lay = find_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(lay)
    inherit_placeholders(slide, lay)
    drop_empty_placeholders(slide)
    assert 12 in [ph.placeholder_format.idx for ph in slide.placeholders]


def test_split_content_area_puts_tables_at_the_bottom():
    area = DEFAULT_CONTENT_AREA
    image_box, table_boxes = split_content_area(area, 5, 4)

    assert len(table_boxes) == 5
    # 표는 본문 영역 아래쪽에 붙는다
    table_top = table_boxes[0][1]
    table_height = table_boxes[0][3]
    assert table_top + table_height == area[1] + area[3]
    # 이미지는 위쪽 나머지를 전부 쓴다
    assert image_box[1] == area[1]
    assert image_box[1] + image_box[3] == table_top


def test_split_content_area_matches_measured_geometry():
    """실측 조합(표 5개 × 4행)에서 표 상단이 원본과 같은 자리에 온다."""
    _, table_boxes = split_content_area(DEFAULT_CONTENT_AREA, 5, 4)
    assert table_boxes[0][1] == 5253244
    assert table_boxes[0][3] == 382457 + 268746 + 496168 + 268746


def test_split_content_area_divides_width_evenly():
    area = (0, 0, 10000, 4000000)
    _, boxes = split_content_area(area, 5, 4)
    assert [b[0] for b in boxes] == [0, 2000, 4000, 6000, 8000]
    assert all(b[2] == 2000 for b in boxes)


def test_split_content_area_last_table_absorbs_the_width_remainder():
    """실측 DEFAULT_CONTENT_AREA는 폭이 5로 나누어떨어지지 않는다(9957099 // 5 == 1991419,
    나머지 4). 그 나머지를 버리면 마지막 표 오른쪽 끝이 본문 영역보다 짧아진다."""
    area = DEFAULT_CONTENT_AREA
    area_left, _, area_width, _ = area
    _, boxes = split_content_area(area, 5, 4)

    # 표 폭의 합은 본문 영역 폭과 정확히 같다 (버려지는 나머지가 없다)
    assert sum(b[2] for b in boxes) == area_width
    # 마지막 표의 오른쪽 끝은 본문 영역의 오른쪽 끝과 정확히 맞는다
    last_left, _, last_width, _ = boxes[-1]
    assert last_left + last_width == area_left + area_width
    # 표들은 서로 겹치지도, 벌어지지도 않는다
    for prev, cur in zip(boxes, boxes[1:]):
        assert cur[0] == prev[0] + prev[2]


def test_split_content_area_grows_tables_with_rows():
    _, four = split_content_area(DEFAULT_CONTENT_AREA, 5, 4)
    _, six = split_content_area(DEFAULT_CONTENT_AREA, 5, 6)
    assert six[0][3] > four[0][3]
    # 마지막 행높이를 반복한다
    assert six[0][3] == four[0][3] + 268746 * 2


def test_split_content_area_raises_when_image_slot_too_small():
    with pytest.raises(ValueError) as exc:
        split_content_area(DEFAULT_CONTENT_AREA, 5, 40)
    assert "이미지" in str(exc.value)


def test_add_detail_table_applies_measured_formatting():
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    frame = add_detail_table(slide, (0, 0, 1971135, 1416117), 4, "상세표1")

    assert frame.name == "상세표1"
    table = frame.table
    assert len(table.rows) == 4
    assert len(table.columns) == 2
    assert [r.height for r in table.rows] == [382457, 268746, 496168, 268746]
    # 열폭은 실측 비율을 표 폭에 맞춰 나눈다
    assert sum(c.width for c in table.columns) == 1971135
    assert table.columns[0].width < table.columns[1].width

    num, txt = table.cell(0, 0), table.cell(0, 1)
    assert num.margin_left == 18000
    assert txt.margin_left == CELL_MARGIN
    assert txt.margin_bottom == 0
    assert num.text_frame.paragraphs[0].runs[0].font.size == Pt(6.5)
    assert num.text_frame.paragraphs[0].runs[0].font.bold is True
    assert txt.text_frame.paragraphs[0].runs[0].font.size == Pt(7)
    assert txt.text_frame.paragraphs[0].runs[0].font.name == "맑은 고딕"


def test_add_detail_table_starts_empty():
    """빈 Excel에서 예시 문구가 산출물에 찍히면 안 된다."""
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    table = add_detail_table(slide, (0, 0, 1971135, 1416117), 4, "상세표1").table
    for r in range(4):
        assert table.cell(r, 0).text == ""
        assert table.cell(r, 1).text == ""


def test_split_content_area_grows_tables_upward_with_row_heights():
    """행 높이를 넘기면 표 상단만 올라가고 아래 끝은 본문 영역 하단에 그대로 붙는다."""
    area = DEFAULT_CONTENT_AREA
    _, base = split_content_area(area, 5, 4)
    tall = [600000, 600000, 600000, 600000]
    image_box, boxes = split_content_area(area, 5, 4, tall)

    assert boxes[0][3] == sum(tall)
    assert boxes[0][3] > base[0][3]
    # 아래 끝은 그대로
    assert boxes[0][1] + boxes[0][3] == area[1] + area[3]
    # 상단이 올라간 만큼 이미지 자리가 줄어든다
    assert boxes[0][1] < base[0][1]
    assert image_box[1] == area[1]
    assert image_box[1] + image_box[3] == boxes[0][1]


def test_split_content_area_ignores_none_row_heights():
    """안 넘기면 실측 고정 높이를 그대로 쓴다 — 기존 동작이 바뀌면 안 된다."""
    assert split_content_area(DEFAULT_CONTENT_AREA, 5, 4) == \
        split_content_area(DEFAULT_CONTENT_AREA, 5, 4, None)


def test_split_content_area_raises_when_row_heights_eat_the_image_slot():
    with pytest.raises(ValueError) as exc:
        split_content_area(DEFAULT_CONTENT_AREA, 5, 4, [1600000] * 4)
    assert "이미지" in str(exc.value)


def test_add_detail_table_uses_given_row_heights():
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    heights = [500000, 400000, 600000, 300000]
    frame = add_detail_table(slide, (0, 0, 1971135, sum(heights)), 4, "상세표1",
                             row_heights=heights)
    assert [r.height for r in frame.table.rows] == heights


def test_add_detail_table_shrinks_only_the_text_column_font():
    """번호 칸은 한두 글자라 넘칠 일이 없다. 설명 칸만 낮춘다."""
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    table = add_detail_table(slide, (0, 0, 1971135, 1416117), 4, "상세표1",
                             size_pt=6.0).table

    num, txt = table.cell(0, 0), table.cell(0, 1)
    assert txt.text_frame.paragraphs[0].runs[0].font.size == Pt(6)
    assert num.text_frame.paragraphs[0].runs[0].font.size == Pt(6.5)
    # 나머지 실측 서식은 그대로다
    assert txt.text_frame.paragraphs[0].runs[0].font.name == "맑은 고딕"
    assert txt.margin_bottom == 0


def test_add_image_anchor_uses_the_box():
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    shp = add_image_anchor(slide, (100, 200, 3000, 4000), "화면이미지")
    assert (shp.left, shp.top, shp.width, shp.height) == (100, 200, 3000, 4000)
    assert shp.name == "화면이미지"


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _new_table(rows=4):
    prs = _prs()
    slide = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    return add_detail_table(slide, (0, 0, 1971135, 1416117), rows, "상세표1")


def test_add_detail_table_uses_the_plain_table_style():
    """python-pptx 기본 스타일은 파란 머리글과 줄무늬가 들어간다.

    원본 화면설계서의 상세표는 스타일 없이 셀마다 테두리를 그린 꼴이다.
    기본값을 그대로 두면 원본과 전혀 다른 표가 나온다.
    """
    frame = _new_table()
    tbl = frame.table._tbl
    style = tbl.find("{%s}tblPr" % A_NS).find("{%s}tableStyleId" % A_NS)
    assert style is not None
    assert style.text == "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def test_add_detail_table_draws_cell_borders():
    """스타일이 없으므로 테두리는 셀마다 직접 그려야 보인다."""
    table = _new_table().table
    for r in range(4):
        for c in (0, 1):
            pr = table.cell(r, c)._tc.find("{%s}tcPr" % A_NS)
            assert pr is not None, (r, c)
            for edge in ("lnL", "lnR", "lnT", "lnB"):
                ln = pr.find("{%s}%s" % (A_NS, edge))
                assert ln is not None, (r, c, edge)
                assert ln.get("w") == "3175", (r, c, edge)
                clr = ln.find(".//{%s}schemeClr" % A_NS)
                assert clr is not None and clr.get("val") == "tx1", (r, c, edge)


def test_add_detail_table_shades_the_number_column_only():
    """번호 칸에만 배경이 깔린다. 내용 칸은 비운다."""
    table = _new_table().table
    for r in range(4):
        no_pr = table.cell(r, 0)._tc.find("{%s}tcPr" % A_NS)
        txt_pr = table.cell(r, 1)._tc.find("{%s}tcPr" % A_NS)
        no_fill = no_pr.find("{%s}solidFill" % A_NS)
        assert no_fill is not None, r
        assert no_fill.find("{%s}schemeClr" % A_NS).get("val") == "tx2"
        assert txt_pr.find("{%s}solidFill" % A_NS) is None, r


def test_add_detail_table_borders_follow_row_count():
    """행 수를 바꿔도 모든 셀이 테두리를 갖는다."""
    table = _new_table(rows=6).table
    assert len(table.rows) == 6
    for r in range(6):
        pr = table.cell(r, 1)._tc.find("{%s}tcPr" % A_NS)
        assert pr.find("{%s}lnB" % A_NS) is not None, r


def test_add_detail_table_sets_line_spacing_on_text_cells():
    """계산에 쓰는 줄간격이 셀 서식에도 들어가야 렌더링과 계산이 맞는다."""
    from text_metrics import LINE_SPACING_RATIO
    table = _new_table().table
    para = table.cell(0, 1).text_frame.paragraphs[0]
    assert para.line_spacing == LINE_SPACING_RATIO


def test_add_detail_table_text_margin_comes_from_the_shared_constant():
    """여백이 세 곳에 흩어져 있으면 계산과 산출물이 어긋난다."""
    table = _new_table().table
    txt = table.cell(0, 1)
    assert txt.margin_left == CELL_MARGIN
    assert txt.margin_top == CELL_MARGIN


# --- 본문 높이에 비례한 행높이 ---

def test_ratio_row_heights_matches_measurements_at_the_reference_height():
    """기준 본문(ROW_HEIGHTS를 잰 높이)에서는 실측값 그대로여야 한다."""
    from slide_layout import MEASURED_CONTENT_HEIGHT, ROW_HEIGHTS, ratio_row_heights
    assert ratio_row_heights(MEASURED_CONTENT_HEIGHT, 4) == ROW_HEIGHTS


def test_ratio_row_heights_scales_with_the_content_area():
    from slide_layout import MEASURED_CONTENT_HEIGHT, ROW_HEIGHTS, ratio_row_heights
    half = ratio_row_heights(MEASURED_CONTENT_HEIGHT // 2, 4)
    assert sum(half) == sum(ROW_HEIGHTS) // 2 or abs(
        sum(half) - sum(ROW_HEIGHTS) // 2) <= 4      # 정수 절삭 오차
    for h, base in zip(half, ROW_HEIGHTS):
        assert abs(h - base // 2) <= 2


def test_split_content_area_puts_tables_at_the_same_ratio_in_any_area():
    """레이아웃·슬라이드 크기가 달라도 표가 본문에서 차지하는 몫이 같다."""
    from slide_layout import split_content_area

    small = split_content_area((0, 0, 9144000, 3000000), 5, 4)[1][0]
    big = split_content_area((0, 0, 9144000, 6000000), 5, 4)[1][0]
    assert abs(small[3] / 3000000 - big[3] / 6000000) < 0.001
