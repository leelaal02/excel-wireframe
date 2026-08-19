# -*- coding: utf-8 -*-
"""사용자가 자기 조직 템플릿을 기본으로 두는 경로를 검증한다.

이 스킬은 코드째 배포되므로 남의 저작권이 박힌 템플릿을 저장소에 담을 수 없다.
대신 설치본에 개인 설정 파일을 두고, --template 없이 실행할 때 그걸 쓴다.
"""
from pathlib import Path

import pytest
from common import write_json
from default_template import (
    DEFAULT_LAYOUT_NAME,
    build_default_template,
    default_meta_table_labels,
)
from user_default import load_user_default, user_default_mapping


def _install(tmp_path: Path, **over) -> Path:
    """설치본 흉내: user-default.json + assets/템플릿."""
    base = tmp_path / "skill"
    (base / "assets").mkdir(parents=True)
    build_default_template(base / "assets" / "default-template.pptx")
    cfg = {
        "template": "assets/default-template.pptx",
        "mode": "layout",
        "layout": DEFAULT_LAYOUT_NAME,
        "placeholders": {"screen_id": 1, "문서제목": 11, "쪽번호": 12},
        "shapes": {
            "screen_id": "화면ID",
            "image": "화면이미지",
            "문서제목": "문서제목",
            "쪽번호": "쪽번호",
            "detail_tables": ["상세표1", "상세표2", "상세표3", "상세표4",
                              "상세표5"],
        },
        "meta_table": {
            "tables": ["메타표1", "메타표2"],
            "labels": default_meta_table_labels(),
        },
        "detail_tables": {"count": 5, "rows": 4},
        "table_columns": {"no": 0, "text": 1},
        "content_area": [-12319, 337940, 9957099, 6331421],
    }
    cfg.update(over)
    write_json(base / "user-default.json", cfg)
    return base


def test_load_returns_none_without_config(tmp_path: Path):
    """설정이 없으면 지금까지의 동작(기본 템플릿 생성)을 유지해야 한다."""
    base = tmp_path / "skill"
    base.mkdir()
    assert load_user_default(base) is None


def test_load_resolves_template_path_against_skill_dir(tmp_path: Path):
    base = _install(tmp_path)
    cfg = load_user_default(base)
    assert cfg is not None
    assert cfg["template_path"].is_absolute()
    assert cfg["template_path"].exists()
    assert cfg["template_path"].name == "default-template.pptx"


def test_load_accepts_absolute_template_path(tmp_path: Path):
    base = _install(tmp_path)
    abs_path = base / "assets" / "default-template.pptx"
    write_json(base / "user-default.json",
               {"template": str(abs_path), "layout": "화면"})
    cfg = load_user_default(base)
    assert cfg["template_path"] == abs_path


def test_load_raises_when_template_missing(tmp_path: Path):
    """설정은 있는데 파일이 없으면 조용히 기본 템플릿으로 떨어지면 안 된다.

    왜 내 템플릿이 안 쓰이는지 사용자가 알 수 없게 된다.
    """
    base = tmp_path / "skill"
    base.mkdir()
    write_json(base / "user-default.json",
               {"template": "assets/없는파일.pptx", "layout": "화면"})
    with pytest.raises(ValueError) as exc:
        load_user_default(base)
    assert "없는파일.pptx" in str(exc.value)


def test_load_raises_without_layout(tmp_path: Path):
    """layout 모드인데 레이아웃 이름이 없으면 build가 인덱스 0을 집는다."""
    base = _install(tmp_path)
    cfg = {"template": "assets/default-template.pptx"}
    write_json(base / "user-default.json", cfg)
    with pytest.raises(ValueError) as exc:
        load_user_default(base)
    assert "layout" in str(exc.value)


def test_mapping_is_a_usable_template_section(tmp_path: Path):
    base = _install(tmp_path)
    cfg = load_user_default(base)
    m = user_default_mapping(cfg)

    assert m["mode"] == "layout"
    assert m["layout"] == DEFAULT_LAYOUT_NAME
    assert Path(m["file"]) == cfg["template_path"]
    assert m["placeholders"] == {"screen_id": 1, "문서제목": 11, "쪽번호": 12}
    # meta_table도 그대로 넘어가야 상단 메타 표가 채워진다
    assert m["meta_table"]["tables"] == ["메타표1", "메타표2"]
    assert m["meta_table"]["labels"]["화면명"] == "title"
    assert m["detail_tables"] == {"count": 5, "rows": 4}
    assert m["content_area"] == [-12319, 337940, 9957099, 6331421]
    # 내부용 키가 새어 나가면 안 된다
    assert "template" not in m
    assert "template_path" not in m


def test_mapping_builds_end_to_end(tmp_path: Path):
    """제안 매핑을 그대로 build에 넘겨 슬라이드가 나와야 한다."""
    from build import build
    from common import Warnings

    base = _install(tmp_path)
    mapping = {
        "version": 1,
        "excel": {"layout": "sheet-per-screen"},
        "template": user_default_mapping(load_user_default(base)),
        "options": {"detail_text_source": "desc", "clear_unused_slots": True},
    }
    screens = {
        "meta": {},
        "screens": [{"id": "SCR001", "name": "목록", "images": [],
                     "fields": {}, "details": [{"no": "1", "desc": "설명"}]}],
    }
    out = tmp_path / "out.pptx"
    report = build(screens, mapping, tmp_path, out, Warnings())

    assert report["slides"] == 1
    from pptx import Presentation
    slide = Presentation(str(out)).slides[0]
    from slide_layout import meta_slot_name

    by_name = {s.name: s for s in slide.shapes}
    assert by_name["화면ID"].text_frame.text == "SCR001"
    # 화면명은 상단 메타 표의 '화면명' 칸 자리에 얹힌 글자로 들어간다
    assert by_name[meta_slot_name("화면명")].text_frame.text == "목록"
    # 원본 메타 표는 레이아웃에 그대로 있고 슬라이드에는 상세표만 있다
    assert len([s for s in slide.shapes if s.has_table]) == 5
