# -*- coding: utf-8 -*-
"""extract.py가 만든 표지 meta가 build.py까지 실제로 이어지는지 확인한다.

test_build.py는 손으로 만든 meta 딕셔너리로 build만 검증하고, test_extract.py는
extract가 meta를 만드는 것만 검증한다 — 실제 워크북 하나로 두 단계를 이어붙여
확인하는 테스트가 없었다(중요 발견 5). 실제 샘플 파일에 의존하는
test_sample_e2e.py는 클론 직후 파일이 없으면 조용히 skip되므로, 픽스처만으로
동작해 클론 직후에도 살아남는 테스트가 따로 필요하다.
"""
from pathlib import Path

from build import build
from common import Warnings, read_json, write_json
from extract import main as extract_main
from fixtures import make_sheet_per_screen_xlsx, make_template_pptx
from openpyxl import load_workbook
from pptx import Presentation

SPEC = [
    {"id": "SCR001", "name": "이용기관 목록", "image": False,
     "details": [{"no": "1", "type": "버튼", "element": "[등록]",
                  "desc": "등록한다", "pos": "상단"}]},
]


def _mapping(template: Path) -> dict:
    return {
        "version": 1,
        "excel": {
            "layout": "sheet-per-screen",
            "sheet_include": "^설계_",
            "screen_meta": {
                "cell": "A1",
                "pattern": r"화면설계서\s*-\s*(?P<id>\S+)\s*\((?P<name>.+)\)",
            },
            "detail": {
                "header_scan_column": "A",
                "header_marker": "No.",
                "columns": {"no": "A", "type": "B", "element": "C", "desc": "D", "pos": "E"},
            },
        },
        "template": {
            "file": str(template),
            "mode": "clone",
            "source_slide": 0,
            "shapes": {
                "title": "제목 13",
                "screen_id": "텍스트 개체 틀 14",
                "image": "그림 18",
                "문서제목": "문서제목",
                "작성일": "작성일",
                "detail_tables": ["표 7", "표 8", "표 9", "표 10", "표 11"],
            },
            "table_columns": {"no": 0, "text": 1},
        },
        "options": {"detail_text_source": "desc"},
    }


def test_extract_then_build_carries_cover_meta_into_shapes(tmp_path: Path):
    xlsx = make_sheet_per_screen_xlsx(tmp_path / "s.xlsx", SPEC)
    wb = load_workbook(xlsx)
    cover = wb["표지"]
    cover["C7"] = "프로젝트명"
    cover["E7"] = "통합관리시스템"
    cover["C8"] = "작성일"
    cover["E8"] = "2026-06-11"
    wb.save(xlsx)

    tpl = make_template_pptx(tmp_path / "t.pptx")
    work = tmp_path / "work"
    mapping_path = work / "mapping.json"
    write_json(mapping_path, _mapping(tpl))

    code = extract_main(["--excel", str(xlsx), "--mapping", str(mapping_path),
                          "--work", str(work)])
    assert code == 0

    screens_data = read_json(work / "screens.json")
    assert screens_data["meta"]["작성일"] == "2026-06-11"
    assert screens_data["meta"]["프로젝트명"] == "통합관리시스템"

    out = work / "output" / "화면설계서.pptx"
    report = build(screens_data, read_json(mapping_path), work, out, Warnings())
    assert out.exists()
    assert report["failed"] == []

    slide = Presentation(str(out)).slides[0]
    date_shape = next(s for s in slide.shapes if s.name == "작성일")
    assert date_shape.text_frame.text == "2026-06-11"
    title_shape = next(s for s in slide.shapes if s.name == "문서제목")
    assert title_shape.text_frame.text == "화면설계서"
