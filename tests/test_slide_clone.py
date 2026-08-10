from pathlib import Path

from fixtures import make_template_pptx, make_png
from pptx import Presentation
from pptx.util import Emu, Inches
from slide_clone import clone_slide


def test_clone_with_embedded_picture_preserves_image_bytes(tmp_path: Path):
    """Verify rId remapping by cloning a slide with an embedded picture and hyperlink.

    This test proves the rId remapping mechanism works. Without proper remapping
    of r:embed and r:id in shape XML, pictures' relationships would point to
    wrong rIds on the cloned slide, causing images to vanish or corrupt the file.

    The test creates both an embedded picture (r:embed) and an external hyperlink (r:id)
    on the text shape to force rId shifts between slides.
    """
    # Create a presentation with a slide containing an embedded picture
    prs = Presentation()
    prs.slide_width = Emu(9906000)
    prs.slide_height = Emu(6858000)
    blank_slide_layout = prs.slide_layouts[6]  # blank layout
    src = prs.slides.add_slide(blank_slide_layout)

    # Add a picture to the source slide
    pic_path = make_png(tmp_path / "test.png", size=(200, 150))
    with open(pic_path, 'rb') as f:
        original_blob = f.read()

    src.shapes.add_picture(str(pic_path), Inches(1), Inches(1), width=Inches(2))

    # Add a text box with a hyperlink to create an external relationship
    # This shifts the image rId so the new slide will have different rIds
    txBox = src.shapes.add_textbox(Inches(0.5), Inches(4), Inches(3), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Test Link"

    # Add hyperlink to the run (creates an external relationship with r:id)
    hlink = run.hyperlink
    hlink.address = "https://example.com"

    # Clone the slide - this should remap both r:embed (picture) and r:id (hyperlink)
    new = clone_slide(prs, src)

    # Save and reload to verify the clone survives serialization
    out = tmp_path / "out.pptx"
    prs.save(out)

    chk = Presentation(str(out))
    cloned_slide = chk.slides[-1]

    # Assert the cloned slide has exactly one picture
    pics = [s for s in cloned_slide.shapes if s.shape_type == 13]  # 13 = PICTURE
    assert len(pics) == 1, f"Expected 1 picture, got {len(pics)}"

    # Assert the picture blob is byte-identical to the original
    cloned_blob = pics[0].image.blob
    assert cloned_blob == original_blob, (
        f"Picture blob mismatch after clone: original {len(original_blob)} bytes, "
        f"cloned {len(cloned_blob)} bytes"
    )


def test_clone_copies_all_shapes(tmp_path: Path):
    pptx = make_template_pptx(tmp_path / "t.pptx")
    prs = Presentation(str(pptx))
    src = prs.slides[0]
    new = clone_slide(prs, src)
    assert len(new.shapes) == len(src.shapes)
    assert [s.name for s in new.shapes] == [s.name for s in src.shapes]


def test_clone_survives_save_and_reload(tmp_path: Path):
    pptx = make_template_pptx(tmp_path / "t.pptx")
    prs = Presentation(str(pptx))
    clone_slide(prs, prs.slides[0])
    out = tmp_path / "out.pptx"
    prs.save(out)

    chk = Presentation(str(out))
    assert len(chk.slides) == 2
    last = chk.slides[-1]
    tables = [s for s in last.shapes if s.has_table]
    assert len(tables) == 5
    assert tables[0].table.cell(0, 0).text == "1"


def test_clone_does_not_share_notes_slide(tmp_path: Path):
    """노트가 달린 슬라이드를 복제해도 사본은 노트를 물려받지 않는다.

    노트 슬라이드는 슬라이드 하나에만 붙는 파트다. 복제해서 여러 슬라이드가
    같은 notesSlide를 가리키면 python-pptx는 저장하지만 PowerPoint는 파일을
    손상으로 판정해 열지 못한다.
    """
    pptx = make_template_pptx(tmp_path / "t.pptx")
    prs = Presentation(str(pptx))
    src = prs.slides[0]
    src.notes_slide.notes_text_frame.text = "템플릿 작성자의 메모"

    clone_slide(prs, src)
    clone_slide(prs, src)
    out = tmp_path / "out.pptx"
    prs.save(out)

    chk = Presentation(str(out))
    for slide in list(chk.slides)[1:]:
        assert not slide.has_notes_slide

    notes = [
        rel.target_part
        for slide in chk.slides
        for rel in slide.part.rels.values()
        if not rel.is_external and rel.reltype.endswith("/notesSlide")
    ]
    assert len(notes) == len(set(id(p) for p in notes))


def test_clone_drops_shape_tags(tmp_path: Path):
    """도형에 달린 태그(custDataLst)는 복제하지 않는다.

    태그 파트도 슬라이드 하나에만 붙는다. 사본들이 원본의 태그 파트를 나눠
    가지면 PowerPoint가 파일을 열지 못한다. 참조(custDataLst)만 남겨도
    없는 rId를 가리켜 깨지므로 파트와 참조를 함께 버려야 한다.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    pptx = make_template_pptx(tmp_path / "t.pptx")
    prs = Presentation(str(pptx))
    src = prs.slides[0]

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rid = src.part.rels.get_or_add_ext_rel(RT.TAGS, "tag1.xml")
    shape = src.shapes[0]
    cust = shape._element.makeelement("{%s}custDataLst" % p_ns, {})
    tags = shape._element.makeelement("{%s}tags" % p_ns, {"{%s}id" % r_ns: rid})
    cust.append(tags)
    shape._element.append(cust)

    clone_slide(prs, src)
    out = tmp_path / "out.pptx"
    prs.save(out)

    chk = Presentation(str(out))
    cloned = list(chk.slides)[-1]
    assert not cloned._element.findall(".//{%s}custDataLst" % p_ns)
    assert not [r for r in cloned.part.rels.values() if r.reltype == RT.TAGS]


def test_clone_twice_produces_three_slides(tmp_path: Path):
    pptx = make_template_pptx(tmp_path / "t.pptx")
    prs = Presentation(str(pptx))
    src = prs.slides[0]
    clone_slide(prs, src)
    clone_slide(prs, src)
    out = tmp_path / "out.pptx"
    prs.save(out)
    assert len(Presentation(str(out)).slides) == 3
