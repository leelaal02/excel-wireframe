# -*- coding: utf-8 -*-
"""레이아웃을 골라 슬라이드를 만들고, 그 레이아웃이 물려주는 자리에 값을 채운다.

디자인은 레이아웃이 담당한다. 이 모듈은 자리를 찾아 이름을 붙이고, 레이아웃에
자리가 없는 두 가지(화면 이미지, 상세표)만 본문 영역 안에 만든다.
"""
from __future__ import annotations

import copy

from lxml.etree import SubElement

from common import EMU_PER_INCH
from text_metrics import CELL_MARGIN, LINE_SPACING_RATIO
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def find_layout(prs, spec):
    """이름 또는 인덱스로 레이아웃을 찾는다.

    마스터가 여러 개인 템플릿이 흔하므로 전부 훑는다. python-pptx의
    prs.slide_layouts는 첫 번째 마스터만 보여줘서 그것만으로는 부족하다.
    """
    layouts = [lay for master in prs.slide_masters for lay in master.slide_layouts]

    if isinstance(spec, int):
        if 0 <= spec < len(layouts):
            return layouts[spec]
        raise ValueError(
            "레이아웃 인덱스 %d가 범위를 벗어납니다 (레이아웃 %d개)"
            % (spec, len(layouts))
        )

    matches = [lay for lay in layouts if lay.name == spec]
    if not matches:
        raise ValueError(
            "레이아웃 '%s'을(를) 찾지 못했습니다. 있는 레이아웃: %s"
            % (spec, ", ".join(lay.name for lay in layouts))
        )
    if len(matches) > 1:
        # 마스터가 여러 개면 같은 이름이 겹칠 수 있다. 멈출 일은 아니지만
        # 어느 것을 골랐는지는 알려야 한다.
        import sys
        print("경고: 레이아웃 '%s'이(가) %d개 있어 첫 번째를 씁니다"
              % (spec, len(matches)), file=sys.stderr)
    return matches[0]


def append_shape_with_new_id(spTree, element):
    """다른 파트에서 복사해 온 도형 XML을 붙이면서 cNvPr/@id를 새로 매긴다.

    cNvPr/@id는 한 파트(슬라이드 하나, 레이아웃 하나) 안에서 유일해야 한다
    (OOXML 규칙). 다른 파트에서 deepcopy로 들고 온 요소의 id는 목적지가 이미
    쓰고 있는 id와 겹칠 수 있다 — 레이아웃과 슬라이드가 둘 다 2부터 번호를
    매기므로 오히려 겹치는 쪽이 흔하다. 겹치면 PowerPoint가 파일을 열 때
    복구를 요구한다. 도형을 하나 붙일 때마다 다시 계산해야 이식하는 도형끼리도
    겹치지 않는다. python-pptx 자신도 clone_placeholder에서 새 id를 매길 때
    같은 방식(파트 전체 최댓값+1)을 쓴다.

    같은 규칙을 두 군데(placeholder 상속, 기본 템플릿의 껍데기 이식)에서
    따로 구현했다가 같은 결함을 두 번 고쳤다. 이식 코드를 새로 쓸 일이 생기면
    이 함수를 쓴다. 단 slide_clone.clone_slide은 예외다 — 그쪽은 목적지
    트리를 통째로 비우고 원본 슬라이드를 그대로 옮기므로 id가 겹칠 상대가
    애초에 없고, 오히려 id를 바꾸면 rId 재매핑과 어긋난다.
    """
    element.find(".//{%s}cNvPr" % P_NS).set("id", str(spTree.max_shape_id + 1))
    spTree.append(element)
    return element


def inherit_placeholders(slide, layout) -> list[int]:
    """레이아웃에 있으나 슬라이드에 없는 placeholder를 복제한다.

    python-pptx의 add_slide는 date/footer/slidenumber 계열을 복제하지 않는다.
    PowerPoint 관례상 그 셋은 마스터 설정으로 표시되기 때문인데, 우리는 거기에
    값을 써야 하므로 직접 옮긴다.

    레이아웃 XML의 cNvPr/@id는 그대로 들고 오면 안 된다 —
    append_shape_with_new_id가 그 이유와 처리를 담고 있다.
    """
    have = {ph.placeholder_format.idx for ph in slide.placeholders}
    spTree = slide.shapes._spTree
    added: list[int] = []
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx in have:
            continue
        append_shape_with_new_id(spTree, copy.deepcopy(ph._element))
        added.append(idx)
    return added


META_SLOT_PREFIX = "메타:"


def meta_slot_name(label: str) -> str:
    """라벨에 대응하는 메타 값 자리의 도형 이름."""
    return META_SLOT_PREFIX + label


def _cell_boxes(frame):
    """표 안 각 셀의 (행, 열) → (left, top, width, height)."""
    left0, top0 = int(frame.left or 0), int(frame.top or 0)
    table = frame.table
    widths = [int(c.width or 0) for c in table.columns]
    heights = [int(r.height or 0) for r in table.rows]
    boxes = {}
    top = top0
    for r, h in enumerate(heights):
        left = left0
        for c, w in enumerate(widths):
            boxes[(r, c)] = (left, top, w, h)
            left += w
        top += h
    return boxes


def add_meta_text_slots(slide, layout, table_names, labels) -> list[str]:
    """레이아웃 메타 표의 라벨 오른쪽 칸 자리에 빈 글자 자리를 얹는다.

    레이아웃의 표는 배경으로 비칠 뿐 슬라이드의 도형이 아니라 값을 쓸 수 없다
    (PowerPoint는 placeholder만 물려준다). 표를 통째로 복제해 덮는 방법도 있지만
    그러면 표가 두 겹이 되어 원본을 클릭할 수 없고 테두리가 두 번 그려진다.
    **원본 표는 그대로 두고 값이 들어갈 칸 자리에 글자만 올린다.**

    글자 서식은 그 칸이 품고 있던 것(런의 rPr, 없으면 endParaRPr)과 문단 속성을
    그대로 복사한다 — 크기·굵기·색·글꼴·정렬이 원본과 같아야 값을 넣어도 표가
    원래 모습으로 남는다.

    만든 자리의 라벨 목록을 돌려준다.
    """
    wanted = set(labels or ())
    names = list(table_names or ())
    tables = [shp for shp in layout.shapes if shp.has_table]
    tables.sort(key=lambda s: (int(s.top or 0), int(s.left or 0)))
    if names:
        tables = tables[:len(names)]

    made: list[str] = []
    for frame in tables:
        table = frame.table
        boxes = _cell_boxes(frame)
        col_count = len(table.columns)
        for r in range(len(table.rows)):
            for c in range(col_count - 1):
                label = table.cell(r, c).text.strip()
                if label not in wanted or label in made:
                    continue
                _add_text_slot(slide, boxes[(r, c + 1)],
                               meta_slot_name(label), table.cell(r, c + 1))
                made.append(label)
    return made


def _add_text_slot(slide, box, name: str, source_cell):
    """칸 자리에 글자만 담을 투명한 텍스트 상자를 만든다."""
    left, top, width, height = box
    shp = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    shp.name = name
    shp.fill.background()
    shp.line.fill.background()

    tf = shp.text_frame
    # 원본 셀은 넘치는 글자를 옆으로 흘린다(horzOverflow="overflow"). 줄바꿈을
    # 켜면 한 줄짜리 칸에서 글이 아래로 삐져나가므로 같은 동작을 따른다.
    tf.word_wrap = False
    tf.margin_left = source_cell.margin_left
    tf.margin_right = source_cell.margin_right
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    src_p = source_cell.text_frame.paragraphs[0]._p
    p = tf.paragraphs[0]._p
    pPr = src_p.find("{%s}pPr" % A_NS)
    if pPr is not None:
        p.insert(0, copy.deepcopy(pPr))
    # 칸이 품고 있는 글자 서식을 그대로 옮긴다. 빈 칸은 그것을 endParaRPr에
    # 담고 있다 — PowerPoint가 "다음에 칠 글자의 서식"으로 저장해 두는 자리다.
    src_rPr = src_p.find(".//{%s}rPr" % A_NS)
    if src_rPr is None:
        src_rPr = src_p.find("{%s}endParaRPr" % A_NS)
    if src_rPr is not None:
        end = p.find("{%s}endParaRPr" % A_NS)
        if end is not None:
            p.remove(end)
        keep = copy.deepcopy(src_rPr)
        keep.tag = "{%s}endParaRPr" % A_NS
        p.append(keep)
    return shp


def name_placeholders(slide, placeholders_cfg, shapes_cfg, warns, screen_id) -> None:
    """placeholder에 mapping이 정한 이름을 붙인다.

    add_slide가 주는 이름은 'Title 1', 'Content Placeholder 2'처럼 그때그때
    달라진다. verify.py는 template.shapes의 *이름*으로 도형을 찾으므로,
    이름을 고정해 두지 않으면 검증이 느슨한 경로로 떨어진다.
    """
    by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    for key, idx in (placeholders_cfg or {}).items():
        ph = by_idx.get(int(idx))
        if ph is None:
            warns.add(screen_id, "shape-not-found",
                      "레이아웃에 placeholder idx=%s('%s')가 없습니다" % (idx, key))
            continue
        name = shapes_cfg.get(key, key)
        if isinstance(name, str):
            ph.name = name


def has_auto_field(shape) -> bool:
    """쪽번호처럼 자동 필드를 담은 도형인가.

    이런 자리는 PowerPoint가 값을 그린다. 글자를 쓰면 필드가 사라지므로
    채우기 대상에서 뺀다.
    """
    return shape._element.find(".//{%s}fld" % A_NS) is not None


_has_field = has_auto_field


def drop_empty_placeholders(slide) -> int:
    """값이 없는 placeholder를 지운다.

    남겨 두면 PowerPoint가 '제목을 입력하십시오' 프롬프트를 그려서 산출물에
    빈 안내 문구가 보인다. 자동 필드(쪽번호)는 텍스트가 비어 보여도 남긴다.
    """
    removed = 0
    for ph in list(slide.placeholders):
        if _has_field(ph):
            continue
        if ph.has_text_frame and ph.text_frame.text.strip():
            continue
        ph._element.getparent().remove(ph._element)
        removed += 1
    return removed


# --- 실측 상수 (원본 화면설계서 화면 페이지) ---
DEFAULT_CONTENT_AREA = (-12319, 337940, 9957099, 6331421)
ROW_HEIGHTS = [382457, 268746, 496168, 268746]
# ROW_HEIGHTS를 잰 본문 영역의 높이. 행높이를 이 값 대비 비율로 환산해 다른
# 크기의 슬라이드에도 같은 모양으로 옮긴다 — DEFAULT_CONTENT_AREA의 높이와 같은
# 값이어야 원본 화면설계서와 똑같은 표가 나온다.
MEASURED_CONTENT_HEIGHT = DEFAULT_CONTENT_AREA[3]
COL_WIDTH_RATIO = (160215, 1810920)
MIN_IMAGE_HEIGHT_EMU = EMU_PER_INCH  # 1인치
TEXT_CELL_MARGIN = CELL_MARGIN       # 설명 칸 여백. 계산과 서식이 같은 값을 쓴다
DETAIL_FONT_PT = 7.0                 # 설명 칸 실측 글자 크기
# 표를 늘려도 이미지 자리가 모자랄 때만 이 순서로 낮춘다. 6pt 밑으로는 안 간다 —
# 인쇄하면 읽히지 않아 설계서 구실을 못 한다.
DETAIL_FONT_STEPS = (7.0, 6.5, 6.0)

SLOT_BORDER = RGBColor(0xBB, 0xBB, 0xBB)
SLOT_FILL = RGBColor(0xF5, 0xF6, 0xF8)

# --- 표 서식 (원본 화면설계서의 상세표에서 그대로 잰 값) ---
# python-pptx가 새 표에 붙이는 기본 스타일은 파란 머리글과 줄무늬가 들어가서
# 원본과 전혀 다른 표가 나온다. 원본은 "스타일 없음"에 셀마다 테두리를 직접
# 그린 꼴이라, 스타일을 지우고 테두리와 음영을 우리가 그려야 같아진다.
TABLE_STYLE_PLAIN = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid
BORDER_WIDTH = 3175
BORDER_SCHEME = ("tx1", 65000, 35000)   # 테두리 회색
NO_COL_SCHEME = ("tx2", 40000, 60000)   # 번호 칸 배경
_EDGES = ("lnL", "lnR", "lnT", "lnB")


def measured_row_heights(rows_per_table: int) -> list[int]:
    """실측 행높이를 쓰되, 4행을 넘으면 마지막 값을 반복한다.

    내용 기반으로 행을 늘릴 때 이 값이 하한이 된다 — 짧은 상세만 있는 화면은
    원본 화면설계서와 똑같은 높이로 나와야 한다.
    """
    if rows_per_table <= len(ROW_HEIGHTS):
        return ROW_HEIGHTS[:rows_per_table]
    tail = [ROW_HEIGHTS[-1]] * (rows_per_table - len(ROW_HEIGHTS))
    return ROW_HEIGHTS + tail


def ratio_row_heights(area_height: int, rows_per_table: int) -> list[int]:
    """본문 영역 높이에 비례한 행높이. 사진·표 자리를 고정하는 기준이다.

    실측 행높이를 잰 본문(`MEASURED_CONTENT_HEIGHT`)과의 비로 환산한다. 그래서
    슬라이드 크기나 레이아웃이 달라져도 표가 본문에서 차지하는 몫이 같고, 상세가
    길든 짧든 모든 화면에서 사진과 표가 같은 자리에 온다 — 내용에 따라 표를 늘리면
    화면마다 사진 크기가 달라진다.

    기준 본문(6264697)에서는 실측값과 정확히 같다.
    """
    scale = area_height / MEASURED_CONTENT_HEIGHT
    return [int(h * scale) for h in measured_row_heights(rows_per_table)]


def detail_text_width(area_width: int, table_count: int) -> int:
    """설명 칸에서 글자가 실제로 놓이는 폭. 줄 수 계산의 기준이다.

    add_detail_table의 열 분할·여백과 같은 식을 써야 계산과 산출물이 어긋나지
    않는다. 마지막 표는 나머지를 흡수해 몇 EMU 넓지만, 한 글자에 못 미치는
    차이라 첫 표 기준으로 잡는다.
    """
    table_w = area_width // table_count
    no_w = table_w * COL_WIDTH_RATIO[0] // sum(COL_WIDTH_RATIO)
    return table_w - no_w - 2 * TEXT_CELL_MARGIN


def split_content_area(area, table_count: int, rows_per_table: int,
                       row_heights=None):
    """본문 영역을 이미지 자리와 상세표 자리로 나눈다.

    표는 아래쪽에 붙이고 폭을 균등 분할한다. 원본의 표 간격은 0.01인치라
    사실상 붙어 있으므로 간격을 두지 않는다. 이미지는 위쪽 나머지를 전부 쓴다.

    row_heights를 주면 그 높이로 표를 잡는다. 표의 아래 끝은 그대로 두고 상단만
    올라가며, 이미지 자리가 그만큼 줄어든다. 안 주면 본문 높이에 비례한 고정
    높이(`ratio_row_heights`)를 쓴다 — 모든 화면에서 사진과 표를 같은 자리에
    두려면 이 경로를 타야 한다.
    """
    left, top, width, height = area
    heights = (list(row_heights) if row_heights
               else ratio_row_heights(height, rows_per_table))
    table_h = sum(heights)
    table_top = top + height - table_h
    image_h = table_top - top

    if image_h < MIN_IMAGE_HEIGHT_EMU:
        if row_heights:
            raise ValueError(
                "상세 텍스트에 필요한 표 높이가 %.2fin이라 이미지 자리가 "
                "%.2fin로 줄어듭니다(최소 %.2fin 필요)."
                % (table_h / EMU_PER_INCH, image_h / EMU_PER_INCH,
                   MIN_IMAGE_HEIGHT_EMU / EMU_PER_INCH)
            )
        max_rows = 0
        while sum(ratio_row_heights(height, max_rows + 1)) <= height - MIN_IMAGE_HEIGHT_EMU:
            max_rows += 1
        raise ValueError(
            "rows_per_table=%d면 이미지 자리 높이가 %.2fin로 너무 작아집니다"
            "(최소 %.2fin 필요). 이 content_area 높이(%.2fin)에서는 "
            "rows_per_table을 %d 이하로 쓰세요."
            % (rows_per_table, image_h / EMU_PER_INCH,
               MIN_IMAGE_HEIGHT_EMU / EMU_PER_INCH, height / EMU_PER_INCH,
               max_rows)
        )

    table_w = width // table_count
    table_boxes = [
        (left + table_w * i, table_top, table_w, table_h)
        for i in range(table_count - 1)
    ]
    # 나머지(정수 나눗셈 몫에서 버려지는 폭)는 마지막 표에 몰아준다 — 그래야
    # 마지막 표의 오른쪽 끝이 본문 영역의 오른쪽 끝과 정확히 맞는다.
    last_left = left + table_w * (table_count - 1)
    last_width = width - table_w * (table_count - 1)
    table_boxes.append((last_left, table_top, last_width, table_h))
    return (left, top, width, image_h), table_boxes


def add_image_anchor(slide, box, name: str):
    """이미지가 들어갈 자리 사각형. place_image가 이것을 지우고 그림으로 바꾼다.

    이미지가 없는 화면에서는 이 사각형이 그대로 남아 '여기에 스크린샷' 자리로
    보인다 — clone 모드에서 템플릿의 이미지 자리 도형이 남는 것과 같은 동작이다.
    """
    left, top, width, height = box
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                 Emu(width), Emu(height))
    shp.name = name
    shp.fill.solid()
    shp.fill.fore_color.rgb = SLOT_FILL
    shp.line.color.rgb = SLOT_BORDER
    shp.text_frame.text = ""
    return shp


def scheme_fill(parent, tag: str, scheme):
    """schemeClr 채움 요소를 parent 아래에 만든다.

    scheme은 `(val, lumMod, lumOff)`이고 밝기 보정은 None으로 뺄 수 있다 —
    메타 표의 값 칸(`bg1` lumMod만)처럼 한쪽만 쓰는 색이 있다.
    """
    val, lum_mod, lum_off = scheme
    fill = SubElement(parent, "{%s}%s" % (A_NS, tag))
    clr = SubElement(fill, "{%s}schemeClr" % A_NS)
    clr.set("val", val)
    if lum_mod is not None:
        SubElement(clr, "{%s}lumMod" % A_NS).set("val", str(lum_mod))
    if lum_off is not None:
        SubElement(clr, "{%s}lumOff" % A_NS).set("val", str(lum_off))
    return fill



def _set_table_style(table, style_id: str) -> None:
    """표 스타일을 갈아끼운다. tableStyleId는 tblPr의 마지막 자식이다."""
    tbl_pr = table._tbl.find("{%s}tblPr" % A_NS)
    if tbl_pr is None:
        tbl_pr = SubElement(table._tbl, "{%s}tblPr" % A_NS)
    for old in tbl_pr.findall("{%s}tableStyleId" % A_NS):
        tbl_pr.remove(old)
    SubElement(tbl_pr, "{%s}tableStyleId" % A_NS).text = style_id


def draw_cell_edges(cell, fill=None, border=BORDER_SCHEME,
                    width: int = BORDER_WIDTH) -> None:
    """셀에 4방향 테두리를 그리고, fill을 주면 배경까지 칠한다.

    tcPr의 자식 순서는 스키마가 정한다 — 테두리(lnL/lnR/lnT/lnB)가 먼저이고
    채움이 그 뒤다. 순서를 어기면 PowerPoint가 파일을 열지 못한다.
    python-pptx에 셀 테두리 API가 없어 XML을 직접 쓴다.

    상세표와 상단 메타 표가 테두리 두께·색이 달라 인자로 뺐다.
    """
    tc = cell._tc
    pr = tc.find("{%s}tcPr" % A_NS)
    if pr is None:
        pr = SubElement(tc, "{%s}tcPr" % A_NS)
    for tag in _EDGES + ("solidFill",):
        for old in pr.findall("{%s}%s" % (A_NS, tag)):
            pr.remove(old)

    for edge in _EDGES:
        ln = SubElement(pr, "{%s}%s" % (A_NS, edge))
        ln.set("w", str(width))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        scheme_fill(ln, "solidFill", border)
        SubElement(ln, "{%s}prstDash" % A_NS).set("val", "solid")
        SubElement(ln, "{%s}round" % A_NS)

    if fill is not None:
        scheme_fill(pr, "solidFill", fill)


def _draw_cell_edges(cell, shaded: bool) -> None:
    """상세표용 짧은 이름. 번호 칸이면 배경을 칠한다."""
    draw_cell_edges(cell, NO_COL_SCHEME if shaded else None)


def _format_cell(cell, size_pt, bold, align, anchor, margin, margin_bottom,
                 font=None, line_spacing=None):
    cell.margin_left = Emu(margin)
    cell.margin_right = Emu(margin)
    cell.margin_top = Emu(margin)
    cell.margin_bottom = Emu(margin_bottom)
    if anchor is not None:
        cell.vertical_anchor = anchor
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    # 줄간격을 명시하면 text_metrics의 LINE_SPACING과 렌더링이 맞는다.
    # 안 넣으면 폰트 기본 줄높이가 쓰여 계산보다 행이 커진다.
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = ""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if font:
        run.font.name = font


def add_detail_table(slide, box, rows_per_table: int, name: str,
                     row_heights=None, size_pt: float = 7.0):
    """상세표 하나를 만들고 실측 서식을 적용한다. 셀은 비운 채로 둔다.

    row_heights를 주면 그 높이로 행을 잡는다(split_content_area와 같은 값을
    넘겨야 표가 자리에 정확히 들어간다). size_pt는 설명 칸 글자 크기다 —
    표를 늘려도 이미지 자리가 모자랄 때만 실측값 7.0에서 낮춘다. 번호 칸은
    한두 글자라 넘칠 일이 없어 6.5pt 그대로 둔다.
    """
    left, top, width, height = box
    frame = slide.shapes.add_table(rows_per_table, 2, Emu(left), Emu(top),
                                   Emu(width), Emu(height))
    frame.name = name
    table = frame.table

    ratio_total = sum(COL_WIDTH_RATIO)
    no_w = width * COL_WIDTH_RATIO[0] // ratio_total
    table.columns[0].width = Emu(no_w)
    table.columns[1].width = Emu(width - no_w)
    heights = (list(row_heights) if row_heights
               else measured_row_heights(rows_per_table))
    for ri, rh in enumerate(heights):
        table.rows[ri].height = Emu(rh)

    _set_table_style(table, TABLE_STYLE_PLAIN)

    for r in range(rows_per_table):
        _format_cell(table.cell(r, 0), 6.5, True, PP_ALIGN.CENTER,
                     MSO_ANCHOR.MIDDLE, 18000, 18000)
        _format_cell(table.cell(r, 1), size_pt, False, PP_ALIGN.LEFT, None,
                     TEXT_CELL_MARGIN, 0, font="맑은 고딕",
                     line_spacing=LINE_SPACING_RATIO)
        # 여백·정렬을 먼저 넣고 테두리를 그린다. _format_cell이 margin을 쓰면
        # python-pptx가 tcPr을 새로 만들 수 있어, 테두리를 먼저 그리면 지워진다.
        _draw_cell_edges(table.cell(r, 0), shaded=True)
        _draw_cell_edges(table.cell(r, 1), shaded=False)
    return frame
