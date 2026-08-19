# -*- coding: utf-8 -*-
"""슬라이드 복제와 삭제. python-pptx에 둘 다 API가 없어 직접 구현한다.

관계(rels)를 새 슬라이드에 재등록할 때 rId를 지정할 수 없고 자동 할당되므로,
복제한 XML 안의 r:embed / r:id 등을 새 rId로 치환해야 한다. 이 재매핑을 빼면
그림이 사라지거나 파일이 손상된다.
"""
from __future__ import annotations

import copy

from pptx.opc.constants import RELATIONSHIP_TYPE as RT

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_ATTRS = [
    "{%s}%s" % (R_NS, n)
    for n in ("id", "embed", "link", "pict", "dm", "lo", "qs", "cs")
]
# 슬라이드 하나에만 붙을 수 있는 파트는 물려주지 않는다. 여러 슬라이드가 같은
# 파트를 가리키면 python-pptx는 그대로 저장하지만 PowerPoint는 파일을 열지
# 못한다("파일이 손상되어 읽을 수 없습니다"). 이미지처럼 여러 슬라이드가
# 나눠 쓰는 것이 정상인 파트와는 다르다.
#
# - 레이아웃·마스터: add_slide가 이미 이어 준다.
# - 노트 슬라이드: 그 파트가 슬라이드를 역참조하므로 사본이 나눠 가질 수 없다.
# - 태그(p:tags): 도형의 custDataLst가 가리키는 사용자 정의 데이터다. 화면에
#   보이지 않고 산출물에 필요하지도 않으므로, 참조하는 custDataLst째로 지운다.
#
# 실제 사례: 슬라이드 노트와 태그 10개가 달린 기획서 템플릿을 clone 모드로
# 14장 복제한 산출물이 PowerPoint에서 열리지 않았다.
SKIP = {RT.SLIDE_LAYOUT, RT.SLIDE_MASTER, RT.NOTES_SLIDE, RT.TAGS}


def _drop_cust_data(element) -> None:
    """복제한 도형에서 custDataLst(태그 참조)를 걷어낸다.

    태그 파트를 물려주지 않기로 했으니 참조도 함께 지워야 한다. 참조만 남으면
    없는 rId를 가리키게 되어 그 자체로 파일이 깨진다.
    """
    for node in list(element.iter("{%s}custDataLst" % P_NS)):
        node.getparent().remove(node)


def clone_slide(prs, src):
    """src 슬라이드를 프레젠테이션 끝에 복제하고 새 슬라이드를 반환한다."""
    new = prs.slides.add_slide(src.slide_layout)

    # 레이아웃이 자동 삽입한 플레이스홀더를 걷어낸다. 원본 도형만 남겨야
    # 템플릿 모양이 정확히 재현된다.
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    rid_map: dict[str, str] = {}
    for old_rid, rel in src.part.rels.items():
        if rel.reltype in SKIP:
            continue
        if rel.is_external:
            new_rid = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new.part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[old_rid] = new_rid

    spTree = new.shapes._spTree
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        _drop_cust_data(el)
        for node in el.iter():
            for attr in R_ATTRS:
                v = node.get(attr)
                if v is not None and v in rid_map:
                    node.set(attr, rid_map[v])
        spTree.append(el)

    return new


def drop_slide(prs, slide) -> None:
    """프레젠테이션에서 슬라이드를 제거한다.

    sldIdLst 항목과 관계를 함께 지운다 — 한쪽만 지우면 파일이 깨진다.
    복제가 모두 끝난 뒤에 원본을 지워야 한다 — 먼저 지우면 복제 소스가 사라진다.
    """
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        if prs.part.rels[sld_id.rId].target_part is slide.part:
            prs.part.drop_rel(sld_id.rId)
            xml_slides.remove(sld_id)
            return
