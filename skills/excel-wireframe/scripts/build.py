# -*- coding: utf-8 -*-
"""3단계: screens.json과 템플릿으로 PPT를 만든다.

screens.json은 추출 단계가 남긴 중간 산출물이고, 이 모듈은 그 파일만 읽는다 —
Excel은 전혀 모른다. mapping.json에서도 template / options 섹션만 읽는다.
덕분에 Excel 픽스처 없이 빌드 로직을 검증할 수 있고, 이미지 추출이 느린 추출
단계를 다시 돌리지 않고 생성만 반복할 수 있다.
"""
from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path

from common import (
    Warnings,
    migrate_legacy_work,
    migration_notice,
    output_stem,
    read_json,
    resolve_output_path,
    resolve_template_path,
    setup_stdio,
    work_dir,
)
from pptx import Presentation
from image_split import split_for_box
from slide_clone import clone_slide
from slide_fill import (
    collect_tables,
    count_slots,
    fill_slots,
    find_shape,
    place_image,
    set_text,
)
from slide_layout import (
    DEFAULT_CONTENT_AREA,
    DETAIL_FONT_PT,
    DETAIL_FONT_STEPS,
    MIN_IMAGE_HEIGHT_EMU,
    add_detail_table,
    add_image_anchor,
    detail_text_width,
    drop_empty_placeholders,
    find_layout,
    inherit_placeholders,
    measured_row_heights,
    name_placeholders,
    split_content_area,
)
from text_metrics import plan_row_heights
from verify import verify_output


def chunk_details(details: list[dict], slot_count: int) -> list[list[dict]]:
    """슬롯 총량을 넘는 상세를 다음 슬라이드 분량으로 쪼갠다."""
    if not details:
        return [[]]
    if slot_count <= 0:
        return [details]
    return [
        details[i : i + slot_count] for i in range(0, len(details), slot_count)
    ]


def page_title(name: str, index: int, total: int) -> str:
    if total <= 1:
        return name
    return "%s (%d/%d)" % (name, index + 1, total)


def _drop_slide(prs, slide) -> None:
    """프레젠테이션에서 슬라이드를 제거한다.

    python-pptx에 삭제 API가 없어 sldIdLst 항목과 관계를 직접 지운다.
    복제가 모두 끝난 뒤에 원본을 지워야 한다 — 먼저 지우면 복제 소스가 사라진다.
    """
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        if prs.part.rels[sld_id.rId].target_part is slide.part:
            prs.part.drop_rel(sld_id.rId)
            xml_slides.remove(sld_id)
            return


def _new_layout_slide(prs, layout, tpl, warns, screen_id,
                      row_heights=None, size_pt=DETAIL_FONT_PT):
    """레이아웃으로 슬라이드를 만들고, 레이아웃에 없는 자리만 채워 넣는다.

    placeholder에 mapping이 정한 이름을 붙이므로, 값을 채우는 _fill_page는
    clone 모드와 똑같은 코드를 쓴다.

    row_heights와 size_pt는 _fit_tables가 그 화면 전체를 보고 정한 값이다.
    안 주면 실측 고정 높이와 7pt를 쓴다.
    """
    shapes_cfg = tpl.get("shapes", {})
    tables_cfg = tpl.get("detail_tables", {}) or {}
    count = int(tables_cfg.get("count", 5))
    rows = int(tables_cfg.get("rows", 4))
    area = tuple(tpl.get("content_area") or DEFAULT_CONTENT_AREA)

    slide = prs.slides.add_slide(layout)
    inherit_placeholders(slide, layout)
    name_placeholders(slide, tpl.get("placeholders", {}), shapes_cfg,
                      warns, screen_id)

    image_box, table_boxes = split_content_area(area, count, rows, row_heights)
    # shapes.image가 없으면 이미지 자리를 아예 그리지 않는다. 값을 채우는
    # _fill_page가 같은 키를 보고 이미지 블록 전체를 건너뛰기 때문이다 —
    # 여기서만 기본 이름으로 그려 두면 place_image가 지울 도형을 못 찾아
    # 본문 영역만 한 회색 사각형이 슬라이드마다 통째로 남는다(경고도 없이).
    # "shapes.image를 지정하지 않은 매핑은 이미지를 배치하지 않는다"는
    # clone 모드의 규칙을 그대로 따르는 셈이다.
    img_name = shapes_cfg.get("image")
    if img_name:
        add_image_anchor(slide, image_box, img_name)

    names = shapes_cfg.get("detail_tables") or []
    for i, box in enumerate(table_boxes):
        name = names[i] if i < len(names) else "상세표%d" % (i + 1)
        add_detail_table(slide, box, rows, name, row_heights, size_pt)
    return slide


def plan_pages(scr: dict, slot_count: int, work_dir: Path, image_box,
               enabled: bool):
    """화면 하나를 (상세 묶음, 이미지 경로) 페이지 목록으로 나눈다.

    이미지가 자리보다 길면 조각으로 나누고, 각 조각의 SoM 뱃지 수만큼 상세를
    배분한다. 조각의 상세가 슬롯을 넘으면 같은 조각을 여러 장에 반복해 싣는다 —
    그 상세가 가리키는 화면이 그 조각 안에 있기 때문이다.

    뱃지가 안 잡히거나 개수가 상세 건수와 어긋나면 분할만 하지 않고 기존처럼
    슬롯 단위로 나눈다. 조용히 엉뚱한 상세가 붙는 것보다 낫다.
    """
    details = scr.get("details", [])
    images = scr.get("images") or []
    plain = [(page, None) for page in chunk_details(details, slot_count)]
    if not enabled or not images or image_box is None:
        return plain, None

    src = Path(work_dir) / images[0]
    if not src.exists():
        return plain, None

    pieces, counts, total = split_for_box(src, image_box[2], image_box[3],
                                          Path(work_dir) / "images")
    if len(pieces) <= 1:
        return plain, None
    if not counts or sum(counts) != len(details):
        # 조각은 쓰되 상세는 순차 배분한다. 첫 조각만 쓰면 나머지 화면이 사라지므로
        # 페이지 수에 맞춰 조각을 돌려 쓴다.
        pages = []
        for i, page in enumerate(chunk_details(details, slot_count)):
            pages.append((page, pieces[min(i, len(pieces) - 1)]))
        return pages, ("뱃지 %d개가 상세 %d건과 맞지 않아 상세를 순서대로 배분했습니다"
                       % (total, len(details)))

    pages = []
    at = 0
    for piece, n in zip(pieces, counts):
        part = details[at:at + n]
        at += n
        if not part:
            pages.append(([], piece))
            continue
        for chunk in chunk_details(part, slot_count):
            pages.append((chunk, piece))
    return pages, None


def _cap_heights(heights, floors, limit: int) -> list[int]:
    """표 높이가 한계를 넘으면 하한을 지키며 비례 축소한다.

    가장 작은 글자로도 안 들어가는 상세가 있을 때 쓴다. 그냥 두면
    split_content_area가 ValueError를 던지고 화면 단위 격리에 걸려 그 화면이
    통째로 '[생성 실패]'가 된다 — 넘침을 막으려다 화면을 잃는 셈이다.

    높이를 깎으면 텍스트가 셀을 넘치지만 표는 슬라이드 안에 남는다. 넘친 항목은
    text-overflow 경고가 잡아 사람이 문장을 고치게 한다.
    """
    total = sum(heights)
    if total <= limit:
        return list(heights)
    slack = total - sum(floors)   # 하한 위로 늘어난 양
    room = limit - sum(floors)    # 하한 위로 허용된 양
    if slack <= 0 or room <= 0:
        return list(floors)
    return [f + (h - f) * room // slack for h, f in zip(heights, floors)]


def _fit_tables(pages, area, count: int, rows: int, text_key: str):
    """화면의 장별 행 높이와, 화면 전체가 공유할 설명 칸 글자 크기를 정한다.

    2패스의 두 번째다. 배분(pages)은 이미 확정됐고, 여기서는 그 배분에 맞는
    표 높이만 구한다 — 배분을 다시 돌리면 이미지 자리가 바뀌어 조각 수가 바뀌고,
    조각 수가 바뀌면 배분이 또 바뀌어 순환에 빠진다.

    높이는 장마다 자기 내용에 맞춰 잡는다. 전 장을 최악값으로 통일하면 긴 상세
    하나가 모든 장의 표를 밀어 올려 이미지 자리를 통째로 잡아먹는다. 배분을
    재계산하지 않으므로 장마다 높이가 달라도 순환은 생기지 않고, 이미지는
    place_image가 비율을 지켜 축소 배치한다.

    글자 크기만은 화면 단위로 통일한다 — 한 장만 작으면 장을 넘길 때 글자가
    커졌다 작아졌다 한다.

    (장별 행 높이 목록, 글자 크기, 낮췄는지)를 돌려준다.
    """
    floors = measured_row_heights(rows)
    width = detail_text_width(area[2], count)
    limit = area[3] - MIN_IMAGE_HEIGHT_EMU
    texts = [[str(d.get(text_key, "") or "") for d in page] for page, _ in pages]

    for size_pt in DETAIL_FONT_STEPS:
        per_page = [plan_row_heights([t], rows, width, size_pt, floors)
                    for t in texts]
        if all(sum(h) <= limit for h in per_page):
            return per_page, size_pt, size_pt != DETAIL_FONT_STEPS[0]
    return ([_cap_heights(h, floors, limit) for h in per_page],
            DETAIL_FONT_STEPS[-1], True)


def _today() -> str:
    """생성일. 표지 작성일과 같은 표기(YYYY-MM-DD)를 쓴다."""
    return date.today().isoformat()


def _doc_meta(screens_data: dict, mapping: dict) -> dict:
    """문서 meta에 생성일을 얹는다.

    화면설계서의 '작성일'은 그 PPT를 만든 날이다. Excel 표지의 작성일은
    Excel을 쓴 날이라 다르다 — 실제 샘플에서 두 달이 벌어져 있었다.
    그래서 표지 값이 있어도 생성일이 이긴다.

    표지 값을 그대로 쓰거나 특정 날짜를 박고 싶으면 options.date_field를
    null로 꺼라. 그러면 표지에서 읽은 값(또는 meta_overrides로 지정한 값)이
    그대로 간다.

    넘겨받은 screens_data는 건드리지 않는다 — 호출자의 것이다.
    """
    meta = dict(screens_data.get("meta") or {})
    field = mapping.get("options", {}).get("date_field", "작성일")
    if field:
        meta[field] = _today()
    return meta


def _fill_page(slide, scr: dict, page: list[dict], title: str, mapping: dict,
               work_dir: Path, warns: Warnings, meta: dict | None = None,
               image_path: Path | None = None) -> None:
    tpl = mapping["template"]
    shapes_cfg = tpl.get("shapes", {})
    opts = mapping.get("options", {})
    cols = tpl.get("table_columns", {"no": 0, "text": 1})
    text_key = opts.get("detail_text_source", "desc")
    clear_unused = bool(opts.get("clear_unused_slots", True))

    title_name = shapes_cfg.get("title")
    if title_name:
        shp = find_shape(slide, title_name)
        if shp is None:
            warns.add(scr["id"], "shape-not-found", "제목 도형 '%s' 없음" % title_name)
        else:
            set_text(shp, title)

    sid_name = shapes_cfg.get("screen_id")
    if sid_name:
        shp = find_shape(slide, sid_name)
        if shp is None:
            warns.add(scr["id"], "shape-not-found", "화면ID 도형 '%s' 없음" % sid_name)
        else:
            set_text(shp, scr["id"])

    # 화면별 fields가 문서 meta를 이긴다. 화면마다 다른 값이 있으면 그게 더 구체적이다.
    reserved = {"title", "screen_id", "image", "detail_tables"}
    doc_meta = meta or {}
    for key, name in shapes_cfg.items():
        if key in reserved:
            continue
        if key in (scr.get("fields") or {}):
            value = (scr.get("fields") or {})[key]
        elif key in doc_meta:
            value = doc_meta[key]
        else:
            continue
        shp = find_shape(slide, name)
        if shp is not None:
            set_text(shp, value)

    img_name = shapes_cfg.get("image")
    if img_name:
        anchor = find_shape(slide, img_name)
        images = scr.get("images") or []
        if anchor is None:
            warns.add(scr["id"], "shape-not-found", "이미지 자리 '%s' 없음" % img_name)
        elif not images and image_path is None:
            warns.add(scr["id"], "no-image", "배치할 이미지가 없습니다")
        else:
            # image_path는 분할된 조각이다. 없으면 원본 한 장을 그대로 쓴다.
            place_image(slide, anchor,
                        image_path if image_path is not None else work_dir / images[0])

    # 표 이름 검증은 build()가 소스 슬라이드를 상대로 한 번만 한다. 여기서 다시
    # warns를 넘기면 화면(슬라이드) 수만큼 같은 shape-not-found 경고가 반복된다.
    tables = collect_tables(slide, shapes_cfg.get("detail_tables"))
    if tables:
        fill_slots(tables, page, cols, text_key, clear_unused, warns, scr["id"])


def build(screens_data: dict, mapping: dict, work_dir: Path, out_path: Path,
          warns: Warnings) -> dict:
    """`work_dir`은 이미지와 템플릿을 푸는 기준이다. 결과물 경로는 `out_path`가
    따로 받는다 — 결과물 폴더와 작업 폴더는 다른 곳이다."""
    tpl = mapping["template"]
    template_path = resolve_template_path(tpl, work_dir)

    prs = Presentation(str(template_path))
    mode = tpl.get("mode", "clone")
    originals = list(prs.slides)

    if mode == "layout":
        layout = find_layout(prs, tpl.get("layout", 0))
        tables_cfg = tpl.get("detail_tables", {}) or {}
        count = int(tables_cfg.get("count", 5))
        rows = int(tables_cfg.get("rows", 4))
        slot_count = count * rows
        src = None
        image_box_for_split, _ = split_content_area(
            tuple(tpl.get("content_area") or DEFAULT_CONTENT_AREA), count, rows)
        area = tuple(tpl.get("content_area") or DEFAULT_CONTENT_AREA)
        # 수직(top/height)만 검사한다. 가로는 일부러 보지 않는다 —
        # DEFAULT_CONTENT_AREA(-12319, 337940, 9957099, 6331421)는 오른쪽으로
        # 슬라이드 폭(9906000)을 38780 EMU 넘어가도록 실측값 그대로 만들어졌다.
        # 원본 화면설계서의 스크린샷 자리가 실제로 그렇게 슬라이드 가장자리를
        # 살짝 넘겨 배치돼 있어서다. 소박한 좌우 경계 검사를 넣으면 이 정상
        # 기본값부터 막히므로 가로는 검사하지 않는다. 수직은 그런 의도된 여유가
        # 없고, top이 음수거나 아래로 슬라이드를 벗어나면 도형이 위아래로
        # 잘려 나가므로 여기서 막는다.
        if area[1] < 0 or area[1] + area[3] > int(prs.slide_height):
            raise ValueError(
                "content_area가 슬라이드 높이를 벗어납니다: top=%d height=%d, "
                "슬라이드 높이=%d" % (area[1], area[3], int(prs.slide_height))
            )
        # rows_per_table이 과도하면 이미지 자리가 1인치 밑으로 내려가는
        # split_content_area의 ValueError가 난다. 이건 화면마다 다른 사고가
        # 아니라 mapping 전체에 걸린 설정 오류다 — 루프 안(_new_layout_slide)
        # 에서만 부르면 화면 단위 예외 격리에 걸려 같은 오류가 화면 수만큼
        # screen-failed 경고로 흩어지고, build()는 조용히 끝나면서 슬라이드
        # 전부가 '[생성 실패]' 표지판이 된다. 루프 밖인 여기서 한 번 불러
        # 잘못된 설정이면 ValueError를 곧장 터뜨린다(같은 인자이므로 루프
        # 안의 호출은 그대로 둬도 다시 통과할 뿐이다).
        split_content_area(area, count, rows)
    else:
        source_slide = tpl.get("source_slide", 0)
        if source_slide is None:
            # analyze.py는 예시 슬라이드가 없는 템플릿에 source_slide: null과 함께
            # mode: layout을 제안한다. 생성 단계는 예시 슬라이드 복제만 지원하므로
            # int(None)이 던지는 원시 TypeError 대신 원인과 대안을 알려준다.
            raise ValueError(
                "template.source_slide가 없습니다. mode를 'layout'으로 두고 "
                "template.layout에 레이아웃 이름을 지정하거나, 예시 슬라이드가 "
                "있는 템플릿을 --template으로 지정하세요."
            )
        src = prs.slides[int(source_slide)]
        # 이미지 자리는 layout 모드처럼 계산할 수 없다. 대신 예시 슬라이드의 앵커
        # 도형이 그 자리다 — place_image가 이미지를 앉힐 때 보는 것과 같은 값을
        # 분할 계산도 봐야 목표 높이와 실제 배치가 어긋나지 않는다. 소스는 한 장뿐
        # 이므로 여기서 한 번만 찾는다. 앵커가 없으면 분할하지 않는다(경고는
        # _fill_page가 화면마다 shape-not-found로 남기므로 여기서 겹쳐 내지 않는다).
        anchor = find_shape(src, tpl.get("shapes", {}).get("image") or "")
        image_box_for_split = None
        if anchor is not None and anchor.width and anchor.height:
            image_box_for_split = (anchor.left, anchor.top,
                                   anchor.width, anchor.height)
        slot_count = count_slots(
            collect_tables(src, tpl.get("shapes", {}).get("detail_tables"), warns)
        )

    made = 0
    split_ids: list[str] = []
    failed_ids: list[str] = []
    screen_count = 0
    doc_meta = _doc_meta(screens_data, mapping)
    # _fill_page가 셀에 넣는 것과 같은 필드를 봐야 표 높이가 내용과 맞는다.
    text_key = mapping.get("options", {}).get("detail_text_source", "desc")

    for scr in screens_data.get("screens", []):
        screen_count += 1
        # id는 화면을 가리키는 핵심 키이지만 screens.json에 빠져 있을 수 있다
        # (추출이 어긋났거나 사람이 중간 산출물을 손댄 경우).
        # try 진입 전에 방어적으로 계산해 둬야, id가 없어서 나는 예외를 처리하는
        # except 블록 안에서 scr["id"]가 또 KeyError를 내며 전체 빌드를 무너뜨리는
        # 일이 없다 — 화면 단위 격리 보장이 바로 이 시나리오를 위해 있다.
        scr_id = scr.get("id") or "(id 없음 #%d)" % screen_count
        try:
            split_on = str(mapping.get("options", {}).get("image_split", "auto")) != "off"
            pages, note = plan_pages(scr, slot_count, work_dir,
                                     image_box_for_split, split_on)
            if note:
                print("  [%s] %s" % (scr_id, note))
            if len(pages) > 1:
                split_ids.append(scr_id)
                warns.add(scr_id, "slide-split",
                          "상세 %d건이 슬롯 %d개를 넘어 %d장으로 나눴습니다"
                          % (len(scr["details"]), slot_count, len(pages)))
            # 2패스: 배분이 확정됐으니 그 상세에 맞는 표 높이를 구한다.
            # 높이는 장별로, 글자 크기는 화면 단위로 정해진다.
            page_heights = size_pt = None
            if mode == "layout":
                page_heights, size_pt, shrunk = _fit_tables(
                    pages, area, count, rows, text_key)
                if shrunk:
                    print("  [%s] 상세가 길어 설명 글자를 %.1fpt로 낮췄습니다"
                          % (scr_id, size_pt))

            for i, (page, page_image) in enumerate(pages):
                if mode == "layout":
                    slide = _new_layout_slide(prs, layout, tpl, warns, scr_id,
                                              page_heights[i], size_pt)
                else:
                    slide = clone_slide(prs, src)
                _fill_page(slide, scr, page,
                           page_title(scr["name"], i, len(pages)),
                           mapping, work_dir, warns, doc_meta, page_image)
                if mode == "layout":
                    drop_empty_placeholders(slide)
                made += 1
        except Exception as exc:
            failed_ids.append(scr_id)
            warns.add(scr_id, "screen-failed", "슬라이드 생성 실패: %s" % exc)
            if len(prs.slides) > len(originals) + made:
                # 예외 직전에 만들어진 슬라이드가 남아 있다. 어디가 실패했는지
                # 결과물에서 바로 보이도록 제목만 실패 표시로 바꾼다.
                title_name = tpl.get("shapes", {}).get("title")
                shp = find_shape(prs.slides[-1], title_name) if title_name else None
                if shp is not None:
                    set_text(shp, "[생성 실패] %s" % scr.get("name", scr_id))
                made += 1

    for slide in originals:
        _drop_slide(prs, slide)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    return {
        "slides": made,
        "screens": screen_count,
        "split": split_ids,
        "failed": failed_ids,
    }


def main(argv: list[str] | None = None) -> int:
    setup_stdio()
    ap = argparse.ArgumentParser(description="screens.json으로 화면설계서 PPT를 만든다")
    ap.add_argument("--screens", default=None,
                    help="생략하면 결과물 폴더의 .work/screens.json을 읽는다")
    ap.add_argument("--mapping", default=None,
                    help="생략하면 결과물 폴더의 .work/mapping.json을 읽는다")
    ap.add_argument("--output", required=True,
                    help="결과물 폴더. pptx가 여기 저장되고 작업 파일은 .work/에 있다")
    ap.add_argument("--out-file", default=None,
                    help="결과물 경로를 직접 지정한다. 생략하면 원본 Excel 파일명을 쓴다")
    args = ap.parse_args(argv)

    output_dir = Path(args.output)
    given = [Path(p) for p in (args.screens, args.mapping) if p]
    notice = migration_notice(migrate_legacy_work(output_dir, keep=given))
    if notice:
        print(notice)
    work = work_dir(output_dir)

    screens_data = read_json(Path(args.screens) if args.screens
                             else work / "screens.json")
    mapping = read_json(Path(args.mapping) if args.mapping
                        else work / "mapping.json")
    warns = Warnings()
    source = screens_data.get("meta", {}).get("source", "")
    if args.out_file:
        # 명시 지정은 그대로 쓴다. 사용자가 파일명을 직접 정했으면 덮어쓰기가
        # 의도한 동작이고, 코드가 이를 뒤집으면 결과를 예측할 수 없다.
        out_path = Path(args.out_file)
        numbered = False
    else:
        out_path = resolve_output_path(output_dir, source)
        numbered = out_path.stem != output_stem(source)

    report = build(screens_data, mapping, work, out_path, warns)

    print("슬라이드 %d장 생성 (화면 %d개)" % (report["slides"], report["screens"]))
    if report["split"]:
        print("분할된 화면: %s" % ", ".join(report["split"]))
    if report["failed"]:
        print("실패한 화면: %s" % ", ".join(report["failed"]))
    if len(warns):
        print(warns.format())
    result = verify_output(out_path, screens_data, mapping, report["slides"], work)
    print("검증: %s" % ("통과" if result["ok"] else "실패"))
    for c in result["checks"]:
        print("  [%s] %s — %s" % ("O" if c["ok"] else "X", c["name"], c["detail"]))
    # 번호가 붙었으면 알린다. 조용히 다른 이름으로 저장하면 사용자가 옛 파일을
    # 열어 보고 반영이 안 됐다고 오해한다.
    if numbered:
        print("저장: %s (같은 이름이 있어 번호를 붙였습니다)" % out_path)
    else:
        print("저장: %s" % out_path)
    return 0


if __name__ == "__main__":
    sys.exit(main())
