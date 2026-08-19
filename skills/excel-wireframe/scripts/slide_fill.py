# -*- coding: utf-8 -*-
"""슬라이드 도형에 값을 채운다.

핵심은 서식 보존이다. 런을 전부 지우고 새로 만들면 폰트·크기·색이 초기화되어
템플릿 디자인이 무너지므로, 첫 런의 텍스트만 갈아끼우는 방식을 쓴다.
"""
from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx.dml.color import RGBColor
from text_metrics import fits_lines, text_lines


def find_shape(slide, name: str):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# 미입력 표시용 "입력필요"(빨강·굵게). 값을 못 찾은 자리를 비워 두면 채워야 할
# 자리가 있다는 사실 자체가 산출물에서 사라진다 — 눈에 띄게 남겨 사람이 채우게 한다.
INPUT_REQUIRED = "입력필요"
INPUT_REQUIRED_COLOR = RGBColor(0xFF, 0x00, 0x00)


def _drop_fields(paragraph) -> None:
    """문단에서 자동 필드(a:fld)를 걷어낸다.

    날짜·쪽번호 같은 자동 필드는 런(a:r)이 아니라서 paragraph.runs에 잡히지
    않는다. 그대로 두고 런만 갈아끼우면 필드가 보여 주던 값 뒤에 새 값이
    이어 붙어 '2026-05-292026-06-25'처럼 나온다. 값을 명시적으로 채우기로
    한 자리에 자동 필드가 남아 있을 이유는 없다.

    채우지 않는 자리(쪽번호 등)는 이 경로를 타지 않으므로 필드가 보존된다.
    """
    p = paragraph._p
    for fld in list(p.findall("{%s}fld" % A_NS)):
        p.remove(fld)


def _emphasize(tf) -> None:
    """텍스트 프레임의 모든 런을 빨강·굵게로 바꾼다.

    런을 새로 만들지 않고 이미 놓인 런의 속성만 건드린다 — 크기·글꼴 같은
    나머지 서식은 템플릿 값 그대로 남아야 한다.
    """
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.bold = True
            run.font.color.rgb = INPUT_REQUIRED_COLOR


def _fill_text_frame(tf, text: str, emphasis: bool = False) -> None:
    lines = str(text).split("\n") if text else [""]
    p0 = tf.paragraphs[0]
    _drop_fields(p0)

    # 빈 자리는 런이 없고 서식을 endParaRPr에 품고 있다 — PowerPoint가 "다음에
    # 칠 글자의 서식"으로 저장해 두는 자리다. 그냥 런을 만들면 그 서식을 잃고
    # 크기가 상속(기본 18pt)으로 떨어져, 0.13인치짜리 메타 표 행이 부풀어 버린다.
    # 원본이 정해 둔 크기를 그대로 물려받아야 자리가 원본 상태로 남는다.
    end_rPr = p0._p.find("{%s}endParaRPr" % A_NS)
    if p0.runs:
        base_run = p0.runs[0]
        rPr = base_run._r.find("{%s}rPr" % A_NS)
    else:
        base_run = p0.add_run()
        rPr = None
        if end_rPr is not None:
            rPr = copy.deepcopy(end_rPr)
            rPr.tag = "{%s}rPr" % A_NS
            base_run._r.insert(0, rPr)
    base_rPr = copy.deepcopy(rPr) if rPr is not None else None

    # 문단 속성(줄간격·정렬)도 물려줘야 한다. 런 속성만 옮기면 둘째 줄부터
    # 줄간격이 폰트 기본값으로 돌아가 행이 계산보다 커진다.
    pPr = p0._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
    )
    base_pPr = copy.deepcopy(pPr) if pPr is not None else None

    base_run.text = lines[0]
    for extra in p0.runs[1:]:
        extra._r.getparent().remove(extra._r)
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)

    for line in lines[1:]:
        p = tf.add_paragraph()
        # pPr은 a:p의 첫 자식이어야 한다.
        if base_pPr is not None:
            p._p.insert(0, copy.deepcopy(base_pPr))
        run = p.add_run()
        run.text = line
        if base_rPr is not None:
            run._r.insert(0, copy.deepcopy(base_rPr))

    if emphasis:
        _emphasize(tf)


def set_text(shape, text: str, emphasis: bool = False) -> None:
    if not shape.has_text_frame:
        return
    _fill_text_frame(shape.text_frame, text, emphasis)


def set_cell_text(cell, text: str, emphasis: bool = False) -> None:
    _fill_text_frame(cell.text_frame, text, emphasis)


def set_text_or_required(target, value, is_cell: bool = False) -> bool:
    """값이 있으면 그대로, 없으면 "입력필요"를 빨강·굵게 쓴다.

    도형과 표 셀 둘 다 받는다. 채웠으면 True, 입력필요를 남겼으면 False.
    """
    text = "" if value is None else str(value).strip()
    write = set_cell_text if is_cell else set_text
    if text:
        write(target, text)
        return True
    write(target, INPUT_REQUIRED, True)
    return False


def estimate_overflow(text, cell_width_emu: int, cell_height_emu: int,
                      size_pt: float) -> bool:
    """셀에 든 글이 행 높이를 넘기는지 본다.

    build가 표를 미리 늘려 두므로 대개는 넘치지 않는다. 가장 작은 글자로도
    안 들어가는 상세만 여기 걸린다 — 그런 문장은 사람이 줄여야 한다.
    """
    if not text:
        return False
    return (text_lines(text, cell_width_emu, size_pt)
            > fits_lines(cell_height_emu, size_pt))


def _cell_font_pt(cell, default: float = 7.0) -> float:
    """셀에 심어 둔 글자 크기. add_detail_table이나 템플릿이 정한 값이다."""
    runs = cell.text_frame.paragraphs[0].runs
    if runs and runs[0].font.size is not None:
        return runs[0].font.size.pt
    return default


def collect_tables(slide, names: list[str] | None, warns=None, screen_id=None):
    """상세 표를 슬롯 순서대로 모은다.

    이름이 지정되면 그 순서를 그대로 따른다. 없으면 좌→우로 정렬한다 —
    화면상 왼쪽 표가 앞 번호를 담는 것이 사람의 읽기 순서와 맞기 때문이다.

    지정된 이름이 슬라이드에 없으면 shape-not-found 경고를 기록한다 (warns 제공 시).
    """
    tables = [s for s in slide.shapes if s.has_table]
    if names:
        by_name = {t.name: t for t in tables}
        result = []
        for n in names:
            if n in by_name:
                result.append(by_name[n])
            elif warns is not None:
                warns.add(screen_id, "shape-not-found",
                          "표 '%s'을(를) 템플릿에서 찾지 못했습니다" % n)
        return result
    return sorted(tables, key=lambda t: (t.left or 0, t.top or 0))


def count_slots(tables) -> int:
    return sum(len(t.table.rows) for t in tables)


def place_image(slide, anchor_shape, image_path: Path):
    """앵커 도형 자리에 종횡비를 유지해 이미지를 넣고 앵커는 지운다.

    세로는 자리 위쪽에 붙이고 가로는 가운데에 둔다. 스크린샷은 위에서부터 읽는
    것이라 장을 넘길 때 상단 기준선이 흔들리지 않아야 하고, 세로로 가운데를
    맞추면 짧은 이미지가 아래 상세표 쪽으로 내려앉아 위가 휑하게 빈다.
    """
    left = anchor_shape.left
    top = anchor_shape.top
    box_w = anchor_shape.width
    box_h = anchor_shape.height

    with Image.open(image_path) as im:
        iw, ih = im.size

    scale = min(box_w / iw, box_h / ih)
    new_w = int(iw * scale)
    new_h = int(ih * scale)
    new_left = left + (box_w - new_w) // 2
    new_top = top

    pic = slide.shapes.add_picture(str(image_path), new_left, new_top, new_w, new_h)
    anchor_shape._element.getparent().remove(anchor_shape._element)
    return pic


def fill_slots(
    tables,
    details: list[dict],
    cols: dict,
    text_key: str,
    clear_unused: bool,
    warns,
    screen_id: str,
) -> int:
    """상세 항목을 표1 r0…rN → 표2 r0…rN 순으로 흘려 넣는다."""
    no_col = int(cols.get("no", 0))
    text_col = int(cols.get("text", 1))

    slots = []
    for t in tables:
        table = t.table
        width = 0
        if text_col < len(table.columns):
            width = int(table.columns[text_col].width or 0)
        for r in range(len(table.rows)):
            slots.append((table, r, width))

    filled = 0
    for i, (table, row, width) in enumerate(slots):
        if i < len(details):
            d = details[i]
            text = str(d.get(text_key, "") or "")
            if no_col < len(table.columns):
                set_cell_text(table.cell(row, no_col), str(d.get("no", "") or ""))
            if text_col < len(table.columns):
                cell = table.cell(row, text_col)
                set_cell_text(cell, text)
                # 폭은 열 폭에서 좌우 여백을 뺀 값이다. 높이와 글자 크기는
                # 그 셀에 실제로 들어간 값을 읽는다 — build가 화면마다 다르게
                # 정할 수 있으므로 고정값으로 재면 어긋난다.
                inner = width - (cell.margin_left or 0) - (cell.margin_right or 0)
                if estimate_overflow(text, inner, table.rows[row].height,
                                     _cell_font_pt(cell)):
                    warns.add(screen_id, "text-overflow",
                              "%s번 항목의 설명이 셀을 넘길 수 있습니다"
                              % d.get("no", "?"))
            filled += 1
        elif clear_unused:
            if no_col < len(table.columns):
                set_cell_text(table.cell(row, no_col), "")
            if text_col < len(table.columns):
                set_cell_text(table.cell(row, text_col), "")

    if len(details) > len(slots):
        warns.add(screen_id, "slot-shortage",
                  "상세 %d건 중 %d건만 이 슬라이드에 들어갔습니다"
                  % (len(details), len(slots)))
    return filled
