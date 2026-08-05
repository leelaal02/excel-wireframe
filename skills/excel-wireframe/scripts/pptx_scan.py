# -*- coding: utf-8 -*-
"""PPTX 구조 스캔. 어느 슬라이드를 복제 소스로 쓸지 판단할 재료를 만든다."""
from __future__ import annotations

from pathlib import Path

from common import EMU_PER_INCH
from pptx import Presentation


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _scan_shape(shape) -> dict:
    table = None
    if shape.has_table:
        table = {"rows": len(shape.table.rows), "cols": len(shape.table.columns)}
    return {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "is_placeholder": bool(shape.is_placeholder),
        "left": int(shape.left) if shape.left is not None else 0,
        "top": int(shape.top) if shape.top is not None else 0,
        "width": int(shape.width) if shape.width is not None else 0,
        "height": int(shape.height) if shape.height is not None else 0,
        "text": _shape_text(shape)[:200],
        "table": table,
    }


def scan_presentation(path: Path) -> dict:
    path = Path(path)
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = [_scan_shape(s) for s in slide.shapes]
        slides.append(
            {
                "index": i,
                "layout": slide.slide_layout.name,
                "shape_count": len(shapes),
                "text_shape_count": sum(1 for s in shapes if s["text"]),
                "shapes": shapes,
            }
        )
    return {
        "file": str(path),
        "slide_width": int(prs.slide_width),
        "slide_height": int(prs.slide_height),
        "slide_size_in": [
            prs.slide_width / EMU_PER_INCH,
            prs.slide_height / EMU_PER_INCH,
        ],
        "slides": slides,
    }


def scan_layouts(path: Path) -> list[dict]:
    """모든 마스터의 레이아웃과 그 안의 자리·도형을 훑는다.

    어느 레이아웃을 쓸지, 본문 영역을 어디로 잡을지 사람이 판단할 재료다.
    """
    prs = Presentation(str(path))
    out = []
    for mi, master in enumerate(prs.slide_masters):
        for li, lay in enumerate(master.slide_layouts):
            out.append({
                "master": mi,
                "index": li,
                "name": lay.name,
                "placeholders": [
                    {
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type),
                        "left": int(ph.left or 0),
                        "top": int(ph.top or 0),
                        "width": int(ph.width or 0),
                        "height": int(ph.height or 0),
                    }
                    for ph in lay.placeholders
                ],
                "shapes": [
                    _scan_shape(s) for s in lay.shapes if not s.is_placeholder
                ],
            })
    return out


def suggest_content_area(layout_info: dict, slide_width: int,
                         slide_height: int) -> list[int]:
    """레이아웃에서 이미지와 상세표를 놓을 만한 가로 띠를 고른다.

    껍데기는 대개 상단과 하단에 가로로 깔린다. 그래서 완전한 빈 영역 탐색 대신
    '슬라이드 폭의 절반 이상을 덮는 도형'만 장애물로 보고, 위아래에서 잠식된
    만큼 깎는다. 추정이 빗나가도 mapping.json에서 고칠 수 있으므로 이 정도면
    충분하다.
    """
    half = slide_width // 2
    blockers = []
    for s in layout_info["shapes"] + layout_info["placeholders"]:
        w = s.get("width", 0)
        if w >= half:
            blockers.append((s.get("top", 0), s.get("top", 0) + s.get("height", 0)))

    top = 0
    bottom = slide_height
    for b_top, b_bottom in blockers:
        if b_bottom <= slide_height // 2:
            top = max(top, b_bottom)      # 상단 껍데기
        elif b_top >= slide_height // 2:
            bottom = min(bottom, b_top)   # 하단 껍데기

    return [0, int(top), int(slide_width), int(bottom - top)]


def suggest_mode(report: dict) -> dict:
    """예시 슬라이드가 있으면 clone, 없으면 layout.

    '텍스트가 채워진 도형 3개 이상 + 표나 그림 1개 이상'을 예시 슬라이드의 신호로 본다.
    빈 레이아웃만 있는 템플릿은 이 조건을 통과하지 못한다.
    """
    for s in report["slides"]:
        has_visual = any(
            sh["table"] is not None or "PICTURE" in sh["shape_type"]
            for sh in s["shapes"]
        )
        if s["text_shape_count"] >= 3 and has_visual:
            return {
                "mode": "clone",
                "source_slide": s["index"],
                "reason": "슬라이드 %d에 채워진 텍스트 %d개와 표/그림이 있어 예시 슬라이드로 판단"
                % (s["index"], s["text_shape_count"]),
            }
    return {
        "mode": "layout",
        "source_slide": None,
        "reason": "예시 슬라이드를 찾지 못해 빈 레이아웃 모드로 판단",
    }
