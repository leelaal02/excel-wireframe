# -*- coding: utf-8 -*-
"""PPT 템플릿을 사용자가 주지 않았을 때 쓸 기본 템플릿을 만든다.

디자인은 레이아웃이 담당한다. 상단 메타 표·본문 박스·하단 바 같은 껍데기는
레이아웃 위 도형이고, 글자가 들어가는 자리는 placeholder이거나 메타 표의 칸이다 —
빌드는 레이아웃을 골라 슬라이드를 추가하고 그 자리에 값을 채운다. 슬라이드는
한 장도 넣지 않는다.

구조는 실제 화면설계서 템플릿의 '내용설명연결' 레이아웃에서 그대로 잰 값이다.
화면명·작성일 같은 문서 정보는 상단 메타 표의 라벨 오른쪽 칸에 들어간다 —
그 자리를 placeholder로 따로 만들지 않는 것이 원본 구조다.

**메타 표는 레이아웃에 두지만, PowerPoint는 레이아웃 도형 중 placeholder만
슬라이드에 물려준다.** 표는 배경으로 비쳐 보일 뿐이므로 빌드가
`slide_layout.inherit_layout_tables`로 복제해 얹은 뒤에 값을 쓴다.

원본 레이아웃의 저작권 문구와 로고는 복제하지 않는다. 이 파일은 스킬에 코드로
담겨 배포되므로, 넣으면 제3자의 저작권 표기가 모든 사용자 산출물에 박힌다.
그 자리(하단 왼쪽)에는 Excel 표지에서 읽은 문서제목이 들어간다.
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt
from slide_clone import drop_slide
from slide_layout import append_shape_with_new_id, draw_cell_edges, scheme_fill

DEFAULT_LAYOUT_NAME = "내용설명연결"

DEFAULT_SHAPE_NAMES = {
    "screen_id": "화면ID",
    "image": "화면이미지",
    "doc_title": "문서제목",
    "page_no": "쪽번호",
}

META_TABLE_NAMES = ["메타표1", "메타표2"]

# python-pptx 기본 템플릿 'Title and Content'의 placeholder idx.
# 제목(0)과 날짜(10)는 메타 표가 그 몫을 하므로 레이아웃에서 지운다.
PLACEHOLDER_IDX = {
    "screen_id": 1,
    "문서제목": 11,
    "쪽번호": 12,
}
DROP_PLACEHOLDER_IDX = (0, 10)

BASE_LAYOUT_INDEX = 1  # 'Title and Content'

# --- 실측값 (실제 화면설계서 '내용설명연결' 레이아웃) ---
MEASURED_SLIDE = (9144000, 6858000)
MEASURED_PLACEHOLDERS = {
    1: (7496632, 188657, 1635346, 144000),   # 화면ID (화면ID배경 위)
    11: (0, 6750725, 2444995, 75085),        # 문서제목 (원본의 저작권 문구 자리)
    12: (4370029, 6716266, 403943, 144000),  # 쪽번호
}
# placeholder 자리가 작아 기본 글자 크기로는 넘친다. 실측 자리에 맞춰 지정한다.
PLACEHOLDER_FONT_PT = {1: 8.0, 11: 5.0, 12: 8.0}

MEASURED_SHELL = {
    "본문박스": (-1, 404664, 9131979, 6264697),
    "화면ID배경": (7496632, 188657, 1635346, 144000),
    "하단바": (0, 6716266, 9144000, 144000),
}
# 껍데기를 그리는 순서. 화면ID배경은 화면ID placeholder보다 먼저 와야 뒤에 깔린다.
SHELL_ORDER = ("본문박스", "화면ID배경", "하단바")

# --- 상단 메타 표 (실측) ---
# 라벨과 값이 번갈아 놓인 1행짜리 표다. 값 칸은 비워 두고 빌드가 채운다.
META_ROW_HEIGHT = 116632
META_TABLES = (
    {
        "name": META_TABLE_NAMES[0],
        "box": (1, 0, 9134931, META_ROW_HEIGHT),
        "labels": ("프로젝트명", "산출물명", "화면명", "버전", "작성자",
                   "검토자", "작성일", "수정일", "ID"),
        "widths": (450927, 1163077, 365539, 1163077, 299077, 1163077,
                   226665, 265847, 299077, 398769, 265877, 398769,
                   299077, 465284, 299077, 465231, 199407, 947077),
    },
    {
        "name": META_TABLE_NAMES[1],
        "box": (1, 195617, 7430165, META_ROW_HEIGHT),
        "labels": ("네비게이션", "화면유형", "알림여부"),
        "widths": (450927, 4254011, 398813, 930566, 398813, 997035),
    },
)
META_FONT_PT = 6.5
META_CELL_MARGIN_H = 16616
META_CELL_MARGIN_V = 19087
META_BORDER_WIDTH = 3175
META_BORDER_SCHEME = ("tx1", None, None)
META_LABEL_FILL = ("tx1", 65000, 35000)   # 라벨 칸 배경(진회색)
META_VALUE_FILL = ("bg1", 85000, None)    # 값 칸 배경(연회색)

# 껍데기 색 (실측)
SHELL_LINE_SCHEME = ("tx1", 75000, 25000)
BAR_FILL_SCHEME = ("tx1", 65000, 35000)
SCREEN_ID_FILL = "EAE9EF"
SHELL_LINE_WIDTH = 6350


def _scaled(geom, sx: float, sy: float):
    left, top, width, height = geom
    return (int(left * sx), int(top * sy), int(width * sx), int(height * sy))


def _set_line(shape, scheme, width: int) -> None:
    """도형 테두리를 schemeClr로 칠한다. python-pptx는 lumMod를 못 다룬다."""
    ln = shape.line._get_or_add_ln()
    for child in list(ln):
        ln.remove(child)
    ln.set("w", str(width))
    scheme_fill(ln, "solidFill", scheme)


def _shell_rect(shapes, name, geom):
    """껍데기 사각형 하나. 이름마다 채움과 테두리가 다르다."""
    left, top, width, height = geom
    shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                           Emu(width), Emu(height))
    shp.name = name
    shp.text_frame.text = ""

    if name == "본문박스":
        shp.fill.background()
        _set_line(shp, SHELL_LINE_SCHEME, SHELL_LINE_WIDTH)
    elif name == "화면ID배경":
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(SCREEN_ID_FILL)
        _set_line(shp, SHELL_LINE_SCHEME, SHELL_LINE_WIDTH)
    else:  # 하단바
        spPr = shp._element.spPr
        for tag in ("solidFill", "noFill"):
            for old in spPr.findall(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag):
                spPr.remove(old)
        prst = spPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
        fill = scheme_fill(spPr, "solidFill", BAR_FILL_SCHEME)
        prst.addnext(fill)
        shp.line.fill.background()
    return shp


def _meta_cell(cell, text: str, fill) -> None:
    """메타 표 셀 하나에 실측 서식을 입힌다. 값 칸은 text가 빈 문자열이다.

    값 칸에도 런을 하나 심어 둔다 — 런이 없으면 빌드가 값을 채울 때 글자 크기가
    기본값(18pt)으로 잡혀 칸을 넘는다.
    """
    cell.margin_left = Emu(META_CELL_MARGIN_H)
    cell.margin_right = Emu(META_CELL_MARGIN_H)
    cell.margin_top = Emu(META_CELL_MARGIN_V)
    cell.margin_bottom = Emu(META_CELL_MARGIN_V)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(META_FONT_PT)
    run.font.bold = True
    run.font.name = "맑은 고딕"

    # 여백·정렬을 먼저 넣고 테두리를 그린다. python-pptx가 margin을 쓰면서
    # tcPr을 새로 만들 수 있어, 테두리를 먼저 그리면 지워진다.
    draw_cell_edges(cell, fill, META_BORDER_SCHEME, META_BORDER_WIDTH)


def _add_meta_table(shapes, spec, sx: float, sy: float):
    """라벨/값이 번갈아 놓인 1행짜리 메타 표를 만든다."""
    left, top, width, height = _scaled(spec["box"], sx, sy)
    cols = len(spec["widths"])
    frame = shapes.add_table(1, cols, Emu(left), Emu(top), Emu(width), Emu(height))
    frame.name = spec["name"]
    table = frame.table

    for i, w in enumerate(spec["widths"]):
        table.columns[i].width = Emu(int(w * sx))
    table.rows[0].height = Emu(int(META_ROW_HEIGHT * sy))

    for i in range(cols):
        is_label = i % 2 == 0
        text = spec["labels"][i // 2] if is_label else ""
        _meta_cell(table.cell(0, i), text,
                   META_LABEL_FILL if is_label else META_VALUE_FILL)
    return frame


def build_default_template(
    path: Path,
    slide_width_emu: int = MEASURED_SLIDE[0],
    slide_height_emu: int = MEASURED_SLIDE[1],
) -> Path:
    """레이아웃 하나짜리 기본 템플릿을 만든다. 슬라이드는 넣지 않는다."""
    prs = Presentation()
    prs.slide_width = Emu(slide_width_emu)
    prs.slide_height = Emu(slide_height_emu)

    sx = slide_width_emu / MEASURED_SLIDE[0]
    sy = slide_height_emu / MEASURED_SLIDE[1]

    layout = prs.slide_layouts[BASE_LAYOUT_INDEX]
    layout.name = DEFAULT_LAYOUT_NAME

    # 1) 껍데기와 메타 표 — LayoutShapes에는 add_shape/add_table이 없어 임시
    #    슬라이드를 거친다. 빈 화면 레이아웃을 써서 그 레이아웃의 placeholder가
    #    섞이지 않게 한다.
    tmp = prs.slides.add_slide(prs.slide_layouts[6])
    for name in SHELL_ORDER:
        _shell_rect(tmp.shapes, name, _scaled(MEASURED_SHELL[name], sx, sy))
    for spec in META_TABLES:
        _add_meta_table(tmp.shapes, spec, sx, sy)

    spTree = layout.shapes._spTree
    for shp in list(tmp.shapes):
        # 임시 슬라이드에서 매겨진 cNvPr/@id를 그대로 들고 오면 레이아웃
        # placeholder의 id와 겹친다 — append_shape_with_new_id가 그 규칙을 담고
        # 있으니 다른 파트에서 도형을 이식할 때는 그 함수를 쓴다.
        append_shape_with_new_id(spTree, copy.deepcopy(shp._element))
    drop_slide(prs, tmp)

    # 2) 쓰지 않는 placeholder를 지운다. 남겨 두면 analyze가 본문 영역을 추정할 때
    #    제목 자리를 상단 껍데기로 오인해 영역이 통째로 줄어든다.
    for ph in list(layout.placeholders):
        if ph.placeholder_format.idx in DROP_PLACEHOLDER_IDX:
            ph._element.getparent().remove(ph._element)

    # 3) 남은 placeholder를 실측 자리로 옮긴다. 껍데기 뒤에 붙어야 위에 그려진다.
    #    list()로 감싸는 이유: 루프 안에서 _spTree를 재정렬하므로 살아 있는
    #    반복자를 그대로 쓰면 요소를 건너뛴다.
    for ph in list(layout.placeholders):
        idx = ph.placeholder_format.idx
        geom = MEASURED_PLACEHOLDERS.get(idx)
        if geom is None:
            continue
        ph.left, ph.top, ph.width, ph.height = (
            Emu(v) for v in _scaled(geom, sx, sy)
        )
        # python-pptx 기본 레이아웃의 placeholder는 안내 문구('Click to edit
        # Master title style')와 날짜 필드('1/27/13')를 갖고 있다.
        # inherit_placeholders는 레이아웃 XML을 그대로 복제하므로 비워 두지
        # 않으면 그 문구가 슬라이드에 그대로 실린다. 쪽번호(12)만 자동 필드를
        # 살려 둔다.
        if idx != PLACEHOLDER_IDX["쪽번호"]:
            ph.text_frame.text = ""
        para = ph.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(PLACEHOLDER_FONT_PT.get(idx, 8.0))
        layout.shapes._spTree.append(ph._element)  # z-order를 맨 위로

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def default_meta_table_labels() -> dict:
    """메타 표 라벨 → 값을 가져올 이름. 라벨과 이름이 같으면 그대로 쓴다.

    `title`과 `screen_id`는 화면의 이름·ID를 가리키는 예약어이고, `문서제목`은
    Excel 표지의 표제다. 나머지는 화면별 fields나 문서 meta에서 같은 이름을 찾는다.
    """
    labels = {
        "화면명": "title",
        "ID": "screen_id",
        "산출물명": "문서제목",
    }
    for name in ("프로젝트명", "버전", "작성자", "검토자", "작성일", "수정일",
                 "네비게이션", "화면유형", "알림여부"):
        labels[name] = name
    return labels


def default_template_mapping(path: Path, table_count: int = 5,
                             rows_per_table: int = 4) -> dict:
    """생성한 기본 템플릿에 대응하는 mapping.json의 template 섹션."""
    return {
        "file": str(path),
        "mode": "layout",
        "layout": DEFAULT_LAYOUT_NAME,
        "placeholders": dict(PLACEHOLDER_IDX),
        "shapes": {
            "screen_id": DEFAULT_SHAPE_NAMES["screen_id"],
            "image": DEFAULT_SHAPE_NAMES["image"],
            "문서제목": DEFAULT_SHAPE_NAMES["doc_title"],
            "쪽번호": DEFAULT_SHAPE_NAMES["page_no"],
            "detail_tables": ["상세표%d" % (i + 1) for i in range(table_count)],
        },
        "meta_table": {
            "tables": list(META_TABLE_NAMES),
            "labels": default_meta_table_labels(),
        },
        "content_area": list(MEASURED_SHELL["본문박스"]),
        "detail_tables": {"count": table_count, "rows": rows_per_table},
        "table_columns": {"no": 0, "text": 1},
    }
