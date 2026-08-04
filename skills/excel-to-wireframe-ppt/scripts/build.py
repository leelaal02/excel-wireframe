# -*- coding: utf-8 -*-
"""3단계: screens.json과 템플릿으로 PPT를 만든다.

screens.json이 SSOT이므로 이 모듈은 Excel을 전혀 모른다. mapping.json에서도
template / options 섹션만 읽는다. 덕분에 Excel 픽스처 없이 빌드 로직을 검증할 수 있다.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from common import Warnings, read_json, setup_stdio
from pptx import Presentation
from slide_clone import clone_slide
from slide_fill import (
    collect_tables,
    count_slots,
    fill_slots,
    find_shape,
    place_image,
    set_text,
)


def chunk_details(details: list[dict], slot_count: int) -> list[list[dict]]:
    """슬롯 총량을 넘는 상세를 다음 슬라이드 분량으로 쪼갠다."""
    if not details:
        return [[]]
    if slot_count <= 0:
        return [details]
    return [
        details[i : i + slot_count] for i in range(0, len(details), slot_count)
    ]


def page_title(name: str, index: int, total: int) -> str:
    if total <= 1:
        return name
    return "%s (%d/%d)" % (name, index + 1, total)


def _drop_slide(prs, slide) -> None:
    """프레젠테이션에서 슬라이드를 제거한다.

    python-pptx에 삭제 API가 없어 sldIdLst 항목과 관계를 직접 지운다.
    복제가 모두 끝난 뒤에 원본을 지워야 한다 — 먼저 지우면 복제 소스가 사라진다.
    """
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        if prs.part.rels[sld_id.rId].target_part is slide.part:
            prs.part.drop_rel(sld_id.rId)
            xml_slides.remove(sld_id)
            return


def _fill_page(slide, scr: dict, page: list[dict], title: str, mapping: dict,
               work_dir: Path, warns: Warnings) -> None:
    tpl = mapping["template"]
    shapes_cfg = tpl.get("shapes", {})
    opts = mapping.get("options", {})
    cols = tpl.get("table_columns", {"no": 0, "text": 1})
    text_key = opts.get("detail_text_source", "desc")
    clear_unused = bool(opts.get("clear_unused_slots", True))

    title_name = shapes_cfg.get("title")
    if title_name:
        shp = find_shape(slide, title_name)
        if shp is None:
            warns.add(scr["id"], "shape-not-found", "제목 도형 '%s' 없음" % title_name)
        else:
            set_text(shp, title)

    sid_name = shapes_cfg.get("screen_id")
    if sid_name:
        shp = find_shape(slide, sid_name)
        if shp is None:
            warns.add(scr["id"], "shape-not-found", "화면ID 도형 '%s' 없음" % sid_name)
        else:
            set_text(shp, scr["id"])

    for key, value in (scr.get("fields") or {}).items():
        name = shapes_cfg.get(key)
        if not name:
            continue
        shp = find_shape(slide, name)
        if shp is not None:
            set_text(shp, value)

    img_name = shapes_cfg.get("image")
    if img_name:
        anchor = find_shape(slide, img_name)
        images = scr.get("images") or []
        if anchor is None:
            warns.add(scr["id"], "shape-not-found", "이미지 자리 '%s' 없음" % img_name)
        elif not images:
            warns.add(scr["id"], "no-image", "배치할 이미지가 없습니다")
        else:
            place_image(slide, anchor, work_dir / images[0])

    tables = collect_tables(slide, shapes_cfg.get("detail_tables"), warns, scr["id"])
    if tables:
        fill_slots(tables, page, cols, text_key, clear_unused, warns, scr["id"])


def build(screens_data: dict, mapping: dict, work_dir: Path, out_path: Path,
          warns: Warnings) -> dict:
    tpl = mapping["template"]
    template_path = Path(tpl["file"])
    if not template_path.is_absolute() and not template_path.exists():
        template_path = Path(work_dir) / tpl["file"]

    prs = Presentation(str(template_path))
    source_index = int(tpl.get("source_slide", 0))
    src = prs.slides[source_index]
    originals = list(prs.slides)

    slot_count = count_slots(
        collect_tables(src, tpl.get("shapes", {}).get("detail_tables"), warns)
    )

    made = 0
    split_ids: list[str] = []
    failed_ids: list[str] = []
    ok_screens = 0

    for scr in screens_data.get("screens", []):
        ok_screens += 1
        try:
            pages = chunk_details(scr.get("details", []), slot_count)
            if len(pages) > 1:
                split_ids.append(scr["id"])
                warns.add(scr["id"], "slide-split",
                          "상세 %d건이 슬롯 %d개를 넘어 %d장으로 나눴습니다"
                          % (len(scr["details"]), slot_count, len(pages)))
            for i, page in enumerate(pages):
                slide = clone_slide(prs, src)
                _fill_page(slide, scr, page,
                           page_title(scr["name"], i, len(pages)),
                           mapping, work_dir, warns)
                made += 1
        except Exception as exc:
            failed_ids.append(scr["id"])
            warns.add(scr["id"], "screen-failed", "슬라이드 생성 실패: %s" % exc)
            if len(prs.slides) > len(originals) + made:
                # 예외 직전에 만들어진 슬라이드가 남아 있다. 어디가 실패했는지
                # 결과물에서 바로 보이도록 제목만 실패 표시로 바꾼다.
                title_name = tpl.get("shapes", {}).get("title")
                shp = find_shape(prs.slides[-1], title_name) if title_name else None
                if shp is not None:
                    set_text(shp, "[생성 실패] %s" % scr.get("name", scr["id"]))
                made += 1

    for slide in originals:
        _drop_slide(prs, slide)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    return {
        "slides": made,
        "screens": ok_screens,
        "split": split_ids,
        "failed": failed_ids,
    }


def main(argv: list[str] | None = None) -> int:
    setup_stdio()
    ap = argparse.ArgumentParser(description="screens.json으로 화면설계서 PPT를 만든다")
    ap.add_argument("--screens", required=True)
    ap.add_argument("--mapping", required=True)
    ap.add_argument("--work", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args(argv)

    screens_data = read_json(Path(args.screens))
    mapping = read_json(Path(args.mapping))
    warns = Warnings()

    report = build(screens_data, mapping, Path(args.work), Path(args.out), warns)

    print("슬라이드 %d장 생성 (화면 %d개)" % (report["slides"], report["screens"]))
    if report["split"]:
        print("분할된 화면: %s" % ", ".join(report["split"]))
    if report["failed"]:
        print("실패한 화면: %s" % ", ".join(report["failed"]))
    if len(warns):
        print(warns.format())
    print("저장: %s" % args.out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
