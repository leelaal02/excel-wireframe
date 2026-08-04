# -*- coding: utf-8 -*-
"""PPTX 구조 스캔. 어느 슬라이드를 복제 소스로 쓸지 판단할 재료를 만든다."""
from __future__ import annotations

from pathlib import Path

from common import EMU_PER_INCH
from pptx import Presentation


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _scan_shape(shape) -> dict:
    table = None
    if shape.has_table:
        table = {"rows": len(shape.table.rows), "cols": len(shape.table.columns)}
    return {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "is_placeholder": bool(shape.is_placeholder),
        "left": int(shape.left) if shape.left is not None else 0,
        "top": int(shape.top) if shape.top is not None else 0,
        "width": int(shape.width) if shape.width is not None else 0,
        "height": int(shape.height) if shape.height is not None else 0,
        "text": _shape_text(shape)[:200],
        "table": table,
    }


def scan_presentation(path: Path) -> dict:
    path = Path(path)
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = [_scan_shape(s) for s in slide.shapes]
        slides.append(
            {
                "index": i,
                "layout": slide.slide_layout.name,
                "shape_count": len(shapes),
                "text_shape_count": sum(1 for s in shapes if s["text"]),
                "shapes": shapes,
            }
        )
    return {
        "file": str(path),
        "slide_width": int(prs.slide_width),
        "slide_height": int(prs.slide_height),
        "slide_size_in": [
            prs.slide_width / EMU_PER_INCH,
            prs.slide_height / EMU_PER_INCH,
        ],
        "slides": slides,
    }


def suggest_mode(report: dict) -> dict:
    """예시 슬라이드가 있으면 clone, 없으면 layout.

    '텍스트가 채워진 도형 3개 이상 + 표나 그림 1개 이상'을 예시 슬라이드의 신호로 본다.
    빈 레이아웃만 있는 템플릿은 이 조건을 통과하지 못한다.
    """
    for s in report["slides"]:
        has_visual = any(
            sh["table"] is not None or "PICTURE" in sh["shape_type"]
            for sh in s["shapes"]
        )
        if s["text_shape_count"] >= 3 and has_visual:
            return {
                "mode": "clone",
                "source_slide": s["index"],
                "reason": "슬라이드 %d에 채워진 텍스트 %d개와 표/그림이 있어 예시 슬라이드로 판단"
                % (s["index"], s["text_shape_count"]),
            }
    return {
        "mode": "layout",
        "source_slide": None,
        "reason": "예시 슬라이드를 찾지 못해 빈 레이아웃 모드로 판단",
    }
