# -*- coding: utf-8 -*-
"""PPT 템플릿을 사용자가 주지 않았을 때 쓸 기본 템플릿을 만든다.

실제 화면설계서의 표준 페이지 구조를 재현한다 — 상단 제목 바, 우측 화면ID,
가운데 큰 이미지 자리, 하단에 가로로 나란한 상세 표. 도형 이름을 의미 있게 붙여
mapping.json을 자동으로 채울 수 있게 한다.
"""
from __future__ import annotations

from pathlib import Path

from common import EMU_PER_INCH
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

DEFAULT_SHAPE_NAMES = {
    "title": "제목",
    "screen_id": "화면ID",
    "image": "화면이미지",
}

BAR_COLOR = RGBColor(0x1F, 0x3B, 0x63)
TEXT_ON_BAR = RGBColor(0xFF, 0xFF, 0xFF)
SLOT_BORDER = RGBColor(0xBB, 0xBB, 0xBB)
SLOT_FILL = RGBColor(0xF5, 0xF6, 0xF8)


def _textbox(slide, name, left, top, width, height, text, size_pt, color=None,
             bold=False):
    box = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)),
                                   Emu(int(width)), Emu(int(height)))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return box


def build_default_template(
    path: Path,
    slide_width_emu: int = 12192000,
    slide_height_emu: int = 6858000,
    table_count: int = 5,
    rows_per_table: int = 4,
) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(slide_width_emu)
    prs.slide_height = Emu(slide_height_emu)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    margin = int(0.25 * EMU_PER_INCH)
    inner_w = slide_width_emu - margin * 2
    bar_h = int(0.45 * EMU_PER_INCH)

    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(slide_width_emu), Emu(bar_h))
    bar.name = "제목바"
    bar.fill.solid()
    bar.fill.fore_color.rgb = BAR_COLOR
    bar.line.fill.background()
    bar.text_frame.text = ""

    _textbox(slide, DEFAULT_SHAPE_NAMES["title"], margin, int(0.06 * EMU_PER_INCH),
             int(inner_w * 0.6), int(0.35 * EMU_PER_INCH),
             "화면명", 18, TEXT_ON_BAR, bold=True)
    _textbox(slide, DEFAULT_SHAPE_NAMES["screen_id"],
             margin + int(inner_w * 0.62), int(0.10 * EMU_PER_INCH),
             int(inner_w * 0.38), int(0.30 * EMU_PER_INCH),
             "화면ID", 11, TEXT_ON_BAR)

    # 표 높이는 rows_per_table에 비례한다(4행일 때 1.55in — 기존 기본값과 동일).
    # 표를 늘릴수록 그만큼 이미지 자리가 줄어들어야 슬라이드 하단을 넘치지 않는다.
    row_h = (1.55 / 4) * EMU_PER_INCH
    table_h = int(row_h * rows_per_table)
    bottom_margin = int(0.30 * EMU_PER_INCH)
    tables_top = int(slide_height_emu - table_h - bottom_margin)
    img_top = bar_h + int(0.15 * EMU_PER_INCH)
    img_h = tables_top - img_top - int(0.15 * EMU_PER_INCH)

    img_slot = slide.shapes.add_shape(1, Emu(margin), Emu(img_top),
                                      Emu(inner_w), Emu(img_h))
    img_slot.name = DEFAULT_SHAPE_NAMES["image"]
    img_slot.fill.solid()
    img_slot.fill.fore_color.rgb = SLOT_FILL
    img_slot.line.color.rgb = SLOT_BORDER
    tf = img_slot.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "[화면 이미지]"
    run.font.size = Pt(14)
    run.font.color.rgb = SLOT_BORDER

    gap = int(0.06 * EMU_PER_INCH)
    table_w = (inner_w - gap * (table_count - 1)) // table_count
    for t in range(table_count):
        left = margin + (table_w + gap) * t
        shp = slide.shapes.add_table(
            rows_per_table, 2, Emu(left), Emu(tables_top), Emu(table_w), Emu(table_h)
        )
        shp.name = "상세표%d" % (t + 1)
        table = shp.table
        table.columns[0].width = Emu(int(table_w * 0.18))
        table.columns[1].width = Emu(table_w - int(table_w * 0.18))
        for r in range(rows_per_table):
            n = t * rows_per_table + r + 1
            table.cell(r, 0).text = str(n)
            table.cell(r, 1).text = "예시 설명 %d" % n
            for c in range(2):
                for p in table.cell(r, c).text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def default_template_mapping(path: Path, table_count: int = 5) -> dict:
    """생성한 기본 템플릿에 대응하는 mapping.json의 template 섹션."""
    return {
        "file": str(path),
        "mode": "clone",
        "source_slide": 0,
        "shapes": {
            "title": DEFAULT_SHAPE_NAMES["title"],
            "screen_id": DEFAULT_SHAPE_NAMES["screen_id"],
            "image": DEFAULT_SHAPE_NAMES["image"],
            "detail_tables": ["상세표%d" % (i + 1) for i in range(table_count)],
        },
        "table_columns": {"no": 0, "text": 1},
    }
